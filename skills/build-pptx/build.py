"""build-pptx — markdown → Jin-branded PPTX via python-pptx.

This file contains:
  - Slide master functions (add_title_slide, add_content_slide, etc.)
  - main() driver that parses markdown and dispatches to masters

Slide masters are independently callable from other Python code.
"""

from __future__ import annotations

import sys
from pathlib import Path

# Wire imports to sibling _shared/
sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "_shared"))

import branding  # noqa: E402
import argparse  # noqa: E402
import datetime as dt  # noqa: E402
import re  # noqa: E402

from md_loader import load_markdown, extract_title  # noqa: E402

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Emu, Inches, Pt
except ImportError:
    raise SystemExit("python-pptx not installed. Run: pip install python-pptx")

# Re-export helpers and constants from layouts package so that any code that
# imported them directly from build still works.
from layouts._common import (  # noqa: E402
    _rgb,
    _HTML_TAG_RE,
    _WS_RUN_RE,
    INK_RGB,
    WHITE_RGB,
    TURQUOISE_RGB,
    DEEPPINK_RGB,
    AMBER_RGB,
    BLUEVIOLET_RGB,
    DIM_RGB,
    MUTED_RGB,
    RULE_RGB,
    DARK_BG_RGB,
    PAPER_RGB,
    _set_bg,
    _add_rect,
    _add_text,
    _blank,
    _add_card,
    _add_runs_from_html,
    _estimate_paragraph_height,
    _render_paragraph_block,
    _get_image_aspect,
    _add_table,
    _render_media_block,
    _strip_html,
    _split_slides,
    _parse_table,
    _DEFLIST_LABEL_MAX_LEN,
    _DEFLIST_BODY_MAX_LEN,
    _detect_def_cards_from_li_html,
    _parse_slide_chunk,
)

from layouts import catalog as _catalog  # noqa: E402


# === Public API ===
def new_presentation() -> "Presentation":
    """Create a 16:9 widescreen presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_title_slide(prs, *, eyebrow: str = "", title: str, subtitle: str = "",
                    name: str = "", org: str = "", date: str = "", bg_rgb=None):
    """Title slide: dark bg (theme canvas if bg_rgb given, else navy), left
    double-rail (turquoise + deeppink).
    Eyebrow turquoise, title white, name turquoise, org deeppink, gray rule + amber date, all Geist Mono."""
    s = _blank(prs)
    _set_bg(s, bg_rgb if bg_rgb is not None else DARK_BG_RGB)

    # Left double-rail
    _add_rect(s, left=0, top=0, width=0.8, height=7.5, fill_rgb=TURQUOISE_RGB)
    _add_rect(s, left=0.8, top=0, width=0.25, height=7.5, fill_rgb=DEEPPINK_RGB)

    # Auto-shrink title font for long titles so a 3-line wrap doesn't run
    # into the subtitle. 60+ chars at 48pt regularly wraps to 3 lines.
    title_size = 48 if len(title) <= 50 else (40 if len(title) <= 80 else 32)
    title_h = 2.7 if len(title) > 50 else 2.0  # extra room for wraps

    if eyebrow:
        _add_text(s, eyebrow, left=1.3, top=1.4, width=11, height=0.4,
                  size=14, color_rgb=TURQUOISE_RGB, font=branding.MONO_FONT, bold=True)
    _add_text(s, title, left=1.3, top=1.9, width=11.0, height=title_h,
              size=title_size, color_rgb=WHITE_RGB, font=branding.MONO_FONT, bold=True)
    subtitle_top = 1.9 + title_h + 0.15
    if subtitle:
        _add_text(s, subtitle, left=1.3, top=subtitle_top, width=11.0, height=0.7,
                  size=18, color_rgb=_rgb("#E5E5EA"), font=branding.SANS_FONT)
        cursor_top = subtitle_top + 0.85
    else:
        cursor_top = subtitle_top
    if name:
        _add_text(s, name, left=1.3, top=cursor_top, width=11, height=0.4,
                  size=22, color_rgb=TURQUOISE_RGB, font=branding.MONO_FONT, bold=True)
        cursor_top += 0.5
    if org:
        _add_text(s, org, left=1.3, top=cursor_top, width=11, height=0.35,
                  size=16, color_rgb=DEEPPINK_RGB, font=branding.MONO_FONT, bold=True)
        cursor_top += 0.45

    # Gray hairline rule above the date
    _add_rect(s, left=1.3, top=cursor_top + 0.1, width=4.0, height=0.02, fill_rgb=RULE_RGB)
    if date:
        _add_text(s, date, left=1.3, top=cursor_top + 0.25, width=11, height=0.3,
                  size=12, color_rgb=AMBER_RGB, font=branding.MONO_FONT)
    return s


def add_content_slide(prs, *, title: str, body_paragraphs: list[str],
                      accent_color_hex: str | None = None,
                      images: list[Path] | None = None,
                      tables: list[list[list[str]]] | None = None,
                      cards: list[dict] | None = None,
                      name: str = "", org: str = "", deck_title: str = "",
                      date: str = ""):
    """Content slide. Dispatches to a layout renderer from layouts/catalog.

    Layout (16:9, 13.33×7.50):
      L=0      T=0    W=0.22 H=7.50  left accent bar (section color)
      L=0.50   T=0.30 W=12.30 H=0.55 title (28pt mono INK bold) — omitted if no title
      L=0.50   T=0.95 W=12.30 H=0.005 hairline rule (RULE color)
      L=0.50   T=1.05 W=12.30 H=0.40 subtitle (13pt sans MUTED) — first short paragraph
      L=0.50   T=1.55 W=12.30 H=5.30 body region — cards/media/text by content type
      L=0.50   T=7.12 W=10.00 H=0.30 footer (9pt mono MUTED): name · org · deck · date
    """
    accent_hex = accent_color_hex or branding.TURQUOISE
    accent = _rgb(accent_hex)
    images = images or []
    tables = tables or []
    cards = cards or []

    s = _blank(prs)

    body = list(body_paragraphs or [])

    has_cards = bool(cards)
    has_media = bool(images) or bool(tables)

    # Promote first paragraph to lede if it is prose and there is more below it.
    lede = ""
    has_more_below = (len(body) > 1) or has_media or has_cards
    if body and has_more_below:
        first = body[0]
        if isinstance(first, dict):
            is_bullet = first.get("kind") == "bullet"
            text = _strip_html(first.get("html", ""))
        else:
            is_bullet = first.startswith("•")
            text = first
        if not is_bullet and len(text) <= 350:
            lede = text
            body = body[1:]

    # Decide whether to use side-by-side layout.
    n_images_for_layout = len(images)
    aspect_for_layout = (_get_image_aspect(images[0])
                        if n_images_for_layout == 1 else None)
    use_side_by_side = (
        n_images_for_layout == 1
        and len(tables) == 0
        and not has_cards
        and aspect_for_layout is not None
        and aspect_for_layout <= 1.0
        and (bool(body) or bool(lede))
    )

    footer_kwargs = {
        "name": name,
        "org": org,
        "deck_title": deck_title,
        "date": date,
    }

    # Dispatch to the correct layout renderer.
    if has_cards:
        kind = "cards-grid"
        params = {
            "title": title,
            "lede": lede,
            "body": body,
            "cards": cards,
        }
    elif has_media:
        kind = "content-text-image"
        params = {
            "title": title,
            "lede": lede,
            "body": body,
            "images": images,
            "tables": tables,
            "use_side_by_side": use_side_by_side,
        }
    elif body:
        kind = "content-text"
        params = {
            "title": title,
            "lede": lede,
            "body": body,
        }
    else:
        # Nothing to render — but we still want the chrome (title + footer).
        # Use content-text with empty body.
        kind = "content-text"
        params = {
            "title": title,
            "lede": lede,
            "body": [],
        }

    _catalog.get(kind)(s, params=params, accent_rgb=accent,
                       footer_kwargs=footer_kwargs)

    return s


def add_section_divider(prs, *, label: str, index: int = 0,
                        accent_color_hex: str | None = None,
                        name: str = "", org: str = "", deck_title: str = "",
                        bg_rgb=None):
    """Section divider: dark bg (theme canvas if bg_rgb given, else navy) +
    full-height left colorblock + DMG-style top title block + bottom footer
    (name turquoise · org deeppink · deck muted).

    Color comes from `accent_color_hex` if provided, else from
    branding.pick_section_color(index) cycling.
    """
    bg_hex = accent_color_hex or branding.pick_section_color(index)
    accent = _rgb(bg_hex)
    s = _blank(prs)
    _set_bg(s, bg_rgb if bg_rgb is not None else DARK_BG_RGB)

    # Full-height left colorblock (results_overview style)
    _add_rect(s, left=0, top=0, width=0.6, height=7.5, fill_rgb=accent)

    # Small DMG-style accent bar at the top, just right of the colorblock.
    _add_rect(s, left=0.85, top=0.7, width=0.18, height=0.45, fill_rgb=accent)

    # Eyebrow = thematic category based on accent color.
    eyebrow_text = branding.category_for_accent(bg_hex)
    _add_text(s, eyebrow_text, left=1.15, top=0.7, width=11.0, height=0.4,
              size=14, color_rgb=accent, font=branding.MONO_FONT, bold=True)

    # Big section title — H1 text in ALL CAPS.
    _add_text(s, label.upper(), left=0.85, top=2.4, width=12.0, height=3.0,
              size=44, color_rgb=WHITE_RGB, font=branding.MONO_FONT, bold=True)

    # Hairline rule
    _add_rect(s, left=0.85, top=5.6, width=2.0, height=0.02, fill_rgb=accent)

    # Bottom footer
    footer_top = 6.7
    cursor_left = 0.95
    if name:
        _add_text(s, name, left=cursor_left, top=footer_top, width=4.0, height=0.35,
                  size=11, color_rgb=TURQUOISE_RGB, font=branding.MONO_FONT, bold=True)
        cursor_left += max(1.4, 0.11 * len(name))
    if org:
        if name:
            _add_text(s, "·", left=cursor_left, top=footer_top, width=0.2, height=0.35,
                      size=11, color_rgb=DIM_RGB, font=branding.MONO_FONT)
            cursor_left += 0.25
        _add_text(s, org, left=cursor_left, top=footer_top, width=5.0, height=0.35,
                  size=11, color_rgb=DEEPPINK_RGB, font=branding.MONO_FONT, bold=True)
        cursor_left += max(1.4, 0.11 * len(org))
    if deck_title:
        if name or org:
            _add_text(s, "·", left=cursor_left, top=footer_top, width=0.2, height=0.35,
                      size=11, color_rgb=DIM_RGB, font=branding.MONO_FONT)
            cursor_left += 0.25
        _add_text(s, deck_title, left=cursor_left, top=footer_top, width=8.0, height=0.35,
                  size=11, color_rgb=DIM_RGB, font=branding.MONO_FONT)

    return s


def add_big_number_slide(prs, *, number: str, caption: str = ""):
    """White background. Giant deeppink number centered, caption underneath."""
    s = _blank(prs)
    _set_bg(s, WHITE_RGB)
    _add_text(s, number, left=0.5, top=2.4, width=12.3, height=2.2,
              size=120, color_rgb=DEEPPINK_RGB, font=branding.MONO_FONT, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if caption:
        _add_text(s, caption, left=1.0, top=4.8, width=11.3, height=0.8,
                  size=18, color_rgb=MUTED_RGB, font=branding.SANS_FONT,
                  align=PP_ALIGN.CENTER)
    return s


def add_two_column_slide(prs, *, title: str,
                         left_title: str, left_body: list[str],
                         right_title: str, right_body: list[str]):
    """White background. Title at top, two side-by-side columns."""
    s = _blank(prs)
    _set_bg(s, WHITE_RGB)
    _add_text(s, title, left=0.6, top=0.4, width=12.5, height=0.8,
              size=32, color_rgb=TURQUOISE_RGB, font=branding.MONO_FONT, bold=True)
    _add_rect(s, left=0.6, top=1.25, width=12.0, height=0.005, fill_rgb=RULE_RGB)

    # Left column
    _add_text(s, left_title, left=0.6, top=1.5, width=5.8, height=0.5,
              size=18, color_rgb=DEEPPINK_RGB, font=branding.MONO_FONT, bold=True)
    _add_text(s, "\n".join(left_body), left=0.6, top=2.1, width=5.8, height=5.0,
              size=18, color_rgb=INK_RGB, font=branding.SANS_FONT)

    # Right column
    _add_text(s, right_title, left=6.95, top=1.5, width=5.8, height=0.5,
              size=18, color_rgb=DEEPPINK_RGB, font=branding.MONO_FONT, bold=True)
    _add_text(s, "\n".join(right_body), left=6.95, top=2.1, width=5.8, height=5.0,
              size=18, color_rgb=INK_RGB, font=branding.SANS_FONT)

    # Vertical hairline between columns
    _add_rect(s, left=6.7, top=1.5, width=0.005, height=5.5, fill_rgb=RULE_RGB)
    return s


def add_quote_slide(prs, *, quote: str, attribution: str = ""):
    """White background. Quote centered in italic sans, attribution mono."""
    s = _blank(prs)
    _set_bg(s, WHITE_RGB)
    _add_text(s, f'"{quote}"', left=1.5, top=2.5, width=10.3, height=2.5,
              size=36, color_rgb=INK_RGB, font=branding.SANS_FONT, italic=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if attribution:
        _add_text(s, f"— {attribution}", left=1.5, top=5.3, width=10.3, height=0.4,
                  size=14, color_rgb=MUTED_RGB, font=branding.MONO_FONT,
                  align=PP_ALIGN.CENTER)
    return s


# Thank-You slide personal defaults (this is Jinchi's branded tool). Override
# per deck via frontmatter `name`, `name_cjk`, `email` (set email: "" to omit).
_DEFAULT_NAME_CJK = "魏晉祺"
_DEFAULT_EMAIL = "jinchikwei@gmail.com"
_CJK_FONT = "Noto Sans CJK TC"
_SKILL_ASSETS = Path(__file__).resolve().parent / "assets"


def _add_endslide_logos(slide, assets_dir, *, height_in=1.0, gap_in=0.22, margin_in=0.5):
    """meng + xiang (夢想) identity logos, side-by-side in the lower-right.
    No-op if the asset files are absent."""
    paths = [Path(assets_dir) / "meng.png", Path(assets_dir) / "xiang.png"]
    avail = [p for p in paths if p.exists()]
    if not avail:
        return
    slide_w, slide_h = Inches(13.333), Inches(7.5)
    gap, margin = Inches(gap_in), Inches(margin_in)
    pics = [slide.shapes.add_picture(str(p), Inches(0), Inches(0), height=Inches(height_in))
            for p in avail]
    total_w = sum(p.width for p in pics) + gap * (len(pics) - 1)
    block_h = max(p.height for p in pics)
    left = slide_w - total_w - margin
    top = slide_h - block_h - margin
    cur = left
    for pic in pics:
        pic.left = int(cur)
        pic.top = int(top + (block_h - pic.height) // 2)
        cur = cur + pic.width + gap


def add_end_slide(prs, *, message: str = "Thank You", contact: str = "",
                  name: str = "", org: str = "", email=None, name_cjk=None,
                  bg_rgb=None, assets_dir=None):
    """Closing 'Thank You' slide — the build-pptx default. Dark bg + left
    double-rail (mirrors the title slide); big deeppink "Thank You"; the
    presenter's name (CJK over English, both turquoise); email in amber; and
    the 夢想 identity logos lower-right. `email`/`name_cjk` fall back to the
    tool owner's defaults; pass "" (or frontmatter `email: ""`) to omit either.
    `contact` is a back-compat alias for `name`."""
    name = name or contact
    email = _DEFAULT_EMAIL if email is None else str(email)
    name_cjk = _DEFAULT_NAME_CJK if name_cjk is None else str(name_cjk)
    assets = assets_dir if assets_dir else _SKILL_ASSETS

    s = _blank(prs)
    _set_bg(s, bg_rgb if bg_rgb is not None else DARK_BG_RGB)
    _add_rect(s, left=0, top=0, width=0.8, height=7.5, fill_rgb=TURQUOISE_RGB)
    _add_rect(s, left=0.8, top=0, width=0.25, height=7.5, fill_rgb=DEEPPINK_RGB)

    _add_text(s, message, left=1.3, top=1.55, width=11.0, height=1.4,
              size=60, color_rgb=DEEPPINK_RGB, font=branding.MONO_FONT, bold=True,
              align=PP_ALIGN.CENTER)
    y = 3.55
    if name_cjk:
        _add_text(s, name_cjk, left=1.3, top=y, width=11.0, height=0.7, size=30,
                  color_rgb=TURQUOISE_RGB, font=_CJK_FONT, align=PP_ALIGN.CENTER)
        y += 0.74
    if name:
        _add_text(s, name, left=1.3, top=y, width=11.0, height=0.5, size=20,
                  color_rgb=TURQUOISE_RGB, font=branding.MONO_FONT, bold=True,
                  align=PP_ALIGN.CENTER)
        y += 0.52
    if org:
        _add_text(s, org, left=1.3, top=y, width=11.0, height=0.4, size=13,
                  color_rgb=WHITE_RGB, font=branding.SANS_FONT, align=PP_ALIGN.CENTER)
        y += 0.44
    if email:
        _add_text(s, email, left=1.3, top=y + 0.04, width=11.0, height=0.4, size=16,
                  color_rgb=AMBER_RGB, font=branding.MONO_FONT, align=PP_ALIGN.CENTER)

    _add_endslide_logos(s, assets)
    return s


def _legacy_main(args) -> int:
    """The v3 rule-based render path. Preserved as the --no-plan fallback."""
    loaded = load_markdown(args.input)
    meta = loaded["meta"]
    today = dt.date.today().isoformat()

    deck_title = extract_title(loaded) or Path(args.input).stem
    deck_name = str(meta.get("name", ""))
    deck_org = str(meta.get("org", ""))
    deck_date = str(meta.get("date") or today)
    md_dir = Path(args.input).resolve().parent

    prs = new_presentation()

    if not args.no_cover:
        add_title_slide(
            prs,
            eyebrow=str(meta.get("eyebrow", "")),
            title=deck_title,
            subtitle=str(meta.get("subtitle", "")),
            name=deck_name,
            org=deck_org,
            date=deck_date,
        )

    def _emit_content(slide_data: dict, slide_title: str) -> None:
        """Emit a content slide with the section's current accent + deck footer.

        Auto-explode behavior: when a slide has 2+ images and no cards/tables,
        each image becomes its own full-bleed slide titled by its alt text,
        with the immediately preceding paragraph as its lede. Matches Jin's
        DMG v2 deck where each chart was on its own slide rather than three
        tiny figures crammed into a row."""
        image_records = slide_data.get("image_records") or []
        body_paras = slide_data.get("body", [])
        if (len(image_records) >= 2
                and not slide_data.get("cards")
                and not slide_data.get("tables")):
            # Walk by source position; pair each image with the most recent
            # preceding paragraph (consumed once it becomes a lede).
            text_items = sorted(
                [p for p in body_paras if p.get("pos") is not None],
                key=lambda p: p["pos"],
            )
            ti = 0
            consumed = set()
            for img_rec in image_records:
                lede_para = None
                while ti < len(text_items) and text_items[ti]["pos"] < img_rec["pos"]:
                    if id(text_items[ti]) not in consumed:
                        lede_para = text_items[ti]
                    ti += 1
                if lede_para is not None:
                    consumed.add(id(lede_para))
                sub_title = (img_rec["alt"] or slide_title or "").strip() or slide_title
                sub_body = [lede_para] if lede_para else []
                add_content_slide(
                    prs,
                    title=sub_title,
                    body_paragraphs=sub_body,
                    images=[img_rec["path"]],
                    tables=[],
                    cards=None,
                    accent_color_hex=current_accent,
                    name=deck_name, org=deck_org,
                    deck_title=deck_title, date=deck_date,
                )
            # Trailing paragraphs after the last image — render as a final
            # text-only slide using the parent title.
            trailing = [p for p in text_items
                        if p["pos"] > image_records[-1]["pos"]
                        and id(p) not in consumed]
            if trailing:
                add_content_slide(
                    prs,
                    title=slide_title,
                    body_paragraphs=trailing,
                    images=[],
                    tables=[],
                    cards=None,
                    accent_color_hex=current_accent,
                    name=deck_name, org=deck_org,
                    deck_title=deck_title, date=deck_date,
                )
            return
        add_content_slide(
            prs,
            title=slide_title,
            body_paragraphs=body_paras,
            images=slide_data["images"],
            tables=slide_data["tables"],
            cards=slide_data.get("cards"),
            accent_color_hex=current_accent,
            name=deck_name, org=deck_org,
            deck_title=deck_title, date=deck_date,
        )

    # Walk slide chunks; track current section accent.
    chunks = _split_slides(loaded["body_html"])
    current_accent = branding.TURQUOISE

    for chunk in chunks:
        h1_match = re.search(r"<h1[^>]*>(.*?)</h1>", chunk)
        if h1_match:
            section_label = _strip_html(h1_match.group(1))
            current_accent = branding.match_section_color(section_label)
            add_section_divider(prs, label=section_label,
                                accent_color_hex=current_accent,
                                name=deck_name, org=deck_org,
                                deck_title=deck_title)
            remaining = chunk[h1_match.end():].strip()
            if remaining:
                slide = _parse_slide_chunk(remaining, base_dir=md_dir)
                if any(slide.get(k) for k in ("title", "body", "images", "tables", "cards")):
                    slide_title = slide["title"] or section_label
                    _emit_content(slide, slide_title)
        else:
            slide = _parse_slide_chunk(chunk, base_dir=md_dir)
            if any(slide.get(k) for k in ("title", "body", "images", "tables", "cards")):
                _emit_content(slide, slide["title"])

    if not args.no_end:
        add_end_slide(prs, name=str(meta.get("name") or ""),
                      org=str(meta.get("org") or ""),
                      email=meta.get("email"), name_cjk=meta.get("name_cjk"))

    prs.save(args.output)
    print(f"wrote {args.output}")
    return 0


_CONCLUSIONS_KEYWORDS = (
    "takeaways", "takeaway", "conclusions", "conclusion",
    "summary", "next steps", "key findings", "closing", "final thoughts",
)


# Pipeline-style keywords for vertical-timeline auto-detection.
# When a slide title (or its parent H1) contains one of these AND the slide
# has 3-7 cards with short labels, prefer a vertical-timeline render so the
# sequential nature is visually carried.
_PIPELINE_KEYWORDS = (
    "pipeline", "workflow", "stages", "phases",
    "process", "procedure", "protocol",
    "steps", "step-by-step",
)


def _is_pipeline_slide(title: str, section_label: str | None) -> bool:
    """Case-insensitive substring match for pipeline-style slides."""
    candidates = [title or "", section_label or ""]
    for candidate in candidates:
        c = candidate.lower()
        for kw in _PIPELINE_KEYWORDS:
            if kw in c:
                return True
    return False


def _extract_descriptor(body_text: str, max_len: int = 40) -> str:
    """Extract a short descriptor (3-6 words / ≤max_len chars) from the start
    of a stat-tile body. Returns just the descriptor; callers keep the FULL
    body as the sub-line so the tile reads as
        {short descriptor}    ← header
        {value}               ← big number
        {full body sentence}  ← sub
    The header gives identification, the sub keeps the explanation.

    Splits at the first strong delimiter in (period, semicolon, em-dash,
    colon, opening paren, comma) within max_len chars; otherwise truncates at
    the last word boundary before max_len.
    """
    if not body_text:
        return ""
    s = body_text.strip()
    if not s:
        return ""
    for sep in (". ", "; ", " — ", " – ", ": ", " (", ", "):
        idx = s.find(sep)
        if 8 <= idx <= max_len:
            return s[:idx].rstrip(" .,;:—–(").strip()
    # No delimiter — truncate at last word boundary before max_len
    if len(s) <= max_len:
        return s.rstrip(" .,").strip()
    cut = s.rfind(" ", 0, max_len)
    if cut <= 8:
        cut = max_len
    return s[:cut].rstrip(" .,;:—–(").strip()


_STAT_PREFIXES = (
    "n ", "n=", "n =", "p ", "p=", "p =", "p<", "p <", "p>", "p >",
    "r ", "r=", "r =", "r<", "r <", "r>", "r >",
    "β", "ρ", "δ", "μ", "σ", "χ",
    "or ", "or=", "or =", "auc", "rr ", "hr ", "ci ", "ci=", "ci =",
    "f=", "f =", "t=", "t =", "z=", "z =",
)


def _looks_like_stat_label(label: str) -> bool:
    """True if a card label reads as a stat token (e.g. '0.91', 'OR = 2.44',
    'p < 0.001', 'r = +0.92', '68.5%'). Used to distinguish
    stats-with-takeaway from cards-grid.

    Heuristic — short (≤ 25 chars), contains a digit, and matches AT LEAST
    ONE of the following stat markers:
      1. Pure numeric label (digits + decimals/operators only, no letters):
         '0.91', '2.44', '11/11', '68.5%'.
      2. Has a comparator / unit / sign character: '=', '<', '>', '≤', '≥',
         '±', '~', '%', '‰'.
      3. Begins with a known stat prefix: 'n', 'p', 'r', 'OR', 'AUC', 'β',
         'ρ', 'F', 't', 'z', 'χ', 'HR', 'RR', 'CI', etc.

    Plain "year-like" tokens such as 'RSNA 2026' or 'Q1 2026' have no
    operator/unit and don't start with a stat prefix, so they correctly
    classify as categorical (NOT stat).
    """
    if not label:
        return False
    s = label.strip()
    if not s or len(s) > 25:
        return False
    n_digits = sum(1 for c in s if c.isdigit())
    if n_digits == 0:
        return False

    # 1. Pure numeric (no letters at all) → always a stat
    n_letters = sum(1 for c in s if c.isalpha())
    if n_letters == 0:
        return True

    # 2. Contains an explicit stat marker (operator / unit / sign)
    if any(ch in s for ch in "=<>≤≥±~%‰"):
        return True

    s_low = s.lower()

    # 3. Starts with a known stat prefix
    for prefix in _STAT_PREFIXES:
        if s_low.startswith(prefix):
            return True

    # 4. Number followed by a known scientific unit word (e.g. '11 regions',
    # '287 scans', '54 subjects'). Distinguishes count-with-unit stats from
    # 'RSNA 2026' / 'Manuscript' phrases.
    _UNIT_WORDS = (
        "region", "regions", "scan", "scans", "subject", "subjects",
        "obs", "observation", "observations", "voxel", "voxels",
        "case", "cases", "sample", "samples", "participant", "participants",
        "patient", "patients", "trial", "trials", "epoch", "epochs",
        "iteration", "iterations", "step", "steps", "minute", "minutes",
        "hour", "hours", "day", "days", "week", "weeks", "month", "months",
        "year", "years", "fold", "folds", "site", "sites",
    )
    for unit in _UNIT_WORDS:
        if f" {unit}" in s_low or s_low.endswith(unit):
            # Plus require the digit to come before the unit (so 'years 2025'
            # wouldn't pass — only 'X years' / 'X regions' shapes do)
            unit_idx = s_low.find(unit)
            if any(c.isdigit() for c in s_low[:unit_idx]):
                return True

    # Otherwise it's something like 'RSNA 2026' or 'Q3 2025' — categorical.
    return False


def _split_bullet_to_label_body(text: str) -> tuple[str, str]:
    """Try to split a bullet's plain text into a (label, body) pair so it can
    render as a card with a header instead of a flat line.

    Recognized shapes (in order):
      - "**bold prefix** rest"       → label="bold prefix", body="rest"
      - "Label: body"                → label="Label", body="body"
                                       (label ≤ 30 chars, no nested colon)
      - "Label — body"               → label="Label", body="body" (em-dash)
      - "Label - body"               → label="Label", body="body" (hyphen,
                                       label ≤ 24 chars, body must exist)

    Falls back to ("", text) when no clean split is found. Empty labels render
    as a body-only card row, which is still denser than content-text bullets.
    """
    if not text:
        return ("", "")
    s = text.strip()

    # Already-stripped HTML so a leading bold span looks like "**foo**". The
    # markdown parser may also have bolded the prefix via inline formatting
    # before _strip_html, leaving the asterisks behind — handle either form.
    m = re.match(r"^\*\*(.+?)\*\*[\s:—–-]*(.*)$", s)
    if m and m.group(1).strip() and m.group(2).strip():
        return (m.group(1).strip(), m.group(2).strip())

    # "Label: body" — label cannot itself contain a colon.
    if ":" in s:
        head, _, tail = s.partition(":")
        head, tail = head.strip(), tail.strip()
        if 2 <= len(head) <= 30 and tail and ":" not in head:
            return (head, tail)

    # "Label — body"
    for sep in (" — ", " – "):
        if sep in s:
            head, _, tail = s.partition(sep)
            head, tail = head.strip(), tail.strip()
            if 2 <= len(head) <= 40 and tail:
                return (head, tail)

    # "Label - body" (single hyphen with spaces, conservative length cap)
    if " - " in s:
        head, _, tail = s.partition(" - ")
        head, tail = head.strip(), tail.strip()
        if 2 <= len(head) <= 24 and tail:
            return (head, tail)

    return ("", s)


def _is_closing_slide(title: str, section_label: str | None) -> bool:
    """Return True when a slide title (or parent H1 section label) matches
    the closing-slide pattern — case-insensitive substring match."""
    candidates = [title or "", section_label or ""]
    for candidate in candidates:
        lower = candidate.lower()
        if any(kw in lower for kw in _CONCLUSIONS_KEYWORDS):
            return True
    return False


def _infer_default_plan(*, md_text: str, chunks: list[str],
                        slide_records: list[dict], deck_md_hash: str,
                        base_dir: Path):
    """Build a Plan from chunks using the same dispatch logic the legacy
    renderer uses. Lets the v4 path produce the same output as v3 when no
    Claude-generated plan exists yet (Task 6 will replace this with real
    layout-picking via in-session reasoning)."""
    from plan import Plan, SlideEntry

    slides: list[SlideEntry] = []
    current_h1: str | None = None
    for rec, chunk in zip(slide_records, chunks):
        slide_id = rec["slide_id"]
        content_hash = rec["content_hash"]

        h1_match = re.search(r"<h1[^>]*>(.*?)</h1>", chunk, re.DOTALL)
        if h1_match:
            section_label = _strip_html(h1_match.group(1))
            current_h1 = section_label
            accent_hex = branding.match_section_color(section_label)
            # Emit a section divider entry FIRST
            slides.append(SlideEntry(
                slide_id=f"divider-{slide_id}",
                kind="section-divider",
                params={"label": section_label, "accent_hex": accent_hex},
                content_hash=content_hash + "-divider",
            ))
            # If the chunk has body content beyond the H1, also emit a content slide
            remaining = chunk[h1_match.end():].strip()
            if not remaining:
                continue
            sd = _parse_slide_chunk(remaining, base_dir=base_dir)
            slide_title = sd["title"] or section_label
        else:
            sd = _parse_slide_chunk(chunk, base_dir=base_dir)
            slide_title = sd["title"]

        # Skip empty chunks
        if not any(sd.get(k) for k in ("title", "body", "images", "tables", "cards")):
            continue

        # Choose layout kind via the same heuristics as add_content_slide
        body = list(sd["body"])
        images = sd["images"]
        tables = sd["tables"]
        cards = sd["cards"]
        image_records = sd.get("image_records", [])

        # Lede extraction (mirrors add_content_slide logic)
        lede = ""
        has_more_below = (len(body) > 1) or bool(images) or bool(tables) or bool(cards)
        if body and has_more_below:
            first = body[0]
            text = _strip_html(first.get("html", ""))
            if first.get("kind") != "bullet" and len(text) <= 350:
                lede = text
                body = body[1:]

        # Multi-image auto-explode (mirrors _emit_content logic)
        if (len(image_records) >= 2 and not cards and not tables):
            text_items = sorted(
                [p for p in sd["body"] if p.get("pos") is not None],
                key=lambda p: p["pos"],
            )
            ti = 0
            consumed = set()
            for j, img_rec in enumerate(image_records):
                lede_para = None
                while ti < len(text_items) and text_items[ti]["pos"] < img_rec["pos"]:
                    if id(text_items[ti]) not in consumed:
                        lede_para = text_items[ti]
                    ti += 1
                if lede_para is not None:
                    consumed.add(id(lede_para))
                sub_lede = _strip_html(lede_para["html"]) if lede_para else ""
                sub_title = (img_rec["alt"] or slide_title or "").strip() or slide_title
                # Score-based dispatch (mirrors the single-image path below):
                # if the lede + image aspect score >= 2, prefer figure-with-aside
                # over content-image-only so multi-image H1 sections still reach
                # for the richer layout per image.
                sub_aspect = _get_image_aspect(str(img_rec["path"]))
                sub_score = 0
                if sub_aspect is not None and sub_lede:
                    if sub_aspect >= 1.3:
                        sub_score += 2
                    elif sub_aspect >= 1.0:
                        sub_score += 1
                    if 50 <= len(sub_lede) < 500:
                        sub_score += 2
                    elif len(sub_lede) < 1000:
                        sub_score += 1
                    sub_score += 1  # body=[] always satisfies bullet count
                if sub_score >= 2 and sub_aspect is not None:
                    sub_kind = ("figure-with-aside-horizontal"
                                if sub_aspect >= 1.8 else "figure-with-aside")
                    # split lede → aside label (first sentence) + body (rest)
                    label, body_rest = "", sub_lede
                    if sub_lede.count(". ") >= 1:
                        first_sent = sub_lede.split(". ", 1)[0].strip().rstrip(".")
                        if 8 <= len(first_sent) <= 90:
                            label = first_sent
                            body_rest = sub_lede[len(first_sent) + 2:].strip()
                    sub_params = {
                        "title": sub_title, "lede": "",
                        "section_label": current_h1 or "",
                        "image": str(img_rec["path"]), "alt": "",
                        "aside": {"label": label, "body": body_rest, "icon": None},
                    }
                else:
                    sub_kind = "content-image-only"
                    sub_params = {
                        "title": sub_title, "lede": sub_lede,
                        "body": [],
                        "images": [str(img_rec["path"])],
                        "tables": [],
                        "section_label": current_h1 or "",
                    }
                slides.append(SlideEntry(
                    slide_id=f"{slide_id}/img-{j}",
                    kind=sub_kind,
                    params=sub_params,
                    content_hash=f"{content_hash}-img-{j}",
                ))
            # Trailing paragraphs after the last image
            trailing = [p for p in text_items
                        if p["pos"] > image_records[-1]["pos"]
                        and id(p) not in consumed]
            if trailing:
                params = {
                    "title": slide_title, "lede": "",
                    "body": trailing, "section_label": current_h1 or "",
                }
                slides.append(SlideEntry(
                    slide_id=f"{slide_id}/trailing",
                    kind="content-text",
                    params=params,
                    content_hash=f"{content_hash}-trailing",
                ))
            continue

        # Choose layout
        if cards and _is_closing_slide(slide_title, current_h1):
            # Auto-fire conclusions for closing slides with cards
            _CLOSING_ACCENTS = [
                branding.TURQUOISE, branding.DEEPPINK,
                branding.AMBER, branding.BLUEVIOLET,
            ]
            _CLOSING_ICONS = [
                "FaChartLine", "FaCheckCircle",
                "FaExclamationTriangle", "FaCrosshairs",
            ]
            card_list = []
            for i, c in enumerate(cards):
                # Same two-line treatment as stat-eligible cards-grid: when a
                # closing-slide card label is a bare stat token (e.g.
                # "r = +0.92" / "11 / 11"), prepend a short descriptor
                # extracted from the body so the card reads as
                #     {descriptor}
                #     · {value}
                # Already-descriptive labels (e.g. "Survives controls",
                # "RSNA 2026") leave as-is.
                raw_label = c.get("label") or ""
                final_label = raw_label
                if _looks_like_stat_label(raw_label):
                    body_clean = _strip_html(c.get("body", "") or "").strip()
                    descriptor = _extract_descriptor(body_clean, max_len=32)
                    if descriptor:
                        final_label = f"{descriptor}\n· {raw_label}"
                card_entry = {
                    "label": final_label,
                    "body": c["body"],
                    "icon": str(c["icon"]) if c.get("icon") else _CLOSING_ICONS[i % len(_CLOSING_ICONS)],
                    "accent_hex": _CLOSING_ACCENTS[i % len(_CLOSING_ACCENTS)],
                }
                card_list.append(card_entry)
            # Determine callout text: prefer trailing body paragraph; fall back
            # to lede (the lede-extractor may have consumed the "Path forward"
            # sentence before we got here). When lede becomes the callout we
            # clear it to avoid duplicate display.
            callout_text = ""
            if body:
                callout_text = " ".join(
                    _strip_html(b.get("html", "")) for b in body if b.get("html")
                ).strip()
            if not callout_text and lede:
                callout_text = lede
                lede = ""  # promote to callout; don't show as lede too

            kind = "conclusions"
            params = {
                "title": slide_title,
                "lede": lede,
                "section_label": current_h1 or "",
                "cards": card_list,
            }
            if callout_text:
                params["callout"] = {"text": callout_text, "tone": "dark"}
        elif cards:
            # If 2+ cards AND every label looks like a stat token, prefer
            # stats-with-takeaway (big-number tiles + dark callout footer)
            # over cards-grid. Trailing prose / lede promotes to callout.
            stat_eligible = (
                2 <= len(cards) <= 8
                and not images
                and not tables
                and all(_looks_like_stat_label(c.get("label", "")) for c in cards)
            )
            if stat_eligible:
                # Prepend a descriptor extracted from the body to each card's
                # numeric label so each card reads as
                #   '{descriptor} · {n = 476}'
                # instead of the bare 'n = 476' that's unreadable without
                # context. Render via cards-with-takeaway when a slide-level
                # takeaway exists (lede or trailing prose) — keeps the dark
                # navy callout footer that anchors the slide's main message.
                # Falls back to plain cards-grid when no callout text exists.
                new_cards = []
                for c in cards:
                    body_clean = _strip_html(c.get("body", "") or "").strip()
                    descriptor = _extract_descriptor(body_clean, max_len=32)
                    if descriptor:
                        # Two-line label: descriptor on line 1, value on line
                        # 2 with leading "· " to indicate the value is a
                        # sub-stat of the descriptor. Avoids horizontal
                        # crowding and overlap with the body text below.
                        combined_label = f"{descriptor}\n· {c['label']}"
                    else:
                        combined_label = c["label"]
                    new_cards.append({
                        "label": combined_label,
                        "body": body_clean,
                        "icon": str(c["icon"]) if c.get("icon") else None,
                    })
                # Resolve takeaway text + decide whether the slide warrants a
                # dark navy footer at all:
                #   - trailing body prose ALWAYS promotes (it's a real takeaway)
                #   - lede only promotes when substantive (≥ 180 chars).
                #     A short lede (e.g. "Singleshell + multishell side-by-side")
                #     reads better as the regular subtitle in chrome — no
                #     footer needed.
                LEDE_TAKEAWAY_MIN_CHARS = 180
                callout_text = ""
                if body:
                    callout_text = " ".join(
                        _strip_html(b.get("html", "")) for b in body if b.get("html")
                    ).strip()
                lede_for_chrome = lede
                if not callout_text and lede and len(lede) >= LEDE_TAKEAWAY_MIN_CHARS:
                    callout_text = lede
                    lede_for_chrome = ""
                if callout_text:
                    kind = "cards-with-takeaway"
                    params = {
                        "title": slide_title,
                        "lede": lede_for_chrome,
                        "section_label": current_h1 or "",
                        "cards": new_cards,
                        "callout": {"text": callout_text, "tone": "dark"},
                    }
                else:
                    kind = "cards-grid"
                    params = {
                        "title": slide_title,
                        "lede": lede,
                        "body": body,
                        "section_label": current_h1 or "",
                        "cards": new_cards,
                    }
            elif (
                3 <= len(cards) <= 7
                and not images
                and not tables
                and not _is_closing_slide(slide_title, current_h1)
                and _is_pipeline_slide(slide_title, current_h1)
                and all(len((c.get("label") or "").strip()) <= 30 for c in cards)
            ):
                # Pipeline-keyword slide (Pipeline / Workflow / Stages /
                # Process / Steps...) with 3-7 short-labeled cards →
                # vertical-timeline. The dots-and-rail visual carries the
                # sequential nature where flat card grids feel disjoint.
                kind = "vertical-timeline"
                params = {
                    "title": slide_title,
                    "lede": lede,
                    "section_label": current_h1 or "",
                    "stages": [
                        {"label": c["label"], "body": c["body"]}
                        for c in cards
                    ],
                }
            elif (
                2 <= len(cards) <= 4
                and not images
                and not tables
                and not _is_closing_slide(slide_title, current_h1)
            ):
                # 2-4 cards (non-closing, non-stat, non-pipeline) →
                # cards-triple. Flat full-width stacked rows that fill
                # sparse-text slides. No hierarchy implication (vs
                # cards-heterogeneous), no false sequence (vs three-pillars
                # with arrows), and bigger cards than cards-grid 1×3.
                # 2-card variant also dispatches here — same row pattern,
                # avoids the awkward 50/50 visual of cards-grid n=2.
                kind = "cards-triple"
                params = {
                    "title": slide_title,
                    "lede": lede,
                    "section_label": current_h1 or "",
                    "cards": [
                        {"label": c["label"], "body": c["body"],
                         "icon": str(c["icon"]) if c.get("icon") else None}
                        for c in cards
                    ],
                }
            else:
                kind = "cards-grid"
                params = {"title": slide_title, "lede": lede, "body": body,
                          "cards": [{"label": c["label"], "body": c["body"],
                                     "icon": str(c["icon"]) if c.get("icon") else None}
                                    for c in cards],
                          "section_label": current_h1 or ""}
        elif images or tables:
            n_images = len(images)
            aspect = _get_image_aspect(images[0]) if n_images == 1 else None
            use_side_by_side = (
                n_images == 1 and len(tables) == 0 and not cards
                and aspect is not None and aspect <= 1.0
                and (bool(body) or bool(lede))
            )
            # figure-with-aside: 1 wide image + light commentary (≤3 short
            # paragraphs/bullets, < 350 chars total). Heavy-caption slides
            # (long bullets, multi-paragraph captions) keep the centered
            # content-text-image layout because their text doesn't fit a
            # 1/3-width aside card.
            body_text = " ".join(
                _strip_html(b.get("html", "") or "") for b in (body or [])
            ).strip() if body else ""
            # Score-based dispatch (replaces the old hard-predicate gate).
            # figure-with-aside accumulates points from:
            #   - aspect: 2 if >= 1.3, 1 if >= 1.0 (square images can still work);
            #   - caption length: 2 for ideal 50-500 chars; 1 for 500-1000 chars;
            #   - body bullet count: 1 if <= 6 bullets;
            # If the cumulative score >= ASIDE_SCORE_THRESHOLD (=2), pick the
            # aside layout; otherwise fall through to content-text-image /
            # content-image-only as before. Tighter rules (`aspect > 1.3 AND
            # len < 500 AND <=6 bullets`) still always score >= 4 → unchanged
            # behavior on previously-aside slides; only previously-rejected
            # slides at the boundary now get aside.
            ASIDE_SCORE_THRESHOLD = 2
            aside_score = 0
            aside_basics_ok = (
                n_images == 1 and len(tables) == 0 and not cards
                and aspect is not None and (bool(body) or bool(lede))
            )
            if aside_basics_ok:
                if aspect >= 1.3:
                    aside_score += 2
                elif aspect >= 1.0:
                    aside_score += 1
                cap_len = len(body_text) if body_text else len(lede or "")
                if 50 <= cap_len < 500:
                    aside_score += 2
                elif cap_len < 1000:
                    aside_score += 1
                if len(body or []) <= 6:
                    aside_score += 1
            aside_eligible = aside_score >= ASIDE_SCORE_THRESHOLD
            if aside_eligible:
                # Very wide images (panel composites, e.g. 3-up subplot rows)
                # don't read well in figure-with-aside's 2:1 horizontal split
                # because the figure gets crushed into 2/3 width. Aspect ≥ 1.8
                # → figure-with-aside-horizontal: figure on top full-width,
                # aside below as a caption strip with top accent stripe.
                kind = ("figure-with-aside-horizontal" if aspect >= 1.8
                        else "figure-with-aside")

                # Resolve aside body + lede without duplicating content.
                # When body_text is empty (lede was promoted from body and
                # no other body items remain), the aside has nothing of its
                # own — pull lede into the aside and clear the lede so chrome
                # doesn't render the same string twice (top + bottom).
                # When body_text exists, lede stays at top and body fills aside.
                if body_text:
                    aside_body_text = body_text
                    lede_for_chrome = lede
                else:
                    aside_body_text = lede
                    lede_for_chrome = ""

                # Synthesize an aside label from the first sentence of the body
                # when the body is multi-sentence — gives the card a header
                # instead of a sparse wall of text.
                aside_label = ""
                if aside_body_text and aside_body_text.count(". ") >= 2:
                    first_sent = aside_body_text.split(". ", 1)[0].strip()
                    if 8 <= len(first_sent) <= 90:
                        aside_label = first_sent.rstrip(".")
                        aside_body_text = aside_body_text[len(first_sent) + 2:].strip()
                        # Lift trailing period from label
                        aside_label = aside_label.rstrip(".")

                params = {
                    "title": slide_title,
                    "lede": lede_for_chrome,
                    "section_label": current_h1 or "",
                    "image": str(images[0]),
                    "alt": "",
                    "aside": {
                        "label": aside_label,
                        "body": aside_body_text,
                        "icon": None,
                    },
                }
            elif tables and not images:
                # Tables-only slide → table-with-takeaway. Picking
                # content-text-image for table-only content (the historical
                # default) leaves an empty image gutter and renders sparse.
                # The table-with-takeaway layout fills the slide with a
                # full-width table + an optional dark accent footer.
                first_table = tables[0] if tables and isinstance(tables[0], list) \
                              and tables[0] and isinstance(tables[0][0], list) \
                              else tables  # accept both nested + flat shapes
                callout_text = ""
                if body:
                    callout_text = " ".join(
                        _strip_html(b.get("html", "") or "") for b in body if b.get("html")
                    ).strip()
                if not callout_text and lede:
                    callout_text = lede
                    lede = ""
                kind = "table-with-takeaway"
                params = {
                    "title": slide_title,
                    "lede": lede,
                    "section_label": current_h1 or "",
                    "rows": first_table,
                }
                if callout_text:
                    params["callout"] = {"text": callout_text, "tone": "dark"}
            else:
                kind = "content-text-image" if (body or lede) else "content-image-only"
                params = {"title": slide_title, "lede": lede,
                          "body": body if kind == "content-text-image" else [],
                          "images": [str(p) for p in images],
                          "tables": tables,
                          "use_side_by_side": use_side_by_side,
                          "section_label": current_h1 or ""}
        elif body:
            # Promote sparse bullet-only body to cards-triple so a 3-bullet
            # TL;DR slide doesn't render as a wall of empty space below the
            # bullets. Triggers when 2-5 bullet-only items, no images / tables
            # / cards, and the body would otherwise spread thinly across the
            # slide. Each bullet becomes a card; if the bullet has a clear
            # "label: body" or "**label** body" shape, split it; otherwise
            # the bullet text becomes the body and label is left empty.
            bullet_only = (
                2 <= len(body) <= 5
                and all(b.get("kind") == "bullet" for b in body)
            )
            if bullet_only:
                bullet_cards = []
                for b in body:
                    text = _strip_html(b.get("html", "") or "").strip()
                    label, card_body = _split_bullet_to_label_body(text)
                    bullet_cards.append({"label": label, "body": card_body, "icon": None})

                # Closing-section bullet slides (Conclusions / Takeaways /
                # Next steps / etc) get the dark-bg conclusions layout
                # instead of cards-triple — same trigger as the cards-detected
                # path at line 811. Each card gets a brand-cycle accent + icon.
                if _is_closing_slide(slide_title, current_h1):
                    _ACCENTS = [branding.TURQUOISE, branding.DEEPPINK,
                                branding.AMBER, branding.BLUEVIOLET]
                    _ICONS = ["FaChartLine", "FaCheckCircle",
                              "FaExclamationTriangle", "FaCrosshairs"]
                    enriched_cards = []
                    for i, c in enumerate(bullet_cards):
                        enriched_cards.append({
                            "label": c["label"],
                            "body": c["body"],
                            "icon": _ICONS[i % len(_ICONS)],
                            "accent_hex": _ACCENTS[i % len(_ACCENTS)],
                        })
                    kind = "conclusions"
                    params = {
                        "title": slide_title,
                        "lede": lede,
                        "section_label": current_h1 or "",
                        "cards": enriched_cards,
                    }
                else:
                    kind = "cards-triple"
                    params = {
                        "title": slide_title,
                        "lede": lede,
                        "section_label": current_h1 or "",
                        "cards": bullet_cards,
                    }
            else:
                kind = "content-text"
                params = {"title": slide_title, "lede": lede, "body": body,
                          "section_label": current_h1 or ""}
        else:
            # Empty chunk — skip
            continue

        slides.append(SlideEntry(
            slide_id=slide_id, kind=kind, params=params,
            content_hash=content_hash,
        ))

    return Plan(version=1, deck_md_hash=deck_md_hash, slides=slides)


def main() -> int:
    ap = argparse.ArgumentParser(description="markdown → Jin-branded PPTX")
    ap.add_argument("--input", required=True)
    ap.add_argument("--output", required=True)
    ap.add_argument("--no-cover", dest="no_cover", action="store_true",
                    help="suppress title slide")
    ap.add_argument("--no-end", dest="no_end", action="store_true",
                    help="suppress closing 'Thanks' slide")
    ap.add_argument("--shake", action="store_true",
                    help="regenerate the layout plan from scratch")
    ap.add_argument("--plan-only", dest="plan_only", action="store_true",
                    help="emit only the plan JSON; do not render pptx")
    ap.add_argument("--no-plan", dest="no_plan", action="store_true",
                    help="bypass v4 plan path; use legacy rule-based renderer")
    ap.add_argument(
        "--use-blocks",
        dest="use_blocks",
        choices=["auto", "never", "always"],
        default="auto",
        help=(
            "Control when composition/freeform layouts are admissible. "
            "'auto' (default): planner may pick composition/freeform when the "
            "decision rubric says appropriate. "
            "'never': forbid composition and freeform entirely; if a sidecar "
            "contains them the build is aborted with a clear error message. "
            "'always': signal that the agent should prefer composition/freeform "
            "(useful for explicit experiments; does not force-rewrite an existing "
            "sidecar, but informs inline plan generation when no sidecar exists)."
        ),
    )
    ap.add_argument(
        "--mode",
        dest="mode",
        choices=["expressive", "strict"],
        default=None,
        help=(
            "Deck construction mode. 'expressive' (default): themed + "
            "guided-freeform, Anthropic-pptx aesthetic. 'strict': rules-based "
            "named-layout behavior, the revert path. If omitted, an existing "
            "sidecar's recorded mode wins; otherwise defaults to expressive."
        ),
    )
    ap.add_argument(
        "--allow-composed", dest="allow_composed", action="store_true",
        help=(
            "Permit auto-composed (agentless-FLOOR) freeform slides without "
            "failing the build. By default an expressive render ABORTS if any "
            "content slide is still composer-templated rather than handcrafted "
            "(see bespoke_design.md). Use this flag ONLY for non-interactive / "
            "cron renders where the deterministic floor is the accepted output; "
            "an agent in the loop must handcraft each slide instead."
        ),
    )
    ap.add_argument(
        "--no-notes-pdf", dest="no_notes_pdf", action="store_true",
        help=(
            "Skip the presenter-handout PDF (<output>_notes.pdf: each slide's "
            "image + its speaker notes). By default it is generated whenever the "
            "deck has notes and a rasterizer (LibreOffice + pdftoppm) is present."
        ),
    )
    ap.add_argument("--qa", action="store_true",
                    help="after rendering, emit per-slide PNGs for visual "
                         "inspection (requires LibreOffice + poppler)")
    args = ap.parse_args()

    if args.no_plan:
        # Legacy path — current main() behavior
        return _legacy_main(args)

    # Plan path
    from plan import (
        Plan, SlideEntry, build_slide_records, derive_slide_ids_from_chunks,
        merge_with_existing, hash_text, assemble_plan_prompt,
    )

    md_path = Path(args.input).resolve()
    output_path = Path(args.output).resolve()
    sidecar_path = md_path.with_suffix(md_path.suffix + ".layout.json")

    md_text = md_path.read_text(encoding="utf-8")
    loaded = load_markdown(str(md_path))
    chunks = _split_slides(loaded["body_html"])
    slide_ids = derive_slide_ids_from_chunks(chunks)
    slide_records = build_slide_records(chunks=chunks, slide_ids=slide_ids)
    deck_md_hash = hash_text(md_text)

    # Read existing sidecar if present and not shaking
    existing_plan = None
    if sidecar_path.exists() and not args.shake:
        try:
            existing_plan = Plan.from_json(sidecar_path.read_text(encoding="utf-8"))
        except Exception as e:
            print(f"warning: could not parse existing sidecar: {e}", file=sys.stderr)
            existing_plan = None

    # Build a default plan (rule-based layout choice from chunk content)
    default_plan = _infer_default_plan(
        md_text=md_text, chunks=chunks, slide_records=slide_records,
        deck_md_hash=deck_md_hash, base_dir=md_path.parent,
    )

    # Merge with existing
    final_plan = merge_with_existing(default_plan, existing_plan)

    # Mode/theme/seed are user-or-prior intent that must survive re-renders
    # (and mode must survive --shake; strict is the deliberate revert path).
    # Read them from any existing sidecar in one parse.
    prior_mode = None
    prior_theme = None
    prior_seed = None
    if sidecar_path.exists():
        try:
            _prev = Plan.from_json(sidecar_path.read_text(encoding="utf-8"))
            prior_mode, prior_theme, prior_seed = _prev.mode, _prev.theme, _prev.shake_seed
        except Exception:
            pass

    # Resolve effective mode: explicit flag > prior sidecar mode > default.
    # prior_mode already covers the non-shake case (it reads the same sidecar
    # existing_plan came from), so the existing_plan.mode term is redundant.
    effective_mode = args.mode or prior_mode or "expressive"
    final_plan.mode = effective_mode

    # Expressive-freeform composer (Option B): the no-agent floor.
    # In expressive mode, rewrite freshly-INFERRED content slides into designed
    # freeform layouts so that even raw `python build.py` (no agent in the loop)
    # emits Anthropic-aesthetic freeform slides rather than named layouts.
    #
    # Gate (preserve agent-authored sidecars): only rewrite slides whose
    # content_hash is NOT present in the existing sidecar — i.e. slides that
    # came fresh from _infer_default_plan this run, never entries that
    # merge_with_existing carried over from an agent-authored existing plan.
    # When there is no existing sidecar, every content slide is eligible.
    # Strict mode skips the composer entirely (named layouts, unchanged).
    composed_ids: list[str] = []
    if effective_mode == "expressive":
        from expressive_compose import compose_expressive_plan
        if existing_plan is not None:
            existing_hashes = {
                s.content_hash for s in existing_plan.slides if s.content_hash
            }
            only_ids = {
                s.slide_id
                for s in final_plan.slides
                if s.content_hash not in existing_hashes
            }
        else:
            only_ids = None  # no sidecar → all content slides eligible
        composed_ids = compose_expressive_plan(
            final_plan.slides, md_dir=md_path.parent, only_ids=only_ids
        )

    # Resolve theme (expressive only). Freeze the chosen theme name + seed in
    # the plan so re-renders are deterministic, and --shake rerolls.
    theme = None
    if effective_mode == "expressive":
        import uuid
        from expressive import resolve_theme
        if not args.shake and prior_theme:
            # Re-render: keep the frozen theme + seed (determinism).
            final_plan.theme = prior_theme
            final_plan.shake_seed = prior_seed
        elif args.shake:
            # Reroll: a fresh random seed selects across all themes.
            final_plan.theme = None
            final_plan.shake_seed = uuid.uuid4().hex
        else:
            # First build (no prior theme): seed from deck content so the
            # same deck reproducibly picks the same theme across machines.
            final_plan.theme = None
            final_plan.shake_seed = deck_md_hash
        theme = resolve_theme(final_plan)
        final_plan.theme = theme.name if theme else None
    else:
        final_plan.theme = None

    # --use-blocks enforcement
    _BLOCK_KINDS = {"composition", "freeform"}
    use_blocks = getattr(args, "use_blocks", "auto")
    if use_blocks == "never":
        forbidden = [
            s.slide_id
            for s in final_plan.slides
            if s.kind in _BLOCK_KINDS
        ]
        if forbidden:
            print(
                f"error: --use-blocks=never rejects composition/freeform slides: "
                f"{', '.join(forbidden)}\n"
                f"Regenerate the sidecar (delete it or run with --shake) or pass "
                f"--use-blocks=auto to allow them.",
                file=sys.stderr,
            )
            return 1

    # Persist sidecar
    sidecar_path.write_text(final_plan.to_json(), encoding="utf-8")
    print(f"wrote plan: {sidecar_path}")

    if args.plan_only:
        return 0

    # --- Bespoke-enforcement gate (cross-machine; see bespoke_design.md) ----
    # Expressive mode REQUIRES an agent in the loop to handcraft each content
    # slide's freeform geometry. `compose_expressive_plan` is only the agentless
    # FLOOR. If a real render (not --plan-only) still carries composer-floor
    # slides, abort loudly — unless the caller explicitly opted into the floor
    # with --allow-composed (cron / non-interactive synth). This is the single
    # mechanical check that makes "always bespoke" enforceable across machines;
    # prose guidance in SKILL.md / memory was repeatedly ignored.
    # Floor = any content slide NOT affirmatively stamped `_provenance=="agent"`.
    # Requiring an affirmative agent stamp (rather than only catching the
    # composer stamp) is airtight: it also blocks un-stamped named-layout
    # leftovers and freshly-inferred slides, not just composer-freeform. The
    # composer stamps "composer"; inference leaves no stamp; only a handcrafting
    # agent writes "agent". The stamp persists in the sidecar, so you cannot
    # dodge the gate by re-running build.py — it stays red until each content
    # slide is genuinely handcrafted (see bespoke_design.md).
    floor_ids = [
        s.slide_id for s in final_plan.slides
        if s.kind != "section-divider"
        and (s.params or {}).get("_provenance") != "agent"
    ]
    if effective_mode == "expressive" and floor_ids and not args.allow_composed:
        content_total = sum(
            1 for s in final_plan.slides if s.kind != "section-divider"
        )
        shown = "".join(f"     - {sid}\n" for sid in floor_ids[:12])
        if len(floor_ids) > 12:
            shown += f"     ... and {len(floor_ids) - 12} more\n"
        sys.stderr.write(
            "\n==================================================================\n"
            "  BESPOKE NOT SATISFIED — build aborted (no .pptx written)\n"
            "==================================================================\n"
            f"  {len(floor_ids)}/{content_total} content slides are NOT handcrafted —\n"
            "  agentless FLOOR (composer templates or un-stamped layouts).\n\n"
            "  Expressive mode requires you (the agent) to design each content\n"
            "  slide's freeform geometry. Recipe (see bespoke_design.md):\n"
            "    1. build.py --input X.md --output Y.pptx --plan-only --shake\n"
            "    2. For each content slide_id in the .layout.json sidecar, write\n"
            "       a bespoke params.code using the sandbox API.\n"
            "    3. Re-render WITHOUT --shake.\n\n"
            "  The sidecar has been written for you to edit:\n"
            f"    {sidecar_path}\n\n"
            "  Floor slides still needing handcraft:\n"
            f"{shown}\n"
            "  Non-interactive / cron render where the floor is acceptable?\n"
            "  Re-run with --allow-composed to bypass this gate.\n"
            "==================================================================\n"
        )
        return 2

    # --- Speaker-notes check (non-fatal) — see speaker_notes.md -------------
    # Every content slide should ship comprehensive notes in params['notes'].
    # This warns (never aborts) so gaps are visible before delivery.
    missing_notes = [
        s.slide_id for s in final_plan.slides
        if s.kind != "section-divider"
        and not str((s.params or {}).get("notes", "")).strip()
    ]
    if missing_notes:
        shown = "".join(f"     - {sid}\n" for sid in missing_notes[:12])
        if len(missing_notes) > 12:
            shown += f"     ... and {len(missing_notes) - 12} more\n"
        sys.stderr.write(
            f"\n[build-pptx] NOTE: {len(missing_notes)} content slide(s) have no "
            "speaker notes (params['notes']).\n"
            "  Decks ship with comprehensive, didactic notes by default "
            "(see speaker_notes.md):\n"
            f"{shown}"
        )

    # Render
    from render import render_from_plan
    render_from_plan(
        md_path=md_path, plan=final_plan, output_path=output_path,
        no_cover=args.no_cover, no_end=args.no_end, theme=theme,
    )
    print(f"wrote {output_path}")
    if args.qa:
        from qa import render_to_images
        qa_dir = output_path.with_suffix("").parent / (output_path.stem + "_qa")
        try:
            pngs = render_to_images(output_path, qa_dir)
            print(f"QA images ({len(pngs)}):")
            for p in pngs:
                print(f"  {p}")
        except RuntimeError as e:
            print(f"QA skipped: {e}", file=sys.stderr)

    # Contrast lint — flag freeform _add_text calls that use MUTED_RGB /
    # DIM_RGB. Those colors are appropriate only for tiny secondary captions
    # on the OPEN canvas; on a SURFACE / PAPER / accent card they produce
    # low-contrast washed-out text (the bug from the AGF 2026-05-29 deck).
    # Non-fatal — the agent decides whether each flagged call is OK.
    try:
        from contrast_lint import lint_sidecar, format_warnings
        _sidecar = Path(args.input).resolve()
        _sidecar = _sidecar.with_suffix(_sidecar.suffix + ".layout.json")
        if _sidecar.exists():
            cw = lint_sidecar(_sidecar)
            if cw:
                print(format_warnings(cw), file=sys.stderr)
    except Exception as e:
        print(f"contrast lint skipped: {e}", file=sys.stderr)

    # Presenter-handout PDF (slide image + speaker notes per page). Canonical
    # output when the deck has notes; skips gracefully without a rasterizer.
    # See speaker_notes.md. Disable with --no-notes-pdf.
    if not args.no_notes_pdf:
        try:
            from notes_pdf import build_notes_pdf
            np = build_notes_pdf(output_path)
            if np:
                print(f"wrote notes PDF: {np}")
            else:
                print("notes PDF skipped (no notes, or no LibreOffice/pdftoppm).",
                      file=sys.stderr)
        except Exception as e:  # never let the handout abort a successful build
            print(f"notes PDF skipped: {e}", file=sys.stderr)
    return 0


if __name__ == "__main__":
    sys.exit(main())
