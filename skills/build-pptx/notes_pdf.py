"""Canonical presenter-handout PDF: each slide's image + its full speaker notes.

One slide per landscape page — slide thumbnail on top (so every figure/image
on the slide is included automatically), comprehensive notes below. The title
line is word-wrapped; long notes auto-shrink fontsize down to a floor, then
overflow to a second/third page with a "(notes continued)" header instead of
clipping off the bottom or running off the right edge.

Notes are read from the rendered pptx's notes panes (populated by render.py
from params['notes']), so this stays in sync with the deck. Rasterization
reuses qa.render_to_images (LibreOffice + pdftoppm); if those are missing
it degrades gracefully (returns None) rather than failing the build.
"""
from __future__ import annotations

import tempfile
import textwrap
from pathlib import Path


# Layout (landscape letter, figsize=(11, 8.5))
# Title wrap width is measured AFTER prepending the "Slide N — " prefix, so
# this is the actual on-page line length, not just the user-text portion.
# Geist Mono (the brand title font) at 13pt is ~0.118 in/char; with 10.12 in
# of usable width (0.92 of an 11-in page), the practical ceiling is ~85 chars.
# Stay under that to leave a small right margin and survive non-mono fallbacks.
_TITLE_WIDTH_CHARS = 78     # tight; accounts for prefix + non-mono fallback
_BODY_WIDTH_CHARS = 110     # body wrap width (slightly tighter than before)
_BODY_FONT_MAX = 10.5       # starting body font (pt)
_BODY_FONT_MIN = 8.0        # floor before overflow to next page
_BODY_LINESPACING = 1.40

# Available body area: from y_top (just below title) down to a bottom margin.
_BODY_Y_TOP = 0.355
_BODY_Y_BOTTOM = 0.045
_BODY_AVAIL_FRAC = _BODY_Y_TOP - _BODY_Y_BOTTOM  # fraction of figure height


def _wrap_lines(text: str, width: int) -> list[str]:
    """Word-wrap each paragraph, preserving blank-line paragraph breaks."""
    out: list[str] = []
    for para in (text or "").split("\n"):
        if not para.strip():
            out.append("")
            continue
        out.extend(textwrap.wrap(para, width=width) or [""])
    return out


def _fit_body(body_lines: list[str], fig_height_in: float) -> tuple[float, list[list[str]]]:
    """Pick the largest fontsize that fits the body in available space; if even
    the floor fontsize would overflow, split across multiple pages.

    Returns (chosen_fontsize_pt, list_of_pages_each_a_list_of_lines).
    """
    body_avail_in = _BODY_AVAIL_FRAC * fig_height_in
    n_lines = len(body_lines)
    if n_lines == 0:
        return _BODY_FONT_MAX, [[]]

    for pt in (10.5, 10.0, 9.5, 9.0, 8.5, _BODY_FONT_MIN):
        line_h_in = pt / 72.0 * _BODY_LINESPACING
        max_lines = int(body_avail_in / line_h_in)
        if n_lines <= max_lines:
            return pt, [body_lines]

    # Even at floor we overflow — paginate at floor font.
    pt = _BODY_FONT_MIN
    line_h_in = pt / 72.0 * _BODY_LINESPACING
    max_lines = max(1, int(body_avail_in / line_h_in))
    pages: list[list[str]] = []
    i = 0
    while i < n_lines:
        pages.append(body_lines[i:i + max_lines])
        i += max_lines
    return pt, pages


def build_notes_pdf(pptx_path, out_pdf=None, *, dpi: int = 130):
    """Build `<deck>_notes.pdf`. Returns the Path, or None if it couldn't run
    (e.g. no LibreOffice/pdftoppm, or the deck has no notes)."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib.backends.backend_pdf import PdfPages
    from PIL import Image
    from pptx import Presentation

    from qa import render_to_images

    pptx_path = Path(pptx_path)
    out_pdf = Path(out_pdf) if out_pdf else pptx_path.with_name(pptx_path.stem + "_notes.pdf")

    prs = Presentation(str(pptx_path))
    notes = []
    for s in prs.slides:
        t = ""
        if s.has_notes_slide:
            t = s.notes_slide.notes_text_frame.text or ""
        notes.append(t)
    if not any(n.strip() for n in notes):
        return None  # nothing to hand out

    with tempfile.TemporaryDirectory() as td:
        try:
            pngs = render_to_images(pptx_path, Path(td), dpi=dpi)
        except RuntimeError:
            return None  # no rasterizer — skip gracefully

        # Brand fonts if available (harmless fallback otherwise).
        try:
            import sys as _sys
            _sys.path.insert(0, "/home/jiwei/arcadia/superstack/skills/_shared")
            from mpl_style import FONT_BODY, FONT_TITLE
        except Exception:
            FONT_BODY, FONT_TITLE = ["DejaVu Sans"], ["DejaVu Sans Mono"]

        n = min(len(pngs), len(notes))
        fig_w, fig_h = 11.0, 8.5  # landscape letter
        with PdfPages(str(out_pdf)) as pdf:
            for i in range(n):
                note = notes[i].strip()
                lines = note.split("\n", 1)
                raw_head = lines[0] if lines else f"Slide {i + 1}"
                body = lines[1].strip() if len(lines) > 1 else ""

                # Word-wrap the FULL title (prefix + head) so the on-page line
                # length — not just the user-text portion — stays under the limit.
                full_head = f"Slide {i + 1} — {raw_head}"
                head_lines = textwrap.wrap(full_head, width=_TITLE_WIDTH_CHARS) or [full_head]
                # Each title line shifts the body anchor down ~1.7% of fig height.
                title_offset = 0.017 * (len(head_lines) - 1)

                # Wrap body + auto-fit (shrink, then paginate at floor font).
                body_lines = _wrap_lines(body, _BODY_WIDTH_CHARS)
                body_font_pt, body_pages = _fit_body(body_lines, fig_h)

                for page_idx, page_lines in enumerate(body_pages):
                    fig = plt.figure(figsize=(fig_w, fig_h))
                    fig.patch.set_facecolor("white")

                    # Slide image — top 50% (always present, every page).
                    ax_img = fig.add_axes([0.04, 0.46, 0.92, 0.50])
                    ax_img.axis("off")
                    try:
                        ax_img.imshow(Image.open(pngs[i]))
                    except Exception:
                        pass

                    # Title — wrapped, multi-line. Pagination suffix on its own
                    # final line so it can't widen any wrapped title line.
                    title_text = "\n".join(head_lines)
                    if len(body_pages) > 1:
                        suffix = f"   ({page_idx + 1}/{len(body_pages)})"
                        if page_idx > 0:
                            suffix += "   (notes continued)"
                        title_text += "\n" + suffix.strip()
                    fig.text(0.04, 0.42, title_text, ha="left", va="top",
                             fontsize=13, fontfamily=FONT_TITLE, fontweight="bold",
                             color="#1A1A1A", linespacing=1.25)

                    # Body anchored just below the (possibly wrapped) title.
                    body_y = _BODY_Y_TOP - title_offset
                    fig.text(0.04, body_y, "\n".join(page_lines), ha="left", va="top",
                             fontsize=body_font_pt, fontfamily=FONT_BODY, color="#1A1A1A",
                             linespacing=_BODY_LINESPACING)
                    pdf.savefig(fig, facecolor="white")
                    plt.close(fig)

    return out_pdf


if __name__ == "__main__":
    import sys
    p = build_notes_pdf(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else None)
    print(f"wrote {p}" if p else "notes PDF skipped (no notes or no rasterizer)")
