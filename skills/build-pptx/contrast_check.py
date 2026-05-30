"""Runtime contrast QA — walks the rendered pptx and flags low-contrast text.

The static AST lint (contrast_lint.py) catches MUTED_RGB / DIM_RGB used as
text color in freeform code, but it can't tell whether the resulting text is
ACTUALLY low-contrast — that depends on what's behind the text, which depends
on the slide's filled shapes + the theme canvas.

This module loads the rendered pptx, walks every text run, finds the visual
background under that text (topmost filled shape that geometrically contains
the text run, falling back to the slide canvas), computes the WCAG luminance
contrast ratio, and flags pairs below threshold (4.5 normal / 3.0 large
text per WCAG AA).

Wired into build.py to fire as a non-fatal warning after rendering.
"""
from __future__ import annotations

import json
from dataclasses import dataclass
from pathlib import Path


# WCAG AA thresholds
WCAG_AA_NORMAL = 4.5
WCAG_AA_LARGE = 3.0  # large = >= 18pt regular or >= 14pt bold
LARGE_PT = 18


# Theme canvas + surface colors (mirror layouts/_shared/mpl_style + build-pptx themes.py)
THEME_CANVAS_HEX = {
    "slate": "#1E293B",
    "midnight": "#14141C",
    "forest": "#0F1E17",
    "paper": "#FFFFFF",
    "bone": "#F6F4EE",
}
THEME_SURFACE_HEX = {
    "slate": {"#27384F", "#343E4E"},      # mpl_style + build-pptx surface variants
    "midnight": {"#1F1F30", "#262638"},
    "forest": {"#152E22", "#1A3528"},
    "paper": {"#F7F7F4", "#FAFAFA"},
    "bone": {"#FFFFFF", "#F0EDE3"},
}


# Brand-4 accents — these carry section / category identity. Brand spec
# accepts them as text on theme canvas/surface and as fills under WHITE/INK
# text (_text_on canonical pairs), even when raw WCAG ratio dips below AA.
# (deeppink #FF1493 white-on-it is 3.64; blueviolet #8A2BE2 on slate is 2.46
# — both inside intentional design choices the brand spec endorses.)
BRAND_ACCENTS_HEX = {
    "#40E0D0",  # turquoise
    "#FF1493",  # deeppink
    "#F0C840",  # amber
    "#8A2BE2",  # blueviolet
}

# Chrome-style muted text colors used in footers + small captions. Brand spec
# treats these as decorative, not body — they're acceptable on the open canvas
# even at sub-4.5 contrast.
CHROME_MUTED_HEX = {
    "#94A3B8",  # MUTED_RGB (slate / midnight muted_text)
    "#888888",  # DIM_RGB (build-pptx default footer separator)
    "#86A192",  # forest muted_text
    "#9aa0a6",  # paper / bone muted_text
}


def _is_brand_approved(text_hex: str, bg_hex: str, theme: str, font_pt: float) -> bool:
    """Pairs the brand spec deliberately permits — skip the WCAG warning
    even when ratio is below AA. Captures the rendering choices a presenter
    has accepted as visually-fine despite raw WCAG.
    """
    canvas = THEME_CANVAS_HEX.get(theme, "")
    surfaces = THEME_SURFACE_HEX.get(theme, set())
    on_canvas_or_surface = (bg_hex == canvas) or (bg_hex in surfaces)

    # 1) Brand accent text on theme canvas or surface — section identity wins
    if text_hex in BRAND_ACCENTS_HEX and on_canvas_or_surface:
        return True
    # 2) WHITE / INK text on a brand-accent fill — these are the canonical
    #    _text_on() outcomes. Deeppink in particular sits on the WCAG knife's
    #    edge with white at 3.64; brand spec keeps white over deeppink anyway.
    if bg_hex in BRAND_ACCENTS_HEX:
        return True
    # 3) Chrome-muted text on canvas at small sizes — footers, separators,
    #    tiny captions. Decorative; not body text.
    if text_hex in CHROME_MUTED_HEX and on_canvas_or_surface and font_pt <= 12:
        return True
    return False


# --- color math ----------------------------------------------------------

def _hex_to_rgb(h: str) -> tuple[int, int, int]:
    h = h.lstrip("#")
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _channel_lum(c: int) -> float:
    """sRGB channel → linear luminance contribution (WCAG)."""
    s = c / 255.0
    return s / 12.92 if s <= 0.03928 else ((s + 0.055) / 1.055) ** 2.4


def luminance(rgb: tuple[int, int, int]) -> float:
    r, g, b = rgb
    return 0.2126 * _channel_lum(r) + 0.7152 * _channel_lum(g) + 0.0722 * _channel_lum(b)


def contrast_ratio(rgb1, rgb2) -> float:
    L1, L2 = luminance(rgb1), luminance(rgb2)
    lighter, darker = max(L1, L2), min(L1, L2)
    return (lighter + 0.05) / (darker + 0.05)


# --- shape introspection -------------------------------------------------

def _rgb_from_pptx_fill(shape) -> tuple[int, int, int] | None:
    """Try to extract a solid-fill RGB from a shape. Returns None if no
    solid fill (theme color, scheme color, no-fill, gradient, etc.).
    """
    try:
        fill = shape.fill
    except Exception:
        return None
    try:
        if fill.type != 1:  # MSO_FILL.SOLID = 1
            return None
        fc = fill.fore_color
        rgb = fc.rgb  # RGBColor
        return (rgb[0], rgb[1], rgb[2])
    except Exception:
        return None


def _rgb_from_run_color(run) -> tuple[int, int, int] | None:
    try:
        rgb = run.font.color.rgb
        if rgb is None:
            return None
        return (rgb[0], rgb[1], rgb[2])
    except Exception:
        return None


def _font_pt(run) -> float:
    try:
        sz = run.font.size
        if sz is None:
            return 14.0
        return sz.pt
    except Exception:
        return 14.0


def _bold(run) -> bool:
    try:
        return bool(run.font.bold)
    except Exception:
        return False


def _bbox(shape) -> tuple[int, int, int, int]:
    """(left, top, width, height) in EMU; falls back to (0,0,0,0)."""
    try:
        return (shape.left or 0, shape.top or 0, shape.width or 0, shape.height or 0)
    except Exception:
        return (0, 0, 0, 0)


def _contains(outer, inner) -> bool:
    """Does the `outer` bbox enclose the center of `inner`?"""
    ol, ot, ow, oh = outer
    il, it, iw, ih = inner
    cx = il + iw // 2
    cy = it + ih // 2
    return (ol <= cx <= ol + ow) and (ot <= cy <= ot + oh)


# --- main check -----------------------------------------------------------

@dataclass
class ContrastIssue:
    slide_index: int
    text_excerpt: str
    text_rgb: tuple[int, int, int]
    bg_rgb: tuple[int, int, int]
    bg_source: str  # "fill" or "canvas"
    ratio: float
    threshold: float
    font_pt: float


def _canvas_rgb(sidecar_path: Path) -> tuple[int, int, int]:
    """Read the resolved theme from the sidecar and return its canvas RGB."""
    try:
        plan = json.loads(sidecar_path.read_text())
        theme = plan.get("theme", "paper")
        hex_color = THEME_CANVAS_HEX.get(theme, "#FFFFFF")
        return _hex_to_rgb(hex_color)
    except Exception:
        return (255, 255, 255)


def _theme_name(sidecar_path: Path) -> str:
    try:
        plan = json.loads(sidecar_path.read_text())
        return plan.get("theme", "paper")
    except Exception:
        return "paper"


def check_pptx(pptx_path, sidecar_path=None) -> list[ContrastIssue]:
    """Walk the pptx; return contrast issues below WCAG AA."""
    from pptx import Presentation
    pptx_path = Path(pptx_path)
    sidecar_path = sidecar_path or pptx_path.with_suffix(".md.layout.json")
    if not Path(sidecar_path).exists():
        # If sidecar isn't where expected, try the .md alongside
        candidate = pptx_path.with_suffix("").with_suffix(".md.layout.json")
        if candidate.exists():
            sidecar_path = candidate
    canvas = _canvas_rgb(Path(sidecar_path)) if Path(sidecar_path).exists() else (255, 255, 255)
    theme = _theme_name(Path(sidecar_path)) if Path(sidecar_path).exists() else "paper"

    prs = Presentation(str(pptx_path))
    issues: list[ContrastIssue] = []

    for i, slide in enumerate(prs.slides, 1):
        # Collect filled rects in z-order (later = topmost in pptx)
        filled = []  # list of (bbox, rgb, shape)
        text_shapes = []
        for shape in slide.shapes:
            rgb = _rgb_from_pptx_fill(shape)
            bbox = _bbox(shape)
            if rgb is not None and bbox[2] > 0 and bbox[3] > 0:
                filled.append((bbox, rgb, shape))
            if getattr(shape, "has_text_frame", False):
                text_shapes.append(shape)

        for tshape in text_shapes:
            tbbox = _bbox(tshape)
            # Walk z-order from topmost downward to find enclosing filled rect
            bg_rgb = None
            bg_source = "canvas"
            for fbbox, frgb, fshape in reversed(filled):
                # Skip self-fills (text shape with its own fill — that's the text bg)
                if fshape is tshape:
                    bg_rgb = frgb
                    bg_source = "self-fill"
                    break
                if _contains(fbbox, tbbox):
                    bg_rgb = frgb
                    bg_source = "fill"
                    break
            if bg_rgb is None:
                bg_rgb = canvas

            # Inspect every run independently
            try:
                paras = tshape.text_frame.paragraphs
            except Exception:
                continue
            for para in paras:
                for run in para.runs:
                    text_rgb = _rgb_from_run_color(run)
                    if text_rgb is None:
                        continue
                    pt = _font_pt(run)
                    is_large = pt >= LARGE_PT or (pt >= 14 and _bold(run))
                    threshold = WCAG_AA_LARGE if is_large else WCAG_AA_NORMAL
                    ratio = contrast_ratio(text_rgb, bg_rgb)
                    if ratio < threshold:
                        # Brand-tolerance exemption — skip warning for
                        # pairs the brand spec deliberately permits.
                        text_hex = "#{:02X}{:02X}{:02X}".format(*text_rgb)
                        bg_hex = "#{:02X}{:02X}{:02X}".format(*bg_rgb)
                        if _is_brand_approved(text_hex, bg_hex, theme, pt):
                            continue
                        excerpt = (run.text or "").strip()[:60]
                        if not excerpt:
                            continue
                        issues.append(ContrastIssue(
                            slide_index=i, text_excerpt=excerpt,
                            text_rgb=text_rgb, bg_rgb=bg_rgb, bg_source=bg_source,
                            ratio=ratio, threshold=threshold, font_pt=pt,
                        ))
    return issues


def format_issues(issues: list[ContrastIssue]) -> str:
    if not issues:
        return ""
    lines = [
        f"contrast check: {len(issues)} text run(s) below WCAG AA",
        "  text rgb on background rgb (from fill / canvas) — ratio vs threshold",
        "",
    ]
    for w in issues:
        tr = "#{:02X}{:02X}{:02X}".format(*w.text_rgb)
        br = "#{:02X}{:02X}{:02X}".format(*w.bg_rgb)
        lines.append(
            f"  slide {w.slide_index:>2}: {tr} on {br} ({w.bg_source}) "
            f"= {w.ratio:.2f} < {w.threshold} @ {w.font_pt:.0f}pt"
        )
        lines.append(f"      \"{w.text_excerpt}\"")
    return "\n".join(lines)


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("usage: contrast_check.py <deck.pptx> [<deck.md.layout.json>]", file=sys.stderr)
        sys.exit(2)
    pptx = sys.argv[1]
    sidecar = sys.argv[2] if len(sys.argv) > 2 else None
    issues = check_pptx(pptx, sidecar)
    if issues:
        print(format_issues(issues), file=sys.stderr)
    else:
        print("contrast check: clean")
