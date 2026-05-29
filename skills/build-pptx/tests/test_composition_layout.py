"""Round-trip tests for the composition layout.

Loads fixture_composition.md + its pre-built sidecar, renders through
build.py, and asserts basic structural expectations.
"""
from __future__ import annotations

import json
import shutil
import subprocess
import sys
from pathlib import Path

import pytest

SKILL_DIR = Path(__file__).resolve().parents[1]
BUILD_PY = SKILL_DIR / "build.py"
FIXTURE_MD = SKILL_DIR / "tests" / "fixture_composition.md"
FIXTURE_SIDECAR = SKILL_DIR / "tests" / "fixture_composition.md.layout.json"


def _run_build(md_path: Path, out_path: Path, extra: list[str] | None = None) -> subprocess.CompletedProcess:
    cmd = [
        sys.executable, str(BUILD_PY), "--allow-composed",
        "--input", str(md_path),
        "--output", str(out_path),
    ] + (extra or [])
    return subprocess.run(cmd, capture_output=True, text=True)


# ---------------------------------------------------------------------------
# Basic round-trip
# ---------------------------------------------------------------------------

def _render_via_plan(md_path: Path, sidecar_path: Path, out_path: Path) -> None:
    """Helper: render using render_from_plan directly (bypasses merge logic)."""
    sys.path.insert(0, str(SKILL_DIR))
    sys.path.insert(0, str(SKILL_DIR.parent / "_shared"))
    from plan import Plan
    from render import render_from_plan
    plan = Plan.from_json(sidecar_path.read_text())
    render_from_plan(md_path=md_path, plan=plan,
                     output_path=out_path, no_cover=False, no_end=False)


def test_composition_renders_pptx(tmp_path):
    """Fixture with composition slides renders to a valid pptx."""
    md = tmp_path / "deck.md"
    sidecar = tmp_path / "deck.md.layout.json"
    md.write_text(FIXTURE_MD.read_text())
    sidecar.write_text(FIXTURE_SIDECAR.read_text())
    out = tmp_path / "out.pptx"

    _render_via_plan(md, sidecar, out)
    assert out.exists()
    assert out.stat().st_size > 5_000


def test_composition_slide_count(tmp_path):
    """Slide count: title + 4 plan slides (1 content-text + 3 composition) + end = 6."""
    md = tmp_path / "deck.md"
    sidecar = tmp_path / "deck.md.layout.json"
    md.write_text(FIXTURE_MD.read_text())
    sidecar.write_text(FIXTURE_SIDECAR.read_text())
    out = tmp_path / "out.pptx"

    _render_via_plan(md, sidecar, out)

    from pptx import Presentation
    prs = Presentation(str(out))
    # title + 4 plan slides + end = 6
    assert len(prs.slides) == 6


def test_composition_slide_has_shapes(tmp_path):
    """Composition slides contain rendered shapes."""
    md = tmp_path / "deck.md"
    sidecar = tmp_path / "deck.md.layout.json"
    md.write_text(FIXTURE_MD.read_text())
    sidecar.write_text(FIXTURE_SIDECAR.read_text())
    out = tmp_path / "out.pptx"

    _render_via_plan(md, sidecar, out)

    from pptx import Presentation
    prs = Presentation(str(out))
    # Slide index 2 is first composition slide (results-overview)
    composition_slide = prs.slides[2]  # title=0, intro=1, results-overview=2
    assert len(composition_slide.shapes) >= 4


def test_composition_deterministic(tmp_path):
    """Two renders of the same plan produce the same slide count."""
    md = tmp_path / "deck.md"
    sidecar = tmp_path / "deck.md.layout.json"
    md.write_text(FIXTURE_MD.read_text())
    sidecar.write_text(FIXTURE_SIDECAR.read_text())

    out1 = tmp_path / "out1.pptx"
    out2 = tmp_path / "out2.pptx"

    for out in (out1, out2):
        _render_via_plan(md, sidecar, out)

    from pptx import Presentation
    prs1 = Presentation(str(out1))
    prs2 = Presentation(str(out2))
    assert len(prs1.slides) == len(prs2.slides)


def test_composition_via_build_py_does_not_crash(tmp_path):
    """build.py with composition sidecar renders without crashing (end-to-end)."""
    md = tmp_path / "deck.md"
    sidecar = tmp_path / "deck.md.layout.json"
    md.write_text(FIXTURE_MD.read_text())
    sidecar.write_text(FIXTURE_SIDECAR.read_text())
    out = tmp_path / "out.pptx"

    proc = _run_build(md, out)
    assert proc.returncode == 0, proc.stderr
    assert out.exists()
    assert out.stat().st_size > 5_000


# ---------------------------------------------------------------------------
# Catalog registration
# ---------------------------------------------------------------------------

def test_composition_registered_in_catalog():
    """composition is in the layout catalog REGISTRY."""
    sys.path.insert(0, str(SKILL_DIR))
    from layouts import catalog
    assert "composition" in catalog.REGISTRY


def test_composition_renderer_callable():
    """catalog.get('composition') returns a callable."""
    sys.path.insert(0, str(SKILL_DIR))
    from layouts import catalog
    renderer = catalog.get("composition")
    assert callable(renderer)


# ---------------------------------------------------------------------------
# Direct unit tests for block dispatch
# ---------------------------------------------------------------------------

def _make_slide():
    """Create a blank slide for testing."""
    from build import new_presentation, _blank
    prs = new_presentation()
    return _blank(prs)


def test_all_block_kinds_render_without_error():
    """All 8 block kinds can be dispatched without raising."""
    sys.path.insert(0, str(SKILL_DIR))
    sys.path.insert(0, str(SKILL_DIR.parent / "_shared"))

    from layouts.blocks import BLOCKS
    from layouts._common import TURQUOISE_RGB

    block_fixtures = {
        "paragraph": {"items": ["Hello world.", "Second item."], "size": 14},
        "figure": {"image_path": "/nonexistent/image.png", "caption": "A figure"},
        "card-row": {"cards": [{"label": "A", "body": "body a"},
                               {"label": "B", "body": "body b"}]},
        "stat-tile": {"value": "0.91", "label": "AUC", "sub": "mean"},
        "accent-callout": {"text": "This is a callout.", "tone": "dark"},
        "table": {"rows": [["A", "B"], ["1", "2"]]},
        "quote": {"text": "A wise saying.", "attribution": "Someone"},
        "left-accent-card": {"label": "Category", "body": "Description text."},
    }

    for kind, params in block_fixtures.items():
        slide = _make_slide()
        renderer = BLOCKS.get(kind)
        assert renderer is not None, f"Missing block kind: {kind}"
        # Should not raise
        renderer(slide, left=0.5, top=1.5, width=5.0, height=2.0,
                 params=params, accent_rgb=TURQUOISE_RGB)


def test_composition_per_row_accent_override():
    """Per-row accent_hex override is applied without error."""
    sys.path.insert(0, str(SKILL_DIR))
    sys.path.insert(0, str(SKILL_DIR.parent / "_shared"))

    from build import new_presentation, _blank
    from layouts.composition import render
    from layouts._common import TURQUOISE_RGB

    prs = new_presentation()
    slide = _blank(prs)

    params = {
        "title": "Accent Override Test",
        "lede": "",
        "rows": [
            {
                "weight": 1,
                "accent_hex": "#FF1493",
                "blocks": [
                    {"kind": "stat-tile", "weight": 1,
                     "params": {"value": "42", "label": "count"}}
                ]
            },
            {
                "weight": 1,
                "blocks": [
                    {"kind": "paragraph", "weight": 1,
                     "params": {"items": ["Some text."]}}
                ]
            }
        ]
    }

    # Should not raise
    render(slide, params=params, accent_rgb=TURQUOISE_RGB,
           footer_kwargs={"name": "Jin", "org": "UCSF",
                          "deck_title": "Test", "date": "2026-05-02"})
    assert len(slide.shapes) > 0


def test_composition_fixed_height_row():
    """Rows with explicit 'height' (inches) are honoured without error."""
    sys.path.insert(0, str(SKILL_DIR))
    sys.path.insert(0, str(SKILL_DIR.parent / "_shared"))

    from build import new_presentation, _blank
    from layouts.composition import render
    from layouts._common import TURQUOISE_RGB

    prs = new_presentation()
    slide = _blank(prs)

    params = {
        "title": "Fixed Height Row",
        "lede": "",
        "rows": [
            {
                "height": 1.5,
                "blocks": [
                    {"kind": "accent-callout", "weight": 1,
                     "params": {"text": "Fixed height callout.", "tone": "accent"}}
                ]
            },
            {
                "blocks": [
                    {"kind": "paragraph", "weight": 1,
                     "params": {"items": ["Auto-height paragraph."]}}
                ]
            }
        ]
    }

    render(slide, params=params, accent_rgb=TURQUOISE_RGB,
           footer_kwargs={})
    assert len(slide.shapes) > 0
