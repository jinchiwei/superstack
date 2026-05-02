"""Security tests for the freeform-layout sandbox.

These tests are the load-bearing piece — if they fail, freeform code can
escape the namespace. A bug here is a security bug. Never relax these
tests; only add to them."""

import sys
from pathlib import Path
import pytest

SKILL_DIR = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(SKILL_DIR))

from layouts._sandbox import validate_code, run, SandboxError


# === validate_code accepts safe code ===

def test_validates_simple_drawing():
    validate_code("_add_rect(slide, left=0, top=0, width=1, height=1, fill_rgb=accent_rgb)")

def test_validates_loop():
    code = """
for i in range(3):
    _add_text(slide, str(i), left=i, top=1, width=1, height=1, size=14, color_rgb=INK_RGB, font=MONO_FONT)
"""
    validate_code(code)

def test_validates_function_def():
    code = """
def helper(label, x):
    _add_text(slide, label, left=x, top=1, width=2, height=1, size=14, color_rgb=INK_RGB, font=MONO_FONT)
helper("foo", 0)
helper("bar", 3)
"""
    validate_code(code)


# === validate_code rejects dangerous patterns ===

def test_rejects_import_statement():
    with pytest.raises(SandboxError, match="import"):
        validate_code("import os\n_add_rect(slide, left=0, top=0, width=1, height=1, fill_rgb=accent_rgb)")

def test_rejects_from_import():
    with pytest.raises(SandboxError, match="import"):
        validate_code("from os import path\n_add_rect(slide, left=0, top=0, width=1, height=1, fill_rgb=accent_rgb)")

def test_rejects_open_call():
    with pytest.raises(SandboxError):
        validate_code("open('/etc/passwd').read()")

def test_rejects_eval_call():
    with pytest.raises(SandboxError):
        validate_code("eval('1+1')")

def test_rejects_exec_call():
    with pytest.raises(SandboxError):
        validate_code("exec('print(1)')")

def test_rejects_dunder_attribute_class():
    with pytest.raises(SandboxError, match="dunder|__class__"):
        validate_code("x = slide.__class__")

def test_rejects_dunder_attribute_subclasses():
    with pytest.raises(SandboxError):
        validate_code("''.__class__.__bases__[0].__subclasses__()")

def test_rejects_dunder_globals():
    with pytest.raises(SandboxError):
        validate_code("def f():\n    return f.__globals__")

def test_rejects_getattr():
    with pytest.raises(SandboxError):
        validate_code("getattr(slide, '_blob')")

def test_rejects_setattr():
    with pytest.raises(SandboxError):
        validate_code("setattr(slide, 'foo', 1)")

def test_rejects_globals_call():
    with pytest.raises(SandboxError):
        validate_code("globals()")

def test_rejects_try_statement():
    with pytest.raises(SandboxError, match="try|with"):
        validate_code("try:\n    _add_text(slide, 'x', left=0, top=0, width=1, height=1, size=14, color_rgb=INK_RGB, font=MONO_FONT)\nexcept Exception:\n    pass")

def test_rejects_with_statement():
    with pytest.raises(SandboxError, match="try|with"):
        validate_code("with open('/tmp/x') as f:\n    pass")

def test_rejects_unparseable_code():
    with pytest.raises(SandboxError, match="parse|syntax"):
        validate_code("def foo(:\n    pass")


# === build_safe_globals exposes the right names ===

def test_safe_globals_includes_drawing_primitives():
    from pptx import Presentation
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    from layouts._sandbox import build_safe_globals
    g = build_safe_globals(slide=s, accent_hex="#40E0D0",
                           body_top=1.5, body_h=5.0,
                           body_l=0.5, body_w=12.3, body_bottom=6.5)
    for name in ("_add_rect", "_add_text", "_add_card",
                 "INK_RGB", "WHITE_RGB", "TURQUOISE_RGB",
                 "Inches", "Pt", "RGBColor",
                 "MSO_SHAPE", "MONO_FONT", "SANS_FONT",
                 "slide", "accent_rgb",
                 "body_top", "body_h", "body_l", "body_w", "body_bottom"):
        assert name in g, f"missing {name} from safe globals"

def test_safe_globals_excludes_dangerous_builtins():
    from pptx import Presentation
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    from layouts._sandbox import build_safe_globals
    g = build_safe_globals(slide=s, accent_hex="#40E0D0",
                           body_top=1.5, body_h=5.0,
                           body_l=0.5, body_w=12.3, body_bottom=6.5)
    builtins = g["__builtins__"]
    for forbidden in ("open", "__import__", "exec", "eval",
                      "compile", "getattr", "setattr",
                      "globals", "locals", "vars", "input"):
        assert forbidden not in builtins, f"forbidden builtin {forbidden} leaked"


# === run() integration ===

def test_run_executes_safe_code():
    from pptx import Presentation
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    code = """
_add_rect(slide, left=1, top=1, width=2, height=2, fill_rgb=accent_rgb)
_add_text(slide, "test", left=1, top=1, width=2, height=1, size=14, color_rgb=INK_RGB, font=MONO_FONT)
"""
    run(code=code, slide=s, accent_hex="#40E0D0",
        body_top=1.5, body_h=5.0, body_l=0.5, body_w=12.3, body_bottom=6.5)
    # Two shapes should have been added
    n_shapes = sum(1 for _ in s.shapes)
    assert n_shapes >= 2

def test_run_blocks_filesystem_attempt():
    from pptx import Presentation
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    with pytest.raises(SandboxError):
        run(code="open('/etc/passwd').read()", slide=s, accent_hex="#40E0D0",
            body_top=1.5, body_h=5.0, body_l=0.5, body_w=12.3, body_bottom=6.5)

def test_run_blocks_import_attempt():
    from pptx import Presentation
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    with pytest.raises(SandboxError):
        run(code="import subprocess\nsubprocess.call(['ls'])",
            slide=s, accent_hex="#40E0D0",
            body_top=1.5, body_h=5.0, body_l=0.5, body_w=12.3, body_bottom=6.5)

def test_run_blocks_dunder_escape_attempt():
    """Classic sandbox escape: ''.__class__.__bases__[0].__subclasses__() to find os."""
    from pptx import Presentation
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    with pytest.raises(SandboxError):
        run(code="x = ''.__class__.__mro__[1].__subclasses__()",
            slide=s, accent_hex="#40E0D0",
            body_top=1.5, body_h=5.0, body_l=0.5, body_w=12.3, body_bottom=6.5)
