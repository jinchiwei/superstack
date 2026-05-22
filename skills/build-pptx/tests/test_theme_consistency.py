"""Theme-consistency regression tests for the remaining hardcoded-color holdouts.

Under a dark theme, conclusions cards must use the theme's card surface
(#2B2B32 for midnight) rather than the legacy hardcoded navy (#1A2D50). Under
no theme (strict/LIGHT) the cards must stay #1A2D50 for byte-parity.
"""
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from pptx import Presentation

import render as render_mod
from plan import Plan, SlideEntry
from themes import get_theme


def _all_shape_fill_hexes(pptx_path):
    """Every solid-fill shape hex across all slides."""
    prs = Presentation(str(pptx_path))
    out = []
    for s in prs.slides:
        for sh in s.shapes:
            try:
                if int(sh.fill.type) == 1:
                    out.append(str(sh.fill.fore_color.rgb))
            except Exception:
                pass
    return out


_CONCLUSIONS_PARAMS = {
    "title": "Takeaways",
    "lede": "what we found",
    "section_label": "Takeaways",
    "cards": [
        {"label": "Headline AUC", "body": "0.91 internal", "icon": "FaChartLine"},
        {"label": "Survivors", "body": "held out", "icon": "FaCheckCircle"},
    ],
}


def test_conclusions_cards_use_theme_surface_on_dark(tmp_path):
    """Midnight theme: conclusions cards fill with theme surface #2B2B32, not navy."""
    md = tmp_path / "deck.md"
    md.write_text("---\ntitle: T\n---\n\n# A\n\nbody\n", encoding="utf-8")
    plan = Plan(mode="expressive", theme="midnight", slides=[
        SlideEntry(slide_id="h1-a", kind="conclusions", params=_CONCLUSIONS_PARAMS),
    ])
    out = tmp_path / "out.pptx"
    render_mod.render_from_plan(md_path=md, plan=plan, output_path=out,
                                theme=get_theme("midnight"))
    fills = _all_shape_fill_hexes(out)
    assert "2B2B32" in fills, f"expected theme surface 2B2B32 in card fills, got {fills}"
    assert "1A2D50" not in fills, f"legacy navy 1A2D50 should be gone, got {fills}"


def test_conclusions_cards_stay_navy_under_strict(tmp_path):
    """No theme (strict/LIGHT): conclusions cards keep legacy navy #1A2D50 (parity)."""
    md = tmp_path / "deck.md"
    md.write_text("---\ntitle: T\n---\n\n# A\n\nbody\n", encoding="utf-8")
    plan = Plan(mode="strict", slides=[
        SlideEntry(slide_id="h1-a", kind="conclusions", params=_CONCLUSIONS_PARAMS),
    ])
    out = tmp_path / "out.pptx"
    render_mod.render_from_plan(md_path=md, plan=plan, output_path=out, theme=None)
    fills = _all_shape_fill_hexes(out)
    assert "1A2D50" in fills, f"strict conclusions cards should stay navy 1A2D50, got {fills}"
    assert "2B2B32" not in fills, f"theme surface 2B2B32 must not appear under strict, got {fills}"
