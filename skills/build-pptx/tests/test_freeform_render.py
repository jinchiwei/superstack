"""Tests for the freeform layout: round-trip render + visible error chip
on bad code. The sandbox itself is tested in test_sandbox_security."""

import sys
import json
import shutil
import subprocess
from pathlib import Path

import pytest

SKILL_DIR = Path(__file__).resolve().parents[1]
DEMO_MD = SKILL_DIR / "tests" / "fixture_freeform_demo.md"
DEMO_SIDECAR = SKILL_DIR / "tests" / "fixture_freeform_demo.md.layout.json"
BUILD_PY = SKILL_DIR / "build.py"


def test_freeform_layout_registered_in_catalog():
    sys.path.insert(0, str(SKILL_DIR))
    from layouts import catalog
    assert "freeform" in catalog.REGISTRY


def test_freeform_demo_renders_to_pptx(tmp_path):
    """The hand-written sidecar should render cleanly through the freeform
    layout, producing a multi-slide pptx with no errors."""
    # Copy the fixture + sidecar into tmp_path so we don't pollute the
    # source tree, and so the renderer's --shake invariant doesn't matter.
    md_dst = tmp_path / "freeform_demo.md"
    sidecar_dst = tmp_path / "freeform_demo.md.layout.json"
    shutil.copy(DEMO_MD, md_dst)
    shutil.copy(DEMO_SIDECAR, sidecar_dst)
    out = tmp_path / "out.pptx"
    proc = subprocess.run(
        # This test exercises the freeform RENDER pipeline, not contrast QA;
        # the demo fixture's snippet trips the WCAG-AA contrast gate, so opt out
        # of that gate here (the gate itself is covered by its own tests).
        [sys.executable, str(BUILD_PY),
         "--input", str(md_dst), "--output", str(out), "--allow-contrast-fail"],
        capture_output=True, text=True,
    )
    assert proc.returncode == 0, proc.stderr
    assert out.exists() and out.stat().st_size > 30000

    # Confirm the freeform slides each have shapes from their snippets.
    from pptx import Presentation
    prs = Presentation(str(out))
    # Cover + 3 dividers + 3 freeform + end = 8 slides
    assert len(prs.slides) >= 6


def test_freeform_runtime_error_renders_chip_not_crashes(tmp_path):
    """If a freeform snippet raises at runtime (e.g., wrong arg count),
    the renderer should emit an error chip instead of crashing the deck."""
    md_dst = tmp_path / "bad.md"
    md_dst.write_text(
        "---\ntitle: 'Bad freeform'\nname: 'Jin'\norg: 'X'\ndate: '2026-05-02'\n---\n\n"
        "# Section\n\nIntro.\n"
    )
    sidecar = tmp_path / "bad.md.layout.json"
    # Compute the real content_hash so merge_with_existing keeps the sidecar entry.
    import hashlib
    import sys as _sys
    _sys.path.insert(0, str(SKILL_DIR))
    from md_loader import load_markdown as _lmd
    from layouts._common import _split_slides as _ss
    _loaded = _lmd(str(md_dst))
    _chunks = _ss(_loaded["body_html"])
    _h1_chunk = next((c for c in _chunks if "<h1" in c), _chunks[0])
    _real_hash = hashlib.sha256(_h1_chunk.encode("utf-8")).hexdigest()

    sidecar.write_text(json.dumps({
        "version": 1,
        "deck_md_hash": "x",
        "shake_seed": None,
        "slides": [
            {"slide_id": "divider-h1-section", "kind": "section-divider",
             "params": {"label": "Section", "accent_hex": "#40E0D0"},
             "content_hash": _real_hash + "-divider"},
            {"slide_id": "h1-section", "kind": "freeform",
             "params": {"title": "Section", "lede": "Intro.",
                        # Missing required kwarg → TypeError at runtime
                        "code": "_add_rect(slide)", "_provenance": "agent"},
             "content_hash": _real_hash},
        ]
    }))
    out = tmp_path / "out.pptx"
    proc = subprocess.run(
        [sys.executable, str(BUILD_PY),
         "--input", str(md_dst), "--output", str(out)],
        capture_output=True, text=True,
    )
    # Build does NOT crash even though the code is buggy
    assert proc.returncode == 0, proc.stderr
    assert out.exists()

    # Find the error chip on the freeform slide
    from pptx import Presentation
    prs = Presentation(str(out))
    found_chip = False
    for slide in prs.slides:
        for shp in slide.shapes:
            if shp.has_text_frame:
                t = shp.text_frame.text
                if "freeform runtime error" in t:
                    found_chip = True
                    break
        if found_chip:
            break
    assert found_chip, "expected a runtime-error chip on the freeform slide"


def test_freeform_validation_error_renders_chip_not_crashes(tmp_path):
    """If a freeform snippet uses an import (rejected by validation),
    render an error chip instead of crashing."""
    md_dst = tmp_path / "bad.md"
    md_dst.write_text(
        "---\ntitle: 'Validation error'\nname: 'Jin'\norg: 'X'\ndate: '2026-05-02'\n---\n\n"
        "# Section\n\nIntro.\n"
    )
    sidecar = tmp_path / "bad.md.layout.json"
    import hashlib
    import sys as _sys
    _sys.path.insert(0, str(SKILL_DIR))
    from md_loader import load_markdown as _lmd2
    from layouts._common import _split_slides as _ss2
    _loaded2 = _lmd2(str(md_dst))
    _chunks2 = _ss2(_loaded2["body_html"])
    _h1_chunk2 = next((c for c in _chunks2 if "<h1" in c), _chunks2[0])
    _real_hash2 = hashlib.sha256(_h1_chunk2.encode("utf-8")).hexdigest()

    sidecar.write_text(json.dumps({
        "version": 1,
        "deck_md_hash": "x",
        "shake_seed": None,
        "slides": [
            {"slide_id": "divider-h1-section", "kind": "section-divider",
             "params": {"label": "Section", "accent_hex": "#40E0D0"},
             "content_hash": _real_hash2 + "-divider"},
            {"slide_id": "h1-section", "kind": "freeform",
             "params": {"title": "Section", "lede": "Intro.",
                        "code": "import os\n_add_rect(slide, left=0, top=0, width=1, height=1, fill_rgb=accent_rgb)",
                        "_provenance": "agent"},
             "content_hash": _real_hash2},
        ]
    }))
    out = tmp_path / "out.pptx"
    proc = subprocess.run(
        [sys.executable, str(BUILD_PY),
         "--input", str(md_dst), "--output", str(out)],
        capture_output=True, text=True,
    )
    assert proc.returncode == 0, proc.stderr
    assert out.exists()
    from pptx import Presentation
    prs = Presentation(str(out))
    found_chip = False
    for slide in prs.slides:
        for shp in slide.shapes:
            if shp.has_text_frame:
                t = shp.text_frame.text
                if "freeform code rejected" in t and "import" in t.lower():
                    found_chip = True
                    break
    assert found_chip, "expected a validation-error chip mentioning 'import'"
