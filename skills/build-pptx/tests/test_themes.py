import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from themes import THEMES, Theme, get_theme, pick_theme

BRAND4 = {"#40E0D0", "#FF1493", "#F0C840", "#8A2BE2"}


def test_registry_nonempty_and_well_formed():
    assert len(THEMES) >= 4
    for name, t in THEMES.items():
        assert isinstance(t, Theme)
        assert t.name == name
        assert t.bg_hex.startswith("#") and len(t.bg_hex) == 7
        assert set(t.accent_order) <= BRAND4
        assert len(t.accent_order) == 4
        for h in t.supplementary:
            assert h.startswith("#") and len(h) == 7


def test_pick_theme_is_deterministic_by_seed():
    a = pick_theme("seed-123")
    b = pick_theme("seed-123")
    assert a.name == b.name


def test_pick_theme_none_seed_returns_a_theme():
    assert isinstance(pick_theme(None), Theme)


def test_pick_theme_reaches_multiple_themes():
    import uuid
    names = {pick_theme(uuid.uuid4().hex).name for _ in range(200)}
    assert len(names) >= 2


def test_get_theme_by_name():
    name = next(iter(THEMES))
    assert get_theme(name).name == name


def test_get_theme_unknown_falls_back_gracefully():
    t = get_theme("does-not-exist")
    assert isinstance(t, Theme)


def test_resolve_theme_strict_is_none():
    from plan import Plan
    from expressive import resolve_theme
    assert resolve_theme(Plan(mode="strict")) is None


def test_resolve_theme_expressive_picks_and_is_stable():
    from plan import Plan
    from expressive import resolve_theme
    p = Plan(mode="expressive", shake_seed="abc")
    t1 = resolve_theme(p)
    t2 = resolve_theme(p)
    assert t1 is not None and t1.name == t2.name


def test_resolve_theme_honors_frozen_name():
    from plan import Plan
    from expressive import resolve_theme
    assert resolve_theme(Plan(mode="expressive", theme="forest")).name == "forest"


def _render_freeform_and_find_canvas(tmp_path, theme_name: str, expected_bg_hex: str):
    """Render a freeform slide under the given theme, open the output pptx, and
    return True iff a full-bleed (~13.33in x 7.5in) shape with a solid fill of
    expected_bg_hex (e.g. '14141C') is found on the content slide."""
    import render as render_mod
    from plan import Plan, SlideEntry
    from themes import get_theme
    from pptx import Presentation
    from pptx.util import Inches

    md = tmp_path / "deck.md"
    md.write_text(
        "---\ntitle: T\n---\n\n# Results\n\n## Headline\n\nLede.\n",
        encoding="utf-8",
    )
    code = ("_add_rect(slide, left=body_l, top=body_top, width=4, height=2, "
            "fill_rgb=THEME_RGBS[0] if THEME_RGBS else accent_rgb)\n"
            "_add_text(slide, 'hi', left=body_l, top=body_top, width=4, "
            "height=1, size=20, color_rgb=WHITE_RGB if ON_DARK else INK_RGB, "
            "font=MONO_FONT)")
    plan = Plan(mode="expressive", theme=theme_name, slides=[
        SlideEntry(slide_id="h1-results/h2-headline", kind="freeform",
                   params={"title": "Headline", "lede": "Lede.",
                           "section_label": "Results", "code": code}),
    ])
    out = tmp_path / "out.pptx"
    render_mod.render_from_plan(
        md_path=md, plan=plan, output_path=out,
        theme=get_theme(theme_name),
    )
    assert out.exists() and out.stat().st_size > 0

    prs = Presentation(str(out))
    # Content slide is the freeform one (title cover is first, end last).
    full_w, full_h = 13.333, 7.5
    tol = Inches(0.1)
    target = expected_bg_hex.upper().lstrip("#")
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if int(shape.fill.type) != 1:  # MSO_FILL.SOLID
                    continue
                rgb = str(shape.fill.fore_color.rgb).upper()
            except (TypeError, ValueError, AttributeError):
                continue
            if rgb != target:
                continue
            if (abs(shape.width - Inches(full_w)) <= tol
                    and abs(shape.height - Inches(full_h)) <= tol):
                return True
    return False


def test_freeform_dark_canvas_is_painted(tmp_path):
    """A freeform slide under a DARK theme (midnight, #14141C) paints a
    full-bleed canvas rect filled with the theme bg."""
    assert _render_freeform_and_find_canvas(tmp_path, "midnight", "14141C"), (
        "expected a full-bleed 13.33x7.5in rect filled #14141C on the slide"
    )


def test_freeform_tinted_canvas_is_painted(tmp_path):
    """A freeform slide under a TINTED light theme (bone, #F6F4EE) paints a
    full-bleed canvas rect filled with the theme bg. This FAILS before Fix 1
    (paint was gated on on_dark) and PASSES after."""
    assert _render_freeform_and_find_canvas(tmp_path, "bone", "F6F4EE"), (
        "expected a full-bleed 13.33x7.5in rect filled #F6F4EE on the slide"
    )
