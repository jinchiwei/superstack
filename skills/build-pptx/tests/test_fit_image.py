"""Unit tests for _common._fit_image — the freeform figure helper that scales
an image to fit a box (preserving aspect, centered). Verifies it never
overflows the box and reads native aspect automatically."""

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

import pytest
from pptx import Presentation
from pptx.util import Inches, Emu
from PIL import Image

from layouts._common import _fit_image


def _blank_slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def _make_png(path, w_px, h_px):
    Image.new("RGB", (w_px, h_px), (128, 128, 128)).save(str(path))
    return path


def test_wide_image_is_width_constrained(tmp_path):
    """A 2:1 image in a 4x4 box should be limited by width (4in wide, 2in tall)."""
    img = _make_png(tmp_path / "wide.png", 400, 200)
    slide = _blank_slide()
    pic = _fit_image(slide, img, left=1.0, top=1.0, max_w=4.0, max_h=4.0)
    assert pic.width == Inches(4.0)
    assert pic.height == Inches(2.0)
    # Centered vertically within the box: (4-2)/2 = 1in offset from top=1.0
    assert pic.top == Inches(2.0)
    assert pic.left == Inches(1.0)


def test_tall_image_is_height_constrained(tmp_path):
    """A 1:2 image in a 4x4 box should be limited by height (2in wide, 4in tall)."""
    img = _make_png(tmp_path / "tall.png", 200, 400)
    slide = _blank_slide()
    pic = _fit_image(slide, img, left=1.0, top=1.0, max_w=4.0, max_h=4.0)
    assert pic.width == Inches(2.0)
    assert pic.height == Inches(4.0)
    assert pic.left == Inches(2.0)  # (4-2)/2 = 1in offset from left=1.0
    assert pic.top == Inches(1.0)


def test_never_overflows_box(tmp_path):
    """Whatever the aspect, the placed picture stays within the box bounds."""
    img = _make_png(tmp_path / "odd.png", 333, 777)
    slide = _blank_slide()
    left, top, max_w, max_h = 0.5, 0.75, 5.0, 3.0
    pic = _fit_image(slide, img, left=left, top=top, max_w=max_w, max_h=max_h)
    assert pic.width <= Inches(max_w)
    assert pic.height <= Inches(max_h)
    assert pic.left >= Inches(left)
    assert pic.top >= Inches(top)
    assert pic.left + pic.width <= Inches(left + max_w) + Emu(1)
    assert pic.top + pic.height <= Inches(top + max_h) + Emu(1)
