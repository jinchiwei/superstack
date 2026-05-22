import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from pptx import Presentation

import render as render_mod
from plan import Plan, SlideEntry
from themes import get_theme


def _full_bleed_fill_hexes(pptx_path):
    """Canvas fill hexes across slides.

    The simple named renderers paint the canvas via the slide *background*
    fill (a ``<p:bg>`` solid fill), not a full-bleed shape — so we read the
    background fill directly. We also keep the full-bleed-shape scan for
    renderers (e.g. freeform) that paint the canvas as an actual rect.
    """
    prs = Presentation(str(pptx_path))
    hexes = []
    for s in prs.slides:
        # Slide-background solid fill (how _set_bg paints the canvas).
        try:
            bg_fill = s.background.fill
            if int(bg_fill.type) == 1:
                hexes.append(str(bg_fill.fore_color.rgb))
        except Exception:
            pass
        # Full-bleed solid rects (canvas-as-shape, e.g. freeform).
        for sh in s.shapes:
            try:
                w = sh.width / 914400.0
                h = sh.height / 914400.0
            except (TypeError, ValueError):
                continue
            if abs(w - 13.333) < 0.2 and abs(h - 7.5) < 0.2:
                try:
                    if int(sh.fill.type) == 1:
                        hexes.append(str(sh.fill.fore_color.rgb))
                except Exception:
                    pass
    return hexes


def test_content_text_named_layout_paints_dark_canvas(tmp_path):
    md = tmp_path / "deck.md"
    md.write_text("---\ntitle: T\n---\n\n# A\n\nsome body text here\n", encoding="utf-8")
    plan = Plan(mode="expressive", theme="midnight", slides=[
        SlideEntry(slide_id="h1-a", kind="content-text",
                   params={"title": "A", "body": [{"kind": "p", "html": "hi"}]}),
    ])
    out = tmp_path / "out.pptx"
    render_mod.render_from_plan(md_path=md, plan=plan, output_path=out,
                                theme=get_theme("midnight"))
    assert "14141C" in _full_bleed_fill_hexes(out)


import pytest


@pytest.mark.parametrize("kind,params", [
    ("three-pillars", {"title": "T", "pillars": [
        {"label": "A", "body": "x"}, {"label": "B", "body": "y"},
        {"label": "C", "body": "z"}]}),
    ("stat-callouts-right", {"title": "T", "stats": [
        {"value": "1", "label": "a"}, {"value": "2", "label": "b"}]}),
    ("timeline", {"title": "T", "milestones": [
        {"label": "A", "body": "x"}, {"label": "B", "body": "y"}]}),
])
def test_color_heavy_named_layout_dark_canvas(tmp_path, kind, params):
    md = tmp_path / "deck.md"
    md.write_text("---\ntitle: T\n---\n\n# A\n\nbody\n", encoding="utf-8")
    plan = Plan(mode="expressive", theme="midnight", slides=[
        SlideEntry(slide_id="h1-a", kind=kind, params=params),
    ])
    out = tmp_path / f"{kind}.pptx"
    render_mod.render_from_plan(md_path=md, plan=plan, output_path=out,
                                theme=get_theme("midnight"))
    assert "14141C" in _full_bleed_fill_hexes(out)


def _all_shape_fill_hexes(pptx_path):
    """Every solid-fill shape hex across all slides."""
    prs = Presentation(str(pptx_path))
    out = []
    for s in prs.slides:
        for sh in s.shapes:
            try:
                if int(sh.fill.type) == 1:
                    out.append(str(sh.fill.fore_color.rgb))
            except Exception:
                pass
    return out


def test_stats_with_takeaway_inner_tiles_themed_dark(tmp_path):
    md = tmp_path / "deck.md"
    md.write_text("---\ntitle: T\n---\n\n# A\n\nbody\n", encoding="utf-8")
    params = {
        "title": "Results",
        "lede": "summary",
        "stats": [
            {"value": "0.91", "label": "Internal AUC", "sub": "5-seed mean"},
            {"value": "0.88", "label": "External AUC", "sub": "held-out"},
        ],
        "callout": {"text": "key takeaway", "tone": "dark"},
    }
    plan = Plan(mode="expressive", theme="midnight", slides=[
        SlideEntry(slide_id="h1-a", kind="stats-with-takeaway", params=params),
    ])
    out = tmp_path / "stats.pptx"
    render_mod.render_from_plan(md_path=md, plan=plan, output_path=out,
                                theme=get_theme("midnight"))
    fills = _all_shape_fill_hexes(out)
    # No light PAPER surface remains on the inner stat tiles under a dark deck.
    assert "FAFAFC" not in fills


def test_table_with_takeaway_inner_table_themed_dark(tmp_path):
    md = tmp_path / "deck.md"
    md.write_text("---\ntitle: T\n---\n\n# A\n\nbody\n", encoding="utf-8")
    params = {
        "title": "Data",
        "lede": "summary",
        "rows": [
            ["Model", "AUC", "N"],
            ["A", "0.91", "120"],
            ["B", "0.88", "98"],
        ],
        "callout": {"text": "key takeaway", "tone": "dark"},
    }
    plan = Plan(mode="expressive", theme="midnight", slides=[
        SlideEntry(slide_id="h1-a", kind="table-with-takeaway", params=params),
    ])
    out = tmp_path / "table.pptx"
    render_mod.render_from_plan(md_path=md, plan=plan, output_path=out,
                                theme=get_theme("midnight"))
    # Table cells are not shapes, so scan cell fills directly.
    prs = Presentation(str(out))
    cell_hexes = []
    for s in prs.slides:
        for sh in s.shapes:
            if not sh.has_table:
                continue
            for row in sh.table.rows:
                for cell in row.cells:
                    try:
                        if int(cell.fill.type) == 1:
                            cell_hexes.append(str(cell.fill.fore_color.rgb))
                    except Exception:
                        pass
    # No light PAPER alternating-row fill remains under a dark deck.
    assert "FAFAFC" not in cell_hexes


def _table_cell_fill_hexes(pptx_path):
    """Every table-cell solid-fill hex across all slides."""
    prs = Presentation(str(pptx_path))
    cell_hexes = []
    for s in prs.slides:
        for sh in s.shapes:
            if not sh.has_table:
                continue
            for row in sh.table.rows:
                for cell in row.cells:
                    try:
                        if int(cell.fill.type) == 1:
                            cell_hexes.append(str(cell.fill.fore_color.rgb))
                    except Exception:
                        pass
    return cell_hexes


def test_table_striping_preserved_under_light(tmp_path):
    """Under no theme (LIGHT/strict), table rows keep alternating FAFAFC/FFFFFF.

    Regression for the block-helper theming break: the delegating renderer
    used to pass palette.surface_rgb unconditionally, collapsing the two
    distinct PAPER/WHITE row fills into one under LIGHT. Gating on
    palette.on_dark restores the striping. This FAILS on pre-fix code
    (only FAFAFC present) and PASSES after.
    """
    md = tmp_path / "deck.md"
    md.write_text("---\ntitle: T\n---\n\n# A\n\nbody\n", encoding="utf-8")
    params = {
        "title": "Data",
        "lede": "summary",
        "rows": [
            ["Model", "AUC", "N"],
            ["A", "0.91", "120"],
            ["B", "0.88", "98"],
            ["C", "0.85", "77"],
        ],
        "callout": {"text": "key takeaway", "tone": "dark"},
    }
    plan = Plan(mode="strict", slides=[
        SlideEntry(slide_id="h1-a", kind="table-with-takeaway", params=params),
    ])
    out = tmp_path / "out.pptx"
    render_mod.render_from_plan(md_path=md, plan=plan, output_path=out, theme=None)
    cell_hexes = _table_cell_fill_hexes(out)
    # Striping preserved: both alternating data-row fills present.
    assert "FAFAFC" in cell_hexes
    assert "FFFFFF" in cell_hexes


def _run_color_for_text(pptx_path, needle):
    """Font color hex of the first run whose text contains ``needle``."""
    prs = Presentation(str(pptx_path))
    for s in prs.slides:
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if needle in (run.text or ""):
                        try:
                            return str(run.font.color.rgb)
                        except Exception:
                            return None
    return None


def test_stat_tile_sub_distinct_from_label_under_light(tmp_path):
    """Under no theme (LIGHT/strict), stat-tile label (MUTED #555560) and
    sub (DIM #888888) stay distinct colors.

    Regression for the block-helper theming break: the delegating renderer
    used to pass palette.muted_rgb unconditionally, collapsing both into
    #555560 under LIGHT. Gating on palette.on_dark restores the distinct
    DIM sub color. FAILS on pre-fix code (no 888888), PASSES after.
    """
    md = tmp_path / "deck.md"
    md.write_text("---\ntitle: T\n---\n\n# A\n\nbody\n", encoding="utf-8")
    params = {
        "title": "Results",
        "lede": "summary",
        "stats": [
            {"value": "0.91", "label": "Internal AUC", "sub": "5-seed mean"},
            {"value": "0.88", "label": "External AUC", "sub": "held-out"},
        ],
        "callout": {"text": "key takeaway", "tone": "dark"},
    }
    plan = Plan(mode="strict", slides=[
        SlideEntry(slide_id="h1-a", kind="stats-with-takeaway", params=params),
    ])
    out = tmp_path / "out.pptx"
    render_mod.render_from_plan(md_path=md, plan=plan, output_path=out, theme=None)
    # Sub text keeps its distinct DIM color (#888888); label stays MUTED
    # (#555560). Pre-fix the sub collapsed to the label's #555560.
    assert _run_color_for_text(out, "5-seed mean") == "888888"
    assert _run_color_for_text(out, "Internal AUC") == "555560"
