"""Canonical speaker-notes support: params['notes'] is embedded into the
PowerPoint notes pane (Presenter View) for content slides and dividers."""

import sys
from pathlib import Path

SKILL_DIR = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(SKILL_DIR))


def _render(tmp_path):
    from plan import Plan, SlideEntry
    from render import render_from_plan
    plan = Plan(mode="strict", slides=[
        SlideEntry(slide_id="divider-x", kind="section-divider",
                   params={"label": "Background", "notes": "Transition note here."},
                   content_hash="d"),
        SlideEntry(slide_id="x", kind="content-text",
                   params={"title": "T", "lede": "L",
                           "body": [{"kind": "bullet", "html": "one"}],
                           "notes": "Comprehensive didactic note for slide X."},
                   content_hash="c"),
        SlideEntry(slide_id="y", kind="content-text",
                   params={"title": "T2", "lede": "L2",
                           "body": [{"kind": "bullet", "html": "two"}]},  # no notes
                   content_hash="e"),
    ])
    md = tmp_path / "deck.md"
    md.write_text("---\ntitle: 'Notes'\nname: 'Jinchi Wei'\n---\n# Background\n")
    out = tmp_path / "out.pptx"
    render_from_plan(md_path=md, plan=plan, output_path=out,
                     no_cover=True, no_end=True, theme=None)
    return out


def test_notes_embedded_in_notes_pane(tmp_path):
    from pptx import Presentation
    out = _render(tmp_path)
    prs = Presentation(str(out))
    texts = [s.notes_slide.notes_text_frame.text if s.has_notes_slide else ""
             for s in prs.slides]
    joined = "\n".join(texts)
    assert "Transition note here." in joined          # divider notes
    assert "Comprehensive didactic note for slide X." in joined  # content notes


def test_missing_notes_does_not_crash(tmp_path):
    """A slide without params['notes'] renders fine (notes are additive)."""
    from pptx import Presentation
    out = _render(tmp_path)
    prs = Presentation(str(out))
    assert len(prs.slides) == 3  # divider + 2 content (no cover/end)


def test_notes_pdf_skips_gracefully_without_notes(tmp_path):
    """build_notes_pdf returns None (no crash) when the deck has no notes."""
    from plan import Plan, SlideEntry
    from render import render_from_plan
    from notes_pdf import build_notes_pdf
    plan = Plan(mode="strict", slides=[
        SlideEntry(slide_id="x", kind="content-text",
                   params={"title": "T", "lede": "L",
                           "body": [{"kind": "bullet", "html": "one"}]},
                   content_hash="c"),
    ])
    md = tmp_path / "deck.md"
    md.write_text("---\ntitle: 'X'\nname: 'Jinchi Wei'\n---\n# S\n")
    out = tmp_path / "out.pptx"
    render_from_plan(md_path=md, plan=plan, output_path=out,
                     no_cover=True, no_end=True, theme=None)
    assert build_notes_pdf(out) is None  # no notes -> nothing to hand out
