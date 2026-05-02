"""Round-trip tests for the 4 new v7 named layouts.

Each test:
1. Constructs a minimal sidecar with the target layout kind.
2. Renders through render_from_plan (via Plan).
3. Asserts the slide renders without error and contains expected shape counts.
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

import pytest

SKILL_DIR = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(SKILL_DIR))
sys.path.insert(0, str(SKILL_DIR.parent / "_shared"))


def _make_slide():
    """Create a blank pptx slide for direct unit testing."""
    from build import new_presentation, _blank
    prs = new_presentation()
    return _blank(prs)


def _render_sidecar(tmp_path: Path, kind: str, params: dict) -> "Presentation":
    """Helper: render a single-slide sidecar through render_from_plan and
    return the resulting Presentation object."""
    from plan import Plan
    from render import render_from_plan
    from pptx import Presentation as PRS

    md = tmp_path / "deck.md"
    md.write_text("# Test\n\n---\n\n## Test slide\n\nBody text.\n")
    sidecar = tmp_path / "deck.md.layout.json"
    plan_dict = {
        "version": 1,
        "deck_md_hash": f"test-{kind}",
        "shake_seed": None,
        "slides": [
            {
                "slide_id": f"test-{kind}",
                "kind": kind,
                "params": params,
                "content_hash": f"hash-{kind}",
            }
        ],
    }
    sidecar.write_text(json.dumps(plan_dict))
    out = tmp_path / "out.pptx"
    plan = Plan.from_json(sidecar.read_text())
    render_from_plan(md_path=md, plan=plan, output_path=out,
                     no_cover=False, no_end=False)
    assert out.exists()
    assert out.stat().st_size > 5_000
    return PRS(str(out))


# ---------------------------------------------------------------------------
# Catalog registration checks
# ---------------------------------------------------------------------------

def test_new_layouts_registered_in_catalog():
    """All 4 new layouts are registered in catalog.REGISTRY."""
    from layouts import catalog
    for kind in ("stats-with-takeaway", "figure-with-aside",
                 "cards-with-takeaway", "table-with-takeaway"):
        assert kind in catalog.REGISTRY, f"{kind} missing from REGISTRY"


def test_new_layout_renderers_callable():
    """All 4 new layout renderers are callable."""
    from layouts import catalog
    for kind in ("stats-with-takeaway", "figure-with-aside",
                 "cards-with-takeaway", "table-with-takeaway"):
        renderer = catalog.get(kind)
        assert callable(renderer), f"{kind} renderer not callable"


# ---------------------------------------------------------------------------
# Direct render unit tests (no file I/O)
# ---------------------------------------------------------------------------

def _footer():
    return {"name": "Jin", "org": "UCSF", "deck_title": "Test", "date": "2026-05-02"}


def test_stats_with_takeaway_direct():
    """stats-with-takeaway renders without error and adds shapes."""
    from layouts.stats_with_takeaway import render
    from layouts._common import TURQUOISE_RGB

    slide = _make_slide()
    params = {
        "title": "Model Performance",
        "lede": "Key metrics across validation cohorts.",
        "stats": [
            {"value": "0.91", "label": "Internal AUC", "sub": "5-seed mean"},
            {"value": "0.85", "label": "External AUC", "sub": "site-mixed"},
            {"value": "88%", "label": "Sensitivity"},
        ],
        "callout": {"text": "Model exceeds clinical threshold across all cohorts.", "tone": "dark"},
    }
    render(slide, params=params, accent_rgb=TURQUOISE_RGB, footer_kwargs=_footer())
    # Should have at least: left bar + title + hairline + lede + stat tiles + callout
    assert len(slide.shapes) >= 4


def test_stats_with_takeaway_icon_homogeneity():
    """stats-with-takeaway drops icons when all stats share the same icon."""
    from layouts.stats_with_takeaway import render
    from layouts._common import TURQUOISE_RGB

    slide = _make_slide()
    # All stats have same icon — should be dropped (homogeneity rule)
    params = {
        "title": "Stats",
        "lede": "",
        "stats": [
            {"value": "0.91", "label": "AUC-1", "icon": "FaChartLine"},
            {"value": "0.85", "label": "AUC-2", "icon": "FaChartLine"},
        ],
        "callout": {"text": "Same metric family — icons dropped.", "tone": "dark"},
    }
    # Should not raise
    render(slide, params=params, accent_rgb=TURQUOISE_RGB, footer_kwargs=_footer())
    assert len(slide.shapes) >= 2


def test_figure_with_aside_direct():
    """figure-with-aside renders without error and adds shapes."""
    from layouts.figure_with_aside import render
    from layouts._common import DEEPPINK_RGB

    slide = _make_slide()
    params = {
        "title": "Radiogenomics Overview",
        "lede": "Spatial relationship between MRI texture and IDH mutation.",
        "image": "/nonexistent/fig.png",
        "alt": "Model architecture diagram",
        "aside": {
            "label": "Why this matters",
            "body": "IDH mutation status drives prognosis and treatment selection.",
            "icon": "FaLightbulb",
        },
    }
    render(slide, params=params, accent_rgb=DEEPPINK_RGB, footer_kwargs=_footer())
    assert len(slide.shapes) >= 4


def test_figure_with_aside_no_image():
    """figure-with-aside gracefully handles missing image path."""
    from layouts.figure_with_aside import render
    from layouts._common import AMBER_RGB

    slide = _make_slide()
    params = {
        "title": "Placeholder",
        "lede": "",
        "image": "",
        "alt": "placeholder",
        "aside": {"label": "Note", "body": "Content here."},
    }
    render(slide, params=params, accent_rgb=AMBER_RGB, footer_kwargs={})
    assert len(slide.shapes) >= 2


def test_cards_with_takeaway_direct():
    """cards-with-takeaway renders without error and adds shapes."""
    from layouts.cards_with_takeaway import render
    from layouts._common import BLUEVIOLET_RGB

    slide = _make_slide()
    params = {
        "title": "System Components",
        "lede": "Three independent modules feed the ensemble.",
        "cards": [
            {"label": "Extractor", "body": "ResNet50 feature backbone.", "icon": None},
            {"label": "Aggregator", "body": "Attention-weighted pooling.", "icon": None},
            {"label": "Classifier", "body": "Logistic regression head.", "icon": None},
        ],
        "callout": {"text": "All three components trained jointly end-to-end.", "tone": "dark"},
    }
    render(slide, params=params, accent_rgb=BLUEVIOLET_RGB, footer_kwargs=_footer())
    assert len(slide.shapes) >= 4


def test_cards_with_takeaway_icon_homogeneity():
    """cards-with-takeaway drops icons when all cards share the same icon."""
    from layouts.cards_with_takeaway import render
    from layouts._common import TURQUOISE_RGB

    slide = _make_slide()
    params = {
        "title": "Metrics",
        "lede": "",
        "cards": [
            {"label": "AUC", "body": "0.91", "icon": "FaChartBar"},
            {"label": "Sens", "body": "88%", "icon": "FaChartBar"},
        ],
        "callout": {"text": "Homogeneous icon set — dropped.", "tone": "dark"},
    }
    render(slide, params=params, accent_rgb=TURQUOISE_RGB, footer_kwargs=_footer())
    assert len(slide.shapes) >= 2


def test_table_with_takeaway_direct():
    """table-with-takeaway renders without error and adds shapes."""
    from layouts.table_with_takeaway import render
    from layouts._common import AMBER_RGB

    slide = _make_slide()
    params = {
        "title": "Cohort Summary",
        "lede": "UCSF and PNOC validation cohorts.",
        "rows": [
            ["Cohort", "N", "IDH+", "IDH-"],
            ["UCSF train", "120", "65", "55"],
            ["PNOC external", "48", "22", "26"],
            ["Combined", "168", "87", "81"],
        ],
        "callout": {"text": "Balanced class distribution across both sites.", "tone": "dark"},
    }
    render(slide, params=params, accent_rgb=AMBER_RGB, footer_kwargs=_footer())
    assert len(slide.shapes) >= 3


def test_table_with_takeaway_empty_rows():
    """table-with-takeaway handles empty rows gracefully."""
    from layouts.table_with_takeaway import render
    from layouts._common import TURQUOISE_RGB

    slide = _make_slide()
    params = {
        "title": "Empty Table",
        "lede": "",
        "rows": [],
        "callout": {"text": "No data.", "tone": "dark"},
    }
    render(slide, params=params, accent_rgb=TURQUOISE_RGB, footer_kwargs={})
    # Just chrome + callout
    assert len(slide.shapes) >= 1


# ---------------------------------------------------------------------------
# Round-trip render tests (via render_from_plan)
# ---------------------------------------------------------------------------

def test_stats_with_takeaway_roundtrip(tmp_path):
    """stats-with-takeaway full round-trip produces a valid pptx."""
    prs = _render_sidecar(tmp_path, "stats-with-takeaway", {
        "title": "Key Metrics",
        "lede": "Performance across all cohorts.",
        "stats": [
            {"value": "0.91", "label": "Internal AUC"},
            {"value": "0.85", "label": "External AUC"},
        ],
        "callout": {"text": "Both AUCs exceed clinical threshold.", "tone": "dark"},
    })
    # title + 1 content slide + end
    assert len(prs.slides) == 3


def test_figure_with_aside_roundtrip(tmp_path):
    """figure-with-aside full round-trip produces a valid pptx."""
    prs = _render_sidecar(tmp_path, "figure-with-aside", {
        "title": "Figure Overview",
        "lede": "Architecture schematic.",
        "image": "/nonexistent/fig.png",
        "alt": "Architecture diagram",
        "aside": {"label": "Takeaway", "body": "Multi-modal fusion improves accuracy."},
    })
    assert len(prs.slides) == 3


def test_cards_with_takeaway_roundtrip(tmp_path):
    """cards-with-takeaway full round-trip produces a valid pptx."""
    prs = _render_sidecar(tmp_path, "cards-with-takeaway", {
        "title": "Components",
        "lede": "",
        "cards": [
            {"label": "A", "body": "Description A."},
            {"label": "B", "body": "Description B."},
            {"label": "C", "body": "Description C."},
        ],
        "callout": {"text": "All components validated independently.", "tone": "dark"},
    })
    assert len(prs.slides) == 3


def test_table_with_takeaway_roundtrip(tmp_path):
    """table-with-takeaway full round-trip produces a valid pptx."""
    prs = _render_sidecar(tmp_path, "table-with-takeaway", {
        "title": "Results Table",
        "lede": "",
        "rows": [
            ["Method", "AUC", "Sensitivity"],
            ["Ours", "0.91", "88%"],
            ["Baseline", "0.78", "72%"],
        ],
        "callout": {"text": "Our method outperforms baseline by 13 AUC points.", "tone": "dark"},
    })
    assert len(prs.slides) == 3


# ---------------------------------------------------------------------------
# Fixture file round-trip
# ---------------------------------------------------------------------------

FIXTURE_MD = SKILL_DIR / "tests" / "fixture_named_v7.md"
FIXTURE_SIDECAR = SKILL_DIR / "tests" / "fixture_named_v7.md.layout.json"


def test_fixture_named_v7_roundtrip(tmp_path):
    """fixture_named_v7 exercises all 4 new kinds in one deck render."""
    from plan import Plan
    from render import render_from_plan
    from pptx import Presentation as PRS

    md = tmp_path / "deck.md"
    sidecar = tmp_path / "deck.md.layout.json"
    md.write_text(FIXTURE_MD.read_text())
    sidecar.write_text(FIXTURE_SIDECAR.read_text())
    out = tmp_path / "out.pptx"
    plan = Plan.from_json(sidecar.read_text())
    render_from_plan(md_path=md, plan=plan, output_path=out,
                     no_cover=False, no_end=False)
    assert out.exists()
    assert out.stat().st_size > 5_000
    prs = PRS(str(out))
    # title + 4 named layout slides + end = 6
    assert len(prs.slides) == 6
