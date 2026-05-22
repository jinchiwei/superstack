import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from pptx import Presentation

import render as render_mod
from plan import Plan, SlideEntry
from themes import get_theme


def _bg_hex(slide):
    try:
        return str(slide.background.fill.fore_color.rgb)
    except Exception:
        return None


def _build(tmp_path, theme):
    md = tmp_path / "deck.md"
    md.write_text("---\ntitle: T\nname: X\n---\n\n# Section\n\n---\n\n## A\n\nbody\n",
                  encoding="utf-8")
    plan = Plan(mode=("expressive" if theme else "strict"),
                theme=(theme.name if theme else None), slides=[
        SlideEntry(slide_id="divider-h1-section", kind="section-divider",
                   params={"label": "Section"}),
        SlideEntry(slide_id="h1-section/h2-a", kind="content-text",
                   params={"title": "A", "body": [{"kind": "p", "html": "hi"}]}),
    ])
    out = tmp_path / "out.pptx"
    render_mod.render_from_plan(md_path=md, plan=plan, output_path=out, theme=theme)
    return Presentation(str(out))


def test_cover_and_divider_match_dark_theme_canvas(tmp_path):
    prs = _build(tmp_path, get_theme("midnight"))  # canvas #14141C
    # slide 0 = cover, slide 1 = section divider
    assert _bg_hex(prs.slides[0]) == "14141C", "cover should use theme canvas, not navy"
    assert _bg_hex(prs.slides[1]) == "14141C", "divider should use theme canvas, not navy"


def test_cover_stays_navy_without_theme(tmp_path):
    prs = _build(tmp_path, None)  # strict / no theme -> hardcoded navy DARK_BG
    assert _bg_hex(prs.slides[0]) == "0E1A35", "no-theme cover must stay navy (parity)"
    assert _bg_hex(prs.slides[1]) == "0E1A35", "no-theme divider must stay navy (parity)"
