"""Unit tests for the 5 new creative layout primitives.

Each test:
  - Builds a fresh Presentation + blank slide
  - Calls <layout>.render(slide, params=..., accent_rgb=TURQUOISE_RGB,
      footer_kwargs={...})
  - Asserts the slide has > N shapes
  - Saves to a tmp file to confirm pptx serializes cleanly
"""

from __future__ import annotations

import sys
from pathlib import Path

import pytest

# Ensure layouts package is importable
sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from pptx import Presentation
from pptx.util import Inches

from layouts._common import _rgb, TURQUOISE_RGB, _blank


FOOTER = {
    "name": "Jin",
    "org": "UCSF",
    "deck_title": "Test Deck",
    "date": "2026-05-02",
}


def _new_slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = _blank(prs)
    return prs, slide


# ─────────────────────────────────────────────────────────────────────────────
# Layout 1: cards-heterogeneous
# ─────────────────────────────────────────────────────────────────────────────

def test_cards_heterogeneous_renders(tmp_path):
    from layouts import cards_heterogeneous

    prs, slide = _new_slide()
    params = {
        "title": "Heterogeneous Cards",
        "lede": "Primary card left, secondary cards stacked right",
        "primary_card": {
            "label": "Primary Label",
            "body": "This is the main body text for the primary card.",
            "icon": None,
        },
        "secondary_cards": [
            {"label": "Secondary A", "body": "Body for secondary A.", "icon": None},
            {"label": "Secondary B", "body": "Body for secondary B.", "icon": None},
            {"label": "Secondary C", "body": "Body for secondary C.", "icon": None},
        ],
    }
    cards_heterogeneous.render(slide, params=params, accent_rgb=TURQUOISE_RGB,
                               footer_kwargs=FOOTER)

    # chrome ~5 + primary card ~2 + 3 secondary cards ~2 each = ≥10 shapes
    assert len(slide.shapes) > 8, f"Expected >8 shapes, got {len(slide.shapes)}"

    out = tmp_path / "cards_heterogeneous.pptx"
    prs.save(str(out))
    assert out.is_file()
    assert out.stat().st_size > 2000


def test_cards_heterogeneous_two_secondary(tmp_path):
    """Works with only 2 secondary cards."""
    from layouts import cards_heterogeneous

    prs, slide = _new_slide()
    params = {
        "title": "Two Secondary",
        "lede": "",
        "primary_card": {"label": "Big Card", "body": "Primary body.", "icon": None},
        "secondary_cards": [
            {"label": "A", "body": "Body A.", "icon": None},
            {"label": "B", "body": "Body B.", "icon": None},
        ],
    }
    cards_heterogeneous.render(slide, params=params, accent_rgb=TURQUOISE_RGB,
                               footer_kwargs=FOOTER)
    assert len(slide.shapes) > 6

    out = tmp_path / "cards_heterogeneous_2sec.pptx"
    prs.save(str(out))
    assert out.is_file()


# ─────────────────────────────────────────────────────────────────────────────
# Layout 2: three-pillars
# ─────────────────────────────────────────────────────────────────────────────

def test_three_pillars_renders(tmp_path):
    from layouts import three_pillars

    prs, slide = _new_slide()
    params = {
        "title": "Three Pillars",
        "lede": "Three side-by-side columns with arrows",
        "pillars": [
            {"label": "Data", "body": "Curated clinical data from UCSF.", "color_role": "primary"},
            {"label": "Model", "body": "Transformer-based classifier.", "color_role": "secondary"},
            {"label": "Deploy", "body": "Real-time inference API.", "color_role": "tertiary"},
        ],
        "show_arrows": True,
    }
    three_pillars.render(slide, params=params, accent_rgb=TURQUOISE_RGB,
                         footer_kwargs=FOOTER)

    # chrome ~5 + 3 pillars ~2 each + 2 arrows = ≥12 shapes
    assert len(slide.shapes) > 10, f"Expected >10 shapes, got {len(slide.shapes)}"

    out = tmp_path / "three_pillars.pptx"
    prs.save(str(out))
    assert out.is_file()
    assert out.stat().st_size > 2000


def test_three_pillars_no_arrows(tmp_path):
    from layouts import three_pillars

    prs, slide = _new_slide()
    params = {
        "title": "Two Pillars No Arrows",
        "lede": "",
        "pillars": [
            {"label": "Phase 1", "body": "Description A.", "color_role": None},
            {"label": "Phase 2", "body": "Description B.", "color_role": None},
        ],
        "show_arrows": False,
    }
    three_pillars.render(slide, params=params, accent_rgb=TURQUOISE_RGB,
                         footer_kwargs=FOOTER)
    assert len(slide.shapes) > 5

    out = tmp_path / "three_pillars_no_arrows.pptx"
    prs.save(str(out))
    assert out.is_file()


# ─────────────────────────────────────────────────────────────────────────────
# Layout 3: stat-callouts-right
# ─────────────────────────────────────────────────────────────────────────────

def test_stat_callouts_right_renders(tmp_path):
    from layouts import stat_callouts_right

    # Use a placeholder path that doesn't exist — the renderer should
    # fall back gracefully when the image can't be loaded
    prs, slide = _new_slide()
    params = {
        "title": "Model Performance",
        "lede": "Chart left, stats right",
        "image": Path("/nonexistent/chart.png"),
        "stats": [
            {"value": "0.91", "label": "Internal val AUC"},
            {"value": "0.87", "label": "External test AUC"},
            {"value": "1,240", "label": "Patients in cohort"},
            {"value": "94%", "label": "Sensitivity at 0.5 cut"},
        ],
    }
    stat_callouts_right.render(slide, params=params, accent_rgb=TURQUOISE_RGB,
                               footer_kwargs=FOOTER)

    # chrome ~5 + error text or image + 4 stats * 2 textboxes each = ≥14
    assert len(slide.shapes) > 10, f"Expected >10 shapes, got {len(slide.shapes)}"

    out = tmp_path / "stat_callouts_right.pptx"
    prs.save(str(out))
    assert out.is_file()
    assert out.stat().st_size > 2000


# ─────────────────────────────────────────────────────────────────────────────
# Layout 4: bg-flip
# ─────────────────────────────────────────────────────────────────────────────

def test_bg_flip_renders(tmp_path):
    from layouts import bg_flip

    prs, slide = _new_slide()
    params = {
        "title": "Key Takeaway",
        "lede": "This slide uses a dark navy background for emphasis.",
        "body": [
            {"kind": "paragraph", "html": "CurieDx demonstrates strong generalization."},
            {"kind": "bullet", "html": "AUC 0.91 internal validation"},
            {"kind": "bullet", "html": "AUC 0.87 external holdout"},
        ],
    }
    bg_flip.render(slide, params=params, accent_rgb=TURQUOISE_RGB,
                   footer_kwargs=FOOTER)

    # left bar rect + title text + hairline rect + lede text + body textbox + footer text = 6
    assert len(slide.shapes) >= 5, f"Expected >=5 shapes, got {len(slide.shapes)}"

    out = tmp_path / "bg_flip.pptx"
    prs.save(str(out))
    assert out.is_file()
    assert out.stat().st_size > 2000


def test_bg_flip_has_dark_background(tmp_path):
    """The bg-flip slide background must be the dark navy colour (via slide.background.fill)."""
    from layouts import bg_flip

    prs, slide = _new_slide()
    params = {
        "title": "Dark Slide",
        "lede": "",
        "body": [{"kind": "paragraph", "html": "Text on dark."}],
    }
    bg_flip.render(slide, params=params, accent_rgb=TURQUOISE_RGB,
                   footer_kwargs=FOOTER)

    # _set_bg sets the slide background fill, not a shape rect
    fill = slide.background.fill
    # fill.type == 1 means solid fill
    assert fill.type == 1, "Slide background is not a solid fill"
    bg_rgb = str(fill.fore_color.rgb).upper()
    assert bg_rgb == "0E1A35", f"Expected dark bg #0E1A35, got #{bg_rgb}"


# ─────────────────────────────────────────────────────────────────────────────
# Layout 5: timeline
# ─────────────────────────────────────────────────────────────────────────────

def test_timeline_renders(tmp_path):
    from layouts import timeline

    prs, slide = _new_slide()
    params = {
        "title": "Project Timeline",
        "lede": "From pilot to deployment",
        "milestones": [
            {"date": "2026-Q1", "label": "Pilot Launch", "body": "IRB approval + data pull."},
            {"date": "2026-Q2", "label": "Model v1", "body": "Train on retrospective cohort."},
            {"date": "2026-Q3", "label": "External Val", "body": "Partner site validation."},
            {"date": "2026-Q4", "label": "Deploy", "body": "EHR integration go-live."},
        ],
    }
    timeline.render(slide, params=params, accent_rgb=TURQUOISE_RGB,
                    footer_kwargs=FOOTER)

    # chrome ~5 + axis line + 4 milestones * (dot + label + date + body) ≥ ~20
    assert len(slide.shapes) > 14, f"Expected >14 shapes, got {len(slide.shapes)}"

    out = tmp_path / "timeline.pptx"
    prs.save(str(out))
    assert out.is_file()
    assert out.stat().st_size > 2000


def test_timeline_two_milestones(tmp_path):
    """Timeline works with just 2 milestones."""
    from layouts import timeline

    prs, slide = _new_slide()
    params = {
        "title": "Mini Timeline",
        "lede": "",
        "milestones": [
            {"date": "2026-Q1", "label": "Start", "body": "Kick-off."},
            {"date": "2026-Q4", "label": "End", "body": "Ship."},
        ],
    }
    timeline.render(slide, params=params, accent_rgb=TURQUOISE_RGB,
                    footer_kwargs=FOOTER)
    assert len(slide.shapes) > 8

    out = tmp_path / "timeline_2.pptx"
    prs.save(str(out))
    assert out.is_file()
