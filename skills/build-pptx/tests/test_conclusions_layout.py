"""Tests for the 'conclusions' named layout.

Covers:
  - Catalog registration
  - Direct render: 4-card with callout (dark bg, distinct accent fills)
  - Direct render: 3-card no callout (no callout-shaped element)
  - Icon resolution: cards with FaName strings get glyph PNGs rendered
  - Fixture file round-trip
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

import pytest

SKILL_DIR = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(SKILL_DIR))
sys.path.insert(0, str(SKILL_DIR.parent / "_shared"))


def _make_slide():
    """Create a blank pptx slide for direct unit testing."""
    from build import new_presentation, _blank
    prs = new_presentation()
    return _blank(prs)


def _footer():
    return {"name": "Jin", "org": "UCSF", "deck_title": "Test", "date": "2026-05-02"}


def _render_sidecar(tmp_path: Path, kind: str, params: dict):
    """Helper: render a single-slide sidecar through render_from_plan and
    return the resulting Presentation object."""
    from plan import Plan
    from render import render_from_plan
    from pptx import Presentation as PRS

    md = tmp_path / "deck.md"
    md.write_text("# Test\n\n---\n\n## Test slide\n\nBody text.\n")
    sidecar = tmp_path / "deck.md.layout.json"
    plan_dict = {
        "version": 1,
        "deck_md_hash": f"test-{kind}",
        "shake_seed": None,
        "slides": [
            {
                "slide_id": f"test-{kind}",
                "kind": kind,
                "params": params,
                "content_hash": f"hash-{kind}",
            }
        ],
    }
    sidecar.write_text(json.dumps(plan_dict))
    out = tmp_path / "out.pptx"
    plan = Plan.from_json(sidecar.read_text())
    render_from_plan(md_path=md, plan=plan, output_path=out,
                     no_cover=False, no_end=False)
    assert out.exists()
    assert out.stat().st_size > 5_000
    return PRS(str(out))


# ---------------------------------------------------------------------------
# Catalog registration
# ---------------------------------------------------------------------------

def test_conclusions_registered_in_catalog():
    """'conclusions' is registered in catalog.REGISTRY."""
    from layouts import catalog
    assert "conclusions" in catalog.REGISTRY, "'conclusions' missing from REGISTRY"


def test_conclusions_renderer_callable():
    """'conclusions' renderer is callable."""
    from layouts import catalog
    renderer = catalog.get("conclusions")
    assert callable(renderer)


# ---------------------------------------------------------------------------
# Direct render: 4-card with callout
# ---------------------------------------------------------------------------

def test_conclusions_4card_with_callout():
    """4-card conclusions slide renders, has dark bg, and 4 distinct accent fills."""
    from layouts.conclusions import render
    from layouts._common import TURQUOISE_RGB, DARK_BG_RGB
    from pptx.dml.color import RGBColor

    slide = _make_slide()
    params = {
        "title": "Takeaways",
        "lede": "Five findings that survive every methodology check.",
        "section_label": "Takeaways",
        "cards": [
            {"label": "Headline AUC",   "body": "AUC 0.848 site-mixed eval.", "icon": "FaChartLine"},
            {"label": "Survivors",      "body": "OLIG2 (0.808) leak-free.",   "icon": "FaCheckCircle"},
            {"label": "Leakage",        "body": "Original inflated by 0.15.", "icon": "FaExclamationTriangle"},
            {"label": "Honest ceiling", "body": "AUC 0.95 was never realistic.", "icon": "FaCrosshairs"},
        ],
        "callout": {"text": "Path forward: more data, not more architectures.", "tone": "dark"},
    }
    render(slide, params=params, accent_rgb=TURQUOISE_RGB, footer_kwargs=_footer())

    # Slide should have rendered shapes
    assert len(slide.shapes) >= 4

    # Background should be dark navy
    bg_color = slide.background.fill.fore_color.rgb
    assert bg_color == DARK_BG_RGB, f"Expected dark bg, got {bg_color}"

    # Collect all fill colors from shapes (card accent stripes + card fills)
    fill_hexes = set()
    for shp in slide.shapes:
        try:
            rgb = shp.fill.fore_color.rgb
            fill_hexes.add("#{:02X}{:02X}{:02X}".format(rgb[0], rgb[1], rgb[2]))
        except Exception:
            pass

    # All 4 brand accent colors should appear as card top-stripe fills
    brand_accents = {"#40E0D0", "#FF1493", "#F0C840", "#8A2BE2"}
    found = brand_accents & fill_hexes
    assert len(found) == 4, (
        f"Expected all 4 brand accents as card stripe fills, found: {found}\n"
        f"All fills: {fill_hexes}"
    )


# ---------------------------------------------------------------------------
# Direct render: 3-card no callout
# ---------------------------------------------------------------------------

def test_conclusions_3card_no_callout():
    """3-card conclusions slide without callout renders and has no callout fill."""
    from layouts.conclusions import render
    from layouts._common import DEEPPINK_RGB, DARK_BG_RGB

    slide = _make_slide()
    params = {
        "title": "Key Findings",
        "lede": "Three findings from the study.",
        "section_label": "Key Findings",
        "cards": [
            {"label": "Finding A", "body": "Strong result on held-out data.", "icon": "FaCheckCircle"},
            {"label": "Finding B", "body": "Robust to site variation.",       "icon": "FaLightbulb"},
            {"label": "Finding C", "body": "Generalizes across grades.",      "icon": "FaCrosshairs"},
        ],
        # No callout
    }
    before_count = len(slide.shapes)
    render(slide, params=params, accent_rgb=DEEPPINK_RGB, footer_kwargs=_footer())
    after_count = len(slide.shapes)

    # Should have rendered
    assert after_count > before_count

    # Collect shape areas — callout would be a very wide, short rectangle near the bottom.
    # We verify no such element has DARK_BG_RGB fill (the callout's tone="dark" fill).
    # Since the background itself is dark, check that we don't have a second DARK_BG_RGB
    # filled rectangle beyond the background fill.
    dark_bg_hex = "#0D1B2A"
    dark_bg_fill_shapes = []
    for shp in slide.shapes:
        try:
            rgb = shp.fill.fore_color.rgb
            hex_color = "#{:02X}{:02X}{:02X}".format(rgb[0], rgb[1], rgb[2])
            if hex_color.upper() == dark_bg_hex.upper():
                dark_bg_fill_shapes.append(shp)
        except Exception:
            pass

    # No DARK_BG_RGB-filled shape should appear since there's no callout
    # (accent_callout with tone="dark" uses DARK_BG_RGB as fill)
    assert len(dark_bg_fill_shapes) == 0, (
        f"Expected no dark-bg callout shape, but found {len(dark_bg_fill_shapes)} shapes with DARK_BG_RGB fill"
    )


# ---------------------------------------------------------------------------
# Icon resolution test
# ---------------------------------------------------------------------------

def test_conclusions_icon_resolution():
    """Cards with FaName icon strings attempt icon resolution without crashing."""
    from layouts.conclusions import render
    from layouts._common import AMBER_RGB

    slide = _make_slide()
    params = {
        "title": "Summary",
        "lede": "Icons should resolve to glyph PNGs or degrade gracefully.",
        "section_label": "Summary",
        "cards": [
            {"label": "Result",  "body": "AUC 0.91 internal.",     "icon": "FaChartLine"},
            {"label": "Caveat",  "body": "External not yet tested.", "icon": "FaExclamationTriangle"},
            {"label": "Next",    "body": "Prospective validation.", "icon": "FaArrowRight"},
        ],
    }
    # Should not raise regardless of whether cairosvg / icon registry is available
    render(slide, params=params, accent_rgb=AMBER_RGB, footer_kwargs=_footer())
    assert len(slide.shapes) >= 3


# ---------------------------------------------------------------------------
# Grid split helper unit tests
# ---------------------------------------------------------------------------

def test_grid_split_sizes():
    """_grid_split returns the correct row partitions."""
    from layouts.conclusions import _grid_split

    assert _grid_split(2) == [2]
    assert _grid_split(3) == [3]
    assert _grid_split(4) == [2, 2]
    assert _grid_split(5) == [3, 2]
    assert _grid_split(6) == [3, 3]
    assert _grid_split(7) == [4, 3]
    assert _grid_split(8) == [4, 4]
    assert _grid_split(9) == [5, 4]

    # All splits must sum to n
    for n in range(2, 12):
        split = _grid_split(n)
        assert sum(split) == n, f"_grid_split({n}) sums to {sum(split)}, expected {n}"


# ---------------------------------------------------------------------------
# Round-trip render test
# ---------------------------------------------------------------------------

def test_conclusions_roundtrip(tmp_path):
    """conclusions full round-trip produces a valid pptx."""
    prs = _render_sidecar(tmp_path, "conclusions", {
        "title": "Takeaways",
        "lede": "Key results that survived every check.",
        "section_label": "Takeaways",
        "cards": [
            {"label": "Headline AUC",   "body": "0.848 site-mixed, leak-free.", "icon": "FaChartLine"},
            {"label": "Survivors",      "body": "OLIG2 and PIK3CA hold up.",    "icon": "FaCheckCircle"},
            {"label": "Leakage found",  "body": "TP53/ATRX collapse to chance.","icon": "FaExclamationTriangle"},
            {"label": "Honest ceiling", "body": "~0.85 site-mixed ceiling.",    "icon": "FaCrosshairs"},
        ],
        "callout": {"text": "Path forward: more data, not more architectures.", "tone": "dark"},
    })
    # title slide + 1 conclusions slide + end slide = 3
    assert len(prs.slides) == 3


# ---------------------------------------------------------------------------
# Fixture file round-trip
# ---------------------------------------------------------------------------

FIXTURE_MD = SKILL_DIR / "tests" / "fixture_conclusions.md"
FIXTURE_SIDECAR = SKILL_DIR / "tests" / "fixture_conclusions.md.layout.json"


def test_fixture_conclusions_roundtrip(tmp_path):
    """fixture_conclusions exercises the conclusions layout end-to-end."""
    from plan import Plan
    from render import render_from_plan
    from pptx import Presentation as PRS

    md = tmp_path / "deck.md"
    sidecar = tmp_path / "deck.md.layout.json"
    md.write_text(FIXTURE_MD.read_text())
    sidecar.write_text(FIXTURE_SIDECAR.read_text())
    out = tmp_path / "out.pptx"
    plan = Plan.from_json(sidecar.read_text())
    render_from_plan(md_path=md, plan=plan, output_path=out,
                     no_cover=False, no_end=False)
    assert out.exists()
    assert out.stat().st_size > 5_000
    prs = PRS(str(out))
    # title + 1 conclusions slide + end = 3
    assert len(prs.slides) == 3
