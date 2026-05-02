"""End-to-end tests for v4 plan→render."""

import sys
import json
import subprocess
from pathlib import Path

import pytest

SKILL_DIR = Path(__file__).resolve().parents[1]
FIXTURE = SKILL_DIR / "tests" / "fixture.md"
BUILD_PY = SKILL_DIR / "build.py"


def test_plan_only_writes_sidecar_no_pptx(tmp_path):
    """--plan-only emits the JSON sidecar but no pptx."""
    md = tmp_path / "deck.md"
    md.write_text(FIXTURE.read_text())
    out = tmp_path / "out.pptx"
    proc = subprocess.run(
        [sys.executable, str(BUILD_PY),
         "--input", str(md), "--output", str(out),
         "--plan-only"],
        capture_output=True, text=True,
    )
    assert proc.returncode == 0, proc.stderr
    sidecar = md.with_suffix(md.suffix + ".layout.json")
    assert sidecar.exists()
    plan = json.loads(sidecar.read_text())
    assert plan["version"] == 1
    assert isinstance(plan["slides"], list)
    assert len(plan["slides"]) > 0
    # pptx not written
    assert not out.exists()


def test_default_plan_renders_pptx(tmp_path):
    """Without --no-plan, the v4 path renders pptx using the inferred plan."""
    md = tmp_path / "deck.md"
    md.write_text(FIXTURE.read_text())
    out = tmp_path / "out.pptx"
    proc = subprocess.run(
        [sys.executable, str(BUILD_PY),
         "--input", str(md), "--output", str(out)],
        capture_output=True, text=True,
    )
    assert proc.returncode == 0, proc.stderr
    assert out.exists()
    assert out.stat().st_size > 5000


def test_no_plan_falls_back_to_legacy(tmp_path):
    """--no-plan skips the sidecar and uses the v3 path."""
    md = tmp_path / "deck.md"
    md.write_text(FIXTURE.read_text())
    out = tmp_path / "out.pptx"
    proc = subprocess.run(
        [sys.executable, str(BUILD_PY),
         "--input", str(md), "--output", str(out),
         "--no-plan"],
        capture_output=True, text=True,
    )
    assert proc.returncode == 0, proc.stderr
    assert out.exists()
    sidecar = md.with_suffix(md.suffix + ".layout.json")
    assert not sidecar.exists()  # no sidecar in legacy mode


def test_shake_regenerates_sidecar(tmp_path):
    """--shake replaces an existing sidecar."""
    md = tmp_path / "deck.md"
    md.write_text(FIXTURE.read_text())
    sidecar = md.with_suffix(md.suffix + ".layout.json")
    # Plant a stale sidecar
    sidecar.write_text('{"version": 1, "deck_md_hash": "stale", "shake_seed": null, "slides": []}')
    out = tmp_path / "out.pptx"
    proc = subprocess.run(
        [sys.executable, str(BUILD_PY),
         "--input", str(md), "--output", str(out),
         "--shake"],
        capture_output=True, text=True,
    )
    assert proc.returncode == 0, proc.stderr
    plan = json.loads(sidecar.read_text())
    assert plan["deck_md_hash"] != "stale"
    assert len(plan["slides"]) > 0


def test_smoke_realistic_fixture_in_each_mode(tmp_path):
    """Run the realistic fixture through default, --shake, --plan-only, --no-plan
    and assert each produces sane output."""
    src = SKILL_DIR / "tests" / "fixture_realistic.md"
    md = tmp_path / "deck.md"
    md.write_text(src.read_text())
    sidecar = md.with_suffix(md.suffix + ".layout.json")

    # Default mode — sidecar absent
    out_default = tmp_path / "default.pptx"
    proc = subprocess.run(
        [sys.executable, str(BUILD_PY),
         "--input", str(md), "--output", str(out_default)],
        capture_output=True, text=True,
    )
    assert proc.returncode == 0, proc.stderr
    assert out_default.exists() and out_default.stat().st_size > 50_000
    assert sidecar.exists()

    # Default mode again — should hit cache (sidecar exists)
    out_cached = tmp_path / "cached.pptx"
    proc = subprocess.run(
        [sys.executable, str(BUILD_PY),
         "--input", str(md), "--output", str(out_cached)],
        capture_output=True, text=True,
    )
    assert proc.returncode == 0, proc.stderr
    assert out_cached.exists()

    # --shake — regenerates the plan
    sidecar_before = sidecar.read_text()
    out_shaken = tmp_path / "shaken.pptx"
    proc = subprocess.run(
        [sys.executable, str(BUILD_PY),
         "--input", str(md), "--output", str(out_shaken),
         "--shake"],
        capture_output=True, text=True,
    )
    assert proc.returncode == 0, proc.stderr
    assert out_shaken.exists()
    # Sidecar still exists; deck_md_hash should still match (markdown unchanged)
    sidecar_after = sidecar.read_text()
    import json
    assert json.loads(sidecar_after)["deck_md_hash"] == json.loads(sidecar_before)["deck_md_hash"]

    # --plan-only — emits JSON, no pptx
    out_plan_only = tmp_path / "plan_only.pptx"
    sidecar.unlink()  # remove sidecar to force fresh write
    proc = subprocess.run(
        [sys.executable, str(BUILD_PY),
         "--input", str(md), "--output", str(out_plan_only),
         "--plan-only"],
        capture_output=True, text=True,
    )
    assert proc.returncode == 0, proc.stderr
    assert sidecar.exists()
    assert not out_plan_only.exists()

    # --no-plan — legacy path, no sidecar written
    sidecar.unlink()
    out_legacy = tmp_path / "legacy.pptx"
    proc = subprocess.run(
        [sys.executable, str(BUILD_PY),
         "--input", str(md), "--output", str(out_legacy),
         "--no-plan"],
        capture_output=True, text=True,
    )
    assert proc.returncode == 0, proc.stderr
    assert out_legacy.exists()
    assert not sidecar.exists()


def test_render_is_deterministic(tmp_path):
    """Same plan → same pptx bytes (modulo timestamp metadata in pptx)."""
    md = tmp_path / "deck.md"
    md.write_text(FIXTURE.read_text())
    out1 = tmp_path / "out1.pptx"
    out2 = tmp_path / "out2.pptx"
    cmd_base = [sys.executable, str(BUILD_PY), "--input", str(md)]
    for out in (out1, out2):
        proc = subprocess.run(
            cmd_base + ["--output", str(out)],
            capture_output=True, text=True,
        )
        assert proc.returncode == 0, proc.stderr
    # Bytes won't be EXACTLY identical because pptx embeds a timestamp.
    # Instead, assert slide count + slide titles are the same.
    from pptx import Presentation
    p1 = Presentation(str(out1))
    p2 = Presentation(str(out2))
    assert len(p1.slides) == len(p2.slides)
    titles1 = []
    titles2 = []
    for s in p1.slides:
        for shp in s.shapes:
            if shp.has_text_frame:
                t = shp.text_frame.text.strip()
                if t:
                    titles1.append(t[:50]); break
    for s in p2.slides:
        for shp in s.shapes:
            if shp.has_text_frame:
                t = shp.text_frame.text.strip()
                if t:
                    titles2.append(t[:50]); break
    assert titles1 == titles2
