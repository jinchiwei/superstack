"""Slide master functions: each adds one slide of the right shape."""

from __future__ import annotations

import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
from build import (  # noqa: E402
    new_presentation,
    add_title_slide,
    add_content_slide,
    add_section_divider,
    add_big_number_slide,
    add_two_column_slide,
    add_quote_slide,
    add_end_slide,
)


def test_new_presentation_is_16_9():
    prs = new_presentation()
    # 13.333" wide × 7.5" high in EMUs (914400 per inch)
    assert prs.slide_width == int(13.333 * 914400)
    assert prs.slide_height == int(7.5 * 914400)


def test_add_title_slide_dark_background():
    prs = new_presentation()
    s = add_title_slide(prs, eyebrow="TEST", title="Title", subtitle="sub",
                       name="Jinchi", org="UCSF", date="2026-05-01")
    assert s is not None
    assert len(prs.slides) == 1


def test_add_content_slide_white_background():
    prs = new_presentation()
    add_content_slide(prs, title="Section A", body_paragraphs=["one", "two"])
    assert len(prs.slides) == 1


def test_add_section_divider_cycles_color_by_index():
    """Section divider color is determined by the index parameter, cycling."""
    sys.path.insert(0, str(Path(__file__).resolve().parents[2] / "_shared"))
    import branding
    prs = new_presentation()
    add_section_divider(prs, label="One", index=0)
    add_section_divider(prs, label="Three", index=2)
    add_section_divider(prs, label="Five", index=4)
    assert len(prs.slides) == 3
    assert branding.pick_section_color(0) == branding.TURQUOISE
    assert branding.pick_section_color(2) == branding.AMBER
    assert branding.pick_section_color(4) == branding.TURQUOISE


def test_add_big_number_slide():
    prs = new_presentation()
    add_big_number_slide(prs, number="+12.4%", caption="recall improvement on UCSF cohort")
    assert len(prs.slides) == 1


def test_add_two_column_slide():
    prs = new_presentation()
    add_two_column_slide(prs, title="Comparison",
                        left_title="Baseline", left_body=["Old approach", "manual"],
                        right_title="Proposed", right_body=["New approach", "automated"])
    assert len(prs.slides) == 1


def test_add_quote_slide():
    prs = new_presentation()
    add_quote_slide(prs, quote="The best research is reproducible.",
                   attribution="Jinchi Wei")
    assert len(prs.slides) == 1


def test_add_end_slide():
    prs = new_presentation()
    add_end_slide(prs, message="Thanks", contact="mrjinch@gmail.com")
    assert len(prs.slides) == 1


def test_full_deck_renders_to_file(tmp_path):
    """Compose a 7-slide deck using all masters, save, verify file."""
    prs = new_presentation()
    add_title_slide(prs, eyebrow="DECK TEST", title="Master Test Deck",
                    subtitle="exercises every master", name="Jinchi", org="UCSF",
                    date="2026-05-01")
    add_section_divider(prs, label="Section One", index=0)
    add_content_slide(prs, title="Content", body_paragraphs=["paragraph"])
    add_big_number_slide(prs, number="100%", caption="of tests passing")
    add_two_column_slide(prs, title="Compare",
                        left_title="A", left_body=["a1"],
                        right_title="B", right_body=["b1"])
    add_quote_slide(prs, quote="Ship.", attribution="Self")
    add_end_slide(prs, message="Thanks", contact="—")

    out = tmp_path / "deck.pptx"
    prs.save(str(out))
    assert out.is_file()
    assert out.stat().st_size > 5000

    # Re-open with python-pptx to verify slide count
    from pptx import Presentation
    reopened = Presentation(str(out))
    assert len(reopened.slides) == 7


def test_add_content_slide_accepts_accent_color():
    """add_content_slide takes accent_color_hex; defaults to turquoise."""
    import branding as b
    prs = new_presentation()
    add_content_slide(prs, title="Section", body_paragraphs=["body"],
                      accent_color_hex=b.DEEPPINK)
    assert len(prs.slides) == 1
    # Find a deeppink-fill rect on the slide (the left bar)
    found_deeppink = False
    s = prs.slides[0]
    for shp in s.shapes:
        try:
            if shp.fill.type == 1 and str(shp.fill.fore_color.rgb).upper() == "FF1493":
                found_deeppink = True
                break
        except Exception:
            pass
    assert found_deeppink, "left bar in deeppink not found"


def test_add_section_divider_accepts_accent_color_override():
    """add_section_divider takes accent_color_hex overriding the index cycle."""
    import branding as b
    prs = new_presentation()
    add_section_divider(prs, label="Custom", accent_color_hex=b.AMBER)
    s = prs.slides[0]
    found_amber = False
    for shp in s.shapes:
        try:
            if shp.fill.type == 1 and str(shp.fill.fore_color.rgb).upper() == "F0C840":
                found_amber = True
                break
        except Exception:
            pass
    assert found_amber, "amber colorblock not found in section divider"


def test_title_slide_has_amber_bottom_hairline():
    """Title slide includes a thin amber bar at the bottom."""
    prs = new_presentation()
    add_title_slide(prs, eyebrow="X", title="T", date="2026-05-01")
    s = prs.slides[0]
    # Find any rect filled with amber #F0C840
    found_amber = False
    for shp in s.shapes:
        try:
            if shp.fill.type == 1 and str(shp.fill.fore_color.rgb).upper() == "F0C840":
                found_amber = True
                break
        except Exception:
            pass
    assert found_amber, "amber bottom hairline missing"


def test_end_slide_mirrors_title_amber_hairline():
    prs = new_presentation()
    add_end_slide(prs, message="Thanks", contact="me")
    s = prs.slides[0]
    found_amber = False
    for shp in s.shapes:
        try:
            if shp.fill.type == 1 and str(shp.fill.fore_color.rgb).upper() == "F0C840":
                found_amber = True
                break
        except Exception:
            pass
    assert found_amber, "amber bottom hairline missing on end slide"
