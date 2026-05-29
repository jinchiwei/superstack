"""Tests for the expressive-freeform composer (Option B).

A fresh expressive build (no agent in the loop) must emit DESIGNED freeform
content slides — not named layouts. Strict mode must still emit named layouts.
Agent-authored freeform sidecar entries must survive a plain re-render.
"""

import sys
import json
import subprocess
from collections import Counter
from pathlib import Path

SKILL_DIR = Path(__file__).resolve().parents[1]
BUILD_PY = SKILL_DIR / "build.py"
FIXTURE_REALISTIC = SKILL_DIR / "tests" / "fixture_realistic.md"
FIXTURE_FREEFORM = SKILL_DIR / "tests" / "fixture_freeform_demo.md"

_BLOCK_KINDS = {"freeform", "composition"}


def _build(md, out, *extra):
    proc = subprocess.run(
        [sys.executable, str(BUILD_PY), "--allow-composed", "--input", str(md),
         "--output", str(out), *extra],
        capture_output=True, text=True,
    )
    assert proc.returncode == 0, proc.stderr
    return proc


def _kinds(sidecar):
    plan = json.loads(sidecar.read_text())
    return plan, Counter(s["kind"] for s in plan["slides"])


def _content_slides(plan):
    return [s for s in plan["slides"] if s["kind"] != "section-divider"]


def test_expressive_default_yields_freeform_content_slides(tmp_path):
    """A fresh expressive build (no sidecar) rewrites every content slide to
    freeform; section-dividers stay dividers."""
    md = tmp_path / "deck.md"
    md.write_text(FIXTURE_REALISTIC.read_text())
    out = tmp_path / "out.pptx"
    _build(md, out)  # expressive is the default
    sidecar = md.with_suffix(md.suffix + ".layout.json")
    plan, counts = _kinds(sidecar)
    assert plan["mode"] == "expressive"
    # Every non-divider slide is freeform.
    content = _content_slides(plan)
    assert content, "fixture should have content slides"
    assert all(s["kind"] == "freeform" for s in content), counts
    # Dividers are preserved.
    assert counts.get("section-divider", 0) >= 1
    # No named content layouts leaked through.
    named = {"content-text", "cards-grid", "table-with-takeaway",
             "figure-with-aside", "conclusions", "content-text-image"}
    assert not (named & set(counts)), counts


def test_expressive_freeform_has_zero_error_chips(tmp_path):
    """The rendered pptx must contain no freeform error chips — every emitted
    snippet validates + executes cleanly in the sandbox."""
    from pptx import Presentation

    md = tmp_path / "deck.md"
    md.write_text(FIXTURE_REALISTIC.read_text())
    out = tmp_path / "out.pptx"
    _build(md, out)
    prs = Presentation(str(out))
    chips = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and "[freeform" in shape.text_frame.text:
                chips.append((i, shape.text_frame.text))
    assert not chips, chips


def test_strict_mode_keeps_named_layouts(tmp_path):
    """Strict mode must skip the composer entirely — named layouts, no
    freeform."""
    md = tmp_path / "deck.md"
    md.write_text(FIXTURE_REALISTIC.read_text())
    out = tmp_path / "out.pptx"
    _build(md, out, "--mode", "strict")
    sidecar = md.with_suffix(md.suffix + ".layout.json")
    plan, counts = _kinds(sidecar)
    assert plan["mode"] == "strict"
    assert not (_BLOCK_KINDS & set(counts)), counts
    # And there ARE named content layouts.
    content = _content_slides(plan)
    assert content and all(s["kind"] not in _BLOCK_KINDS for s in content)


def test_agent_authored_freeform_sidecar_is_preserved(tmp_path):
    """A plain expressive re-render over an agent-authored freeform sidecar
    must NOT overwrite the hand-written freeform code."""
    md = tmp_path / "deck.md"
    md.write_text(FIXTURE_FREEFORM.read_text())
    sidecar = md.with_suffix(md.suffix + ".layout.json")
    sidecar.write_text(
        (FIXTURE_FREEFORM.with_suffix(".md.layout.json")).read_text()
    )
    # Snapshot the agent-authored freeform code by slide_id.
    before = {
        s["slide_id"]: s["params"].get("code", "")
        for s in json.loads(sidecar.read_text())["slides"]
        if s["kind"] == "freeform"
    }
    assert before, "fixture should be agent-authored freeform"
    out = tmp_path / "out.pptx"
    _build(md, out, "--mode", "expressive")
    after = {
        s["slide_id"]: s["params"].get("code", "")
        for s in json.loads(sidecar.read_text())["slides"]
        if s["kind"] == "freeform"
    }
    for sid, code in before.items():
        assert after.get(sid) == code, f"freeform code for {sid} was overwritten"


def test_compose_unit_rewrites_named_entry():
    """Unit-level: compose_expressive_plan rewrites a named entry to freeform
    with parseable code and leaves a divider alone."""
    import ast
    sys.path.insert(0, str(SKILL_DIR))
    from plan import SlideEntry
    from expressive_compose import compose_expressive_plan

    slides = [
        SlideEntry(slide_id="divider-x", kind="section-divider",
                   params={"label": "Sec"}, content_hash="d"),
        SlideEntry(slide_id="x", kind="content-text",
                   params={"title": "T", "lede": "L",
                           "body": [{"kind": "bullet", "html": "one"},
                                    {"kind": "bullet", "html": "two"}]},
                   content_hash="c"),
    ]
    n = compose_expressive_plan(slides)
    assert n == ["x"]  # returns the list of composed slide_ids (the floor)
    assert slides[1].params.get("_provenance") == "composer"
    assert slides[0].kind == "section-divider"
    assert slides[1].kind == "freeform"
    assert "code" in slides[1].params
    ast.parse(slides[1].params["code"])  # emitted code must parse
