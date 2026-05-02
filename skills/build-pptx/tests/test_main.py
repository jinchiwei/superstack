"""build-pptx end-to-end markdown→pptx test."""

from __future__ import annotations

import subprocess
import sys
from pathlib import Path

import pytest

SKILL_DIR = Path(__file__).resolve().parents[1]
FIXTURE = SKILL_DIR / "tests" / "fixture.md"
BUILD_PY = SKILL_DIR / "build.py"


def _render(out_path: Path, *extra_args: str) -> None:
    cmd = [sys.executable, str(BUILD_PY),
           "--input", str(FIXTURE),
           "--output", str(out_path),
           *extra_args]
    proc = subprocess.run(cmd, capture_output=True, text=True)
    assert proc.returncode == 0, f"build.py failed:\nSTDERR: {proc.stderr}"


def test_renders_pptx(tmp_path):
    out = tmp_path / "out.pptx"
    _render(out)
    assert out.is_file()
    assert out.stat().st_size > 5000


def test_pptx_starts_with_zip_magic_bytes(tmp_path):
    """PPTX is a zip; starts with PK."""
    out = tmp_path / "out.pptx"
    _render(out)
    assert out.read_bytes()[:2] == b"PK"


def test_pptx_has_correct_slide_count(tmp_path):
    """Title slide + 3 content slides + end slide = 5 slides."""
    out = tmp_path / "out.pptx"
    _render(out)
    from pptx import Presentation
    prs = Presentation(str(out))
    assert len(prs.slides) == 5


def test_pptx_slide_titles_match_fixture(tmp_path):
    """Slide 1 = title, slides 2-4 = Background/Methods/Results, slide 5 = end."""
    out = tmp_path / "out.pptx"
    _render(out)
    from pptx import Presentation
    prs = Presentation(str(out))
    titles = []
    for slide in prs.slides:
        for shp in slide.shapes:
            if shp.has_text_frame:
                titles.extend(p.text for p in shp.text_frame.paragraphs if p.text.strip())
    joined = "\n".join(titles)
    assert "Build-PPTX Smoke Test" in joined
    assert "Background" in joined
    assert "Methods" in joined
    assert "Results" in joined


def test_no_cover_flag(tmp_path):
    """--no-cover suppresses title slide → 4 slides instead of 5."""
    out = tmp_path / "out.pptx"
    _render(out, "--no-cover")
    from pptx import Presentation
    prs = Presentation(str(out))
    assert len(prs.slides) == 4


def test_no_end_flag(tmp_path):
    """--no-end suppresses closing slide → 4 slides instead of 5."""
    out = tmp_path / "out.pptx"
    _render(out, "--no-end")
    from pptx import Presentation
    prs = Presentation(str(out))
    assert len(prs.slides) == 4


def test_section_color_cascade(tmp_path):
    """Each content slide's left bar matches its section's expected color."""
    fixture = SKILL_DIR / "tests" / "fixture_sections.md"
    out = tmp_path / "cascade.pptx"
    cmd = [sys.executable, str(BUILD_PY),
           "--input", str(fixture),
           "--output", str(out)]
    proc = subprocess.run(cmd, capture_output=True, text=True)
    assert proc.returncode == 0, proc.stderr

    from pptx import Presentation
    prs = Presentation(str(out))
    # Skip title slide (idx 0) and end slide (last). Slides between are:
    # divider(Background, turquoise) → content(Motivation slide, turquoise)
    # divider(Methods, deeppink)     → content(Pipeline overview, deeppink)
    # divider(Results, amber)        → content(Headline numbers, amber)
    # divider(Limitations, blueviolet) → content(What we couldn't measure, blueviolet)
    # divider(Conclusions, turquoise) → content(Next steps, turquoise)
    expected = [
        # (slide_idx, expected_brand_hex)
        (1, "40E0D0"),  # Background section divider
        (2, "40E0D0"),  # Motivation content slide
        (3, "FF1493"),  # Methods divider
        (4, "FF1493"),  # Pipeline content
        (5, "F0C840"),  # Results divider
        (6, "F0C840"),  # Headline content
        (7, "8A2BE2"),  # Limitations divider
        (8, "8A2BE2"),  # What we couldn't measure content
        (9, "40E0D0"),  # Conclusions divider
        (10, "40E0D0"), # Next steps content
    ]
    for slide_idx, hex_color in expected:
        slide = prs.slides[slide_idx]
        # Find any solid-filled rect on the slide that uses this brand color
        found = False
        for shp in slide.shapes:
            try:
                if shp.fill.type == 1 and str(shp.fill.fore_color.rgb).upper() == hex_color:
                    found = True
                    break
            except Exception:
                pass
        assert found, f"slide {slide_idx} missing expected brand color #{hex_color}"
