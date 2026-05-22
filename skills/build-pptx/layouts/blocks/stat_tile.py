"""stat-tile block — big number + label + optional sub + optional icon.

params:
    value   str     — the main statistic (e.g. "0.91", "120")
    label   str     — descriptor below value (e.g. "Internal AUC")
    sub     str     — optional smaller sub-label (e.g. "5-seed mean")
    icon    str     — optional FA icon name (e.g. "FaChartLine")
    icon_path str   — optional raw icon path fallback
"""
from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.util import Inches

from ._base import (
    _add_text, _add_rect,
    PAPER_RGB, MUTED_RGB, DIM_RGB,
    branding,
    _resolve_icon_path, _rgb_to_hex,
)
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def render(slide, *, left: float, top: float, width: float, height: float,
           params: dict, accent_rgb: RGBColor,
           surface_rgb: RGBColor | None = None,
           text_rgb: RGBColor | None = None,
           muted_rgb: RGBColor | None = None) -> None:
    value = str(params.get("value", ""))
    label = str(params.get("label", ""))
    sub = str(params.get("sub", ""))
    icon_name = params.get("icon")
    icon_path_str = params.get("icon_path")

    # Tile bg + secondary text — palette overrides fall back to today's
    # PAPER / MUTED / DIM colors. (The big value stays on the brand accent.)
    tile_fill = surface_rgb if surface_rgb is not None else PAPER_RGB
    label_color = muted_rgb if muted_rgb is not None else MUTED_RGB
    sub_color = muted_rgb if muted_rgb is not None else DIM_RGB

    # Background (paper)
    _add_rect(slide, left=left, top=top, width=width, height=height,
              fill_rgb=tile_fill)

    # Build accent hex for icon color
    accent_hex = _rgb_to_hex(accent_rgb)

    # Optional icon — top-right corner
    icon_size = 0.38
    if icon_name or icon_path_str:
        icon_path = _resolve_icon_path(icon_name, icon_path_str, accent_hex, size_px=128)
        if icon_path is not None:
            try:
                slide.shapes.add_picture(
                    str(icon_path),
                    Inches(left + width - icon_size - 0.12),
                    Inches(top + 0.12),
                    width=Inches(icon_size),
                    height=Inches(icon_size),
                )
            except Exception:
                pass

    # Value — big mono accent number, vertically centered in the upper 60%
    value_size = _pick_value_size(len(value))
    value_h = min(height * 0.55, 1.2)
    value_top = top + height * 0.08

    tb = slide.shapes.add_textbox(
        Inches(left + 0.12), Inches(value_top),
        Inches(width - 0.24), Inches(value_h),
    )
    tf = tb.text_frame
    tf.word_wrap = False
    from pptx.util import Emu, Pt
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = value
    r.font.name = branding.MONO_FONT
    r.font.size = Pt(value_size)
    r.font.color.rgb = accent_rgb
    r.font.bold = True

    # Label — 12pt sans muted, below value
    if label:
        label_top = value_top + value_h + 0.03
        label_h = 0.32
        _add_text(slide, label,
                  left=left + 0.10, top=label_top,
                  width=width - 0.20, height=label_h,
                  size=12, color_rgb=label_color, font=branding.SANS_FONT,
                  align=PP_ALIGN.CENTER)

        if sub:
            sub_top = label_top + label_h
            _add_text(slide, sub,
                      left=left + 0.10, top=sub_top,
                      width=width - 0.20, height=0.26,
                      size=10, color_rgb=sub_color, font=branding.MONO_FONT,
                      align=PP_ALIGN.CENTER)


def _pick_value_size(char_count: int) -> int:
    """Scale value font size based on length to fit the tile."""
    if char_count <= 4:
        return 44
    if char_count <= 6:
        return 36
    if char_count <= 10:
        return 28
    return 22
