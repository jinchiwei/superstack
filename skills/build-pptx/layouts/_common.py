"""Shared helpers, colour constants, and parse utilities for the layouts package.

All geometry primitives (_add_text, _add_rect, _set_bg, _blank, _add_card,
_render_paragraph_block, _render_media_block), the HTML-strip/parse helpers,
the colour constants, and the slide-chrome helper (_add_chrome) live here.
Individual layout modules import from this module only — they never reach into
build.py.
"""

from __future__ import annotations

import math
import re
import sys
from pathlib import Path

# Wire imports to sibling _shared/
sys.path.insert(0, str(Path(__file__).resolve().parents[2] / "_shared"))

import branding  # noqa: E402

try:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Emu, Inches, Pt
except ImportError:
    raise SystemExit("python-pptx not installed. Run: pip install python-pptx")

# === Regex constants ===
_HTML_TAG_RE = re.compile(r"<[^>]+>")
_WS_RUN_RE = re.compile(r"\s+")

# === Colour helpers ===
def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# Frequently-used RGBColor instances
INK_RGB        = _rgb(branding.INK)
WHITE_RGB      = _rgb(branding.WHITE)
TURQUOISE_RGB  = _rgb(branding.TURQUOISE)
DEEPPINK_RGB   = _rgb(branding.DEEPPINK)
AMBER_RGB      = _rgb(branding.AMBER)
BLUEVIOLET_RGB = _rgb(branding.BLUEVIOLET)
DIM_RGB        = _rgb(branding.DIM)
MUTED_RGB      = _rgb(branding.MUTED)
RULE_RGB       = _rgb(branding.RULE)
DARK_BG_RGB    = _rgb(branding.DARK_BG)
PAPER_RGB      = _rgb(branding.PAPER)


# === Parse constants ===
_DEFLIST_LABEL_MAX_LEN = 80   # bold prefix qualifies as definition label
_DEFLIST_BODY_MAX_LEN  = 350  # too-long body suggests prose, not a card


# === HTML helpers ===
def _strip_html(text: str) -> str:
    """Remove HTML tags AND collapse any whitespace runs to single spaces."""
    return _WS_RUN_RE.sub(" ", _HTML_TAG_RE.sub("", text)).strip()


def _split_slides(body_html: str) -> list[str]:
    """Split rendered body HTML on <hr> elements."""
    parts = re.split(r"<hr\s*/?>", body_html)
    return [p.strip() for p in parts if p.strip()]


def _parse_table(html_table: str) -> list[list[str]]:
    """Parse a rendered HTML <table> into a list of rows of cell strings."""
    rows = []
    for tr_match in re.finditer(r"<tr[^>]*>(.*?)</tr>", html_table, re.DOTALL):
        cells = []
        for cell_match in re.finditer(r"<t[hd][^>]*>(.*?)</t[hd]>",
                                       tr_match.group(1), re.DOTALL):
            cells.append(_strip_html(cell_match.group(1)))
        if cells:
            rows.append(cells)
    return rows


def _detect_def_cards_from_li_html(li_blocks: list[str]) -> list[dict] | None:
    """Detect a definition-list pattern in <li> blocks; return cards or None."""
    cards = []
    for raw in li_blocks:
        m = re.match(r"^\s*<strong>(.+?)</strong>\s*", raw, re.DOTALL)
        if not m:
            return None
        label = _strip_html(m.group(1)).strip().rstrip(":.")
        if len(label) == 0 or len(label) > _DEFLIST_LABEL_MAX_LEN:
            return None
        rest = raw[m.end():]
        rest = re.sub(r"^\s*[:—\-–]\s*", "", rest)
        body = _strip_html(rest).strip()
        if len(body) == 0 or len(body) > _DEFLIST_BODY_MAX_LEN:
            return None
        cards.append({"label": label, "body": body})
    return cards if len(cards) >= 2 else None


def _parse_slide_chunk(html_chunk: str, *, base_dir: Path | None = None) -> dict:
    """Extract slide title, body paragraphs, images, tables, and cards."""
    title_match = re.search(r"<(h[12])[^>]*>(.*?)</\1>", html_chunk)
    if title_match:
        title = _strip_html(title_match.group(2))
        rest = html_chunk[title_match.end():]
    else:
        title = ""
        rest = html_chunk

    tables: list[list[list[str]]] = []
    for tbl_match in re.finditer(r"<table[^>]*>.*?</table>", rest, re.DOTALL):
        parsed = _parse_table(tbl_match.group(0))
        if parsed:
            tables.append(parsed)
    rest_no_tables = re.sub(r"<table[^>]*>.*?</table>", "", rest, flags=re.DOTALL)

    images: list[Path] = []
    image_records: list[dict] = []
    for img_match in re.finditer(r'<img\s+([^>]+)>', rest_no_tables):
        attrs = img_match.group(1)
        src_m = re.search(r'src="([^"]+)"', attrs)
        if not src_m:
            continue
        src = src_m.group(1)
        alt_m = re.search(r'alt="([^"]*)"', attrs)
        alt = alt_m.group(1) if alt_m else ""
        path = Path(src)
        if not path.is_absolute() and base_dir is not None:
            path = (base_dir / path).resolve()
        if path.exists():
            images.append(path)
            image_records.append({"path": path, "alt": alt,
                                  "pos": img_match.start()})
    rest_no_media = re.sub(r"<img[^>]*/?>", "", rest_no_tables)

    cards: list[dict] = []
    h3_pattern = re.compile(r"<h3[^>]*>(.*?)</h3>", re.DOTALL)
    h3_positions = [(m.start(), m.end(), m.group(1))
                    for m in h3_pattern.finditer(rest_no_tables)]
    for i, (start, end, h3_inner) in enumerate(h3_positions):
        icon_path = None
        img_m = re.search(r'<img[^>]+src="([^"]+)"', h3_inner)
        if img_m:
            src = img_m.group(1)
            p = Path(src)
            if not p.is_absolute() and base_dir is not None:
                p = (base_dir / p).resolve()
            if p.exists():
                icon_path = p
            h3_inner_no_img = re.sub(r'<img[^>]*/?>', '', h3_inner)
        else:
            h3_inner_no_img = h3_inner
        label = _strip_html(h3_inner_no_img)

        next_start = h3_positions[i + 1][0] if i + 1 < len(h3_positions) else len(rest_no_tables)
        block = rest_no_tables[end:next_start]
        block_no_img = re.sub(r'<img[^>]*/?>', '', block)
        body_lines = []
        for m in re.finditer(r"<(p|li)\b[^>]*>(.*?)</\1>", block_no_img, re.DOTALL):
            text = _strip_html(m.group(2)).strip()
            if text:
                body_lines.append(text)
        cards.append({"label": label, "body": " ".join(body_lines),
                      "icon": icon_path})

    body_region = rest_no_media[:h3_positions[0][0]] if h3_positions else rest_no_media

    li_blocks = [m.group(1) for m in
                 re.finditer(r"<li[^>]*>(.*?)</li>", body_region, re.DOTALL)]
    auto_cards = None
    if li_blocks and not cards:
        auto_cards = _detect_def_cards_from_li_html(li_blocks)

    paragraphs: list[dict] = []
    if auto_cards:
        cards = auto_cards
        for m in re.finditer(r"<p\b[^>]*>(.*?)</p>", body_region, re.DOTALL):
            html = m.group(1).strip()
            if _strip_html(html):
                paragraphs.append({"kind": "paragraph", "html": html,
                                   "pos": m.start()})
    else:
        for m in re.finditer(r"<(p|li)\b[^>]*>(.*?)</\1>", body_region, re.DOTALL):
            html = m.group(2).strip()
            if _strip_html(html):
                kind = "bullet" if m.group(1) == "li" else "paragraph"
                paragraphs.append({"kind": kind, "html": html,
                                   "pos": m.start()})

    return {"title": title, "body": paragraphs, "images": images,
            "tables": tables, "cards": cards,
            "image_records": image_records}


# === Primitive drawing helpers ===
def _set_bg(slide, color_rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb


def _add_rect(slide, *, left, top, width, height,
              fill_rgb: RGBColor | None = None) -> None:
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    shp.shadow.inherit = False
    if fill_rgb is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_rgb
    else:
        shp.fill.background()
    shp.line.fill.background()


def _add_text(slide, text, *, left, top, width, height, size=18,
              color_rgb: RGBColor = INK_RGB, font: str = branding.SANS_FONT,
              bold=False, italic=False, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color_rgb
    r.font.bold = bold
    r.font.italic = italic
    return tb


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_card(slide, *, label: str, body: str, left: float, top: float,
              width: float, height: float, accent_rgb: RGBColor,
              icon_path: Path | None = None) -> None:
    """A bordered tile (paper bg + thin accent top stripe) with label + body."""
    _add_rect(slide, left=left, top=top, width=width, height=height,
              fill_rgb=PAPER_RGB)
    _add_rect(slide, left=left, top=top, width=width, height=0.06,
              fill_rgb=accent_rgb)
    if icon_path is not None and icon_path.exists():
        try:
            slide.shapes.add_picture(
                str(icon_path),
                Inches(left + 0.15), Inches(top + 0.18),
                width=Inches(0.32), height=Inches(0.32),
            )
            label_left = left + 0.58
            label_w = width - 0.76
        except Exception:
            label_left = left + 0.18
            label_w = width - 0.36
    else:
        label_left = left + 0.18
        label_w = width - 0.36
    _add_text(slide, label, left=label_left, top=top + 0.18,
              width=label_w, height=0.4,
              size=13, color_rgb=accent_rgb, font=branding.MONO_FONT, bold=True)
    _add_text(slide, body, left=left + 0.18, top=top + 0.65,
              width=width - 0.36, height=height - 0.75,
              size=12, color_rgb=INK_RGB, font=branding.SANS_FONT)


def _strip_html_keep_edges(text: str) -> str:
    """Like _strip_html but does NOT strip leading/trailing whitespace.
    Used inside _add_runs_from_html where edge spaces between adjacent runs
    are meaningful — e.g., 'run <code>foo</code> on' must preserve the
    spaces around the inline code so the output isn't 'runfooon'."""
    return _WS_RUN_RE.sub(" ", _HTML_TAG_RE.sub("", text))


def _add_runs_from_html(p, *, html_text: str, size: float,
                        text_color: RGBColor | None = None) -> None:
    """Append runs to paragraph p parsing inline <strong>/<em>/<code> spans.

    Edge whitespace between runs is preserved (we use _strip_html_keep_edges
    instead of _strip_html) so spaces around inline tags survive — without
    this, 'run <code>foo</code> on' renders as 'runfooon'.

    text_color: optional override for run colour; defaults to INK_RGB.
    """
    run_color = text_color if text_color is not None else INK_RGB
    parts = re.split(r"(<strong>.*?</strong>|<em>.*?</em>|<code>.*?</code>)",
                     html_text, flags=re.DOTALL)
    for part in parts:
        if not part:
            continue
        bold = italic = mono = False
        if part.startswith("<strong>"):
            bold = True
        elif part.startswith("<em>"):
            italic = True
        elif part.startswith("<code>"):
            mono = True
        text = _strip_html_keep_edges(part)
        if not text:
            continue
        r = p.add_run()
        r.text = text
        r.font.name = branding.MONO_FONT if mono else branding.SANS_FONT
        r.font.size = Pt(size)
        r.font.color.rgb = run_color
        r.font.bold = bold
        r.font.italic = italic


def _estimate_paragraph_height(text: str, *, width: float, size: float,
                               line_spacing: float = 1.35) -> float:
    """Rough estimate of rendered paragraph height in inches."""
    if not text:
        return 0.0
    char_w = size * 0.0095
    chars_per_line = max(1, int(width / char_w))
    n_lines = max(1, math.ceil(len(text) / chars_per_line))
    line_h = size * line_spacing / 72.0
    return n_lines * line_h * 1.15


def _render_paragraph_block(slide, *, items: list, left: float, top: float,
                            width: float, height: float, accent_rgb: RGBColor,
                            size: float = 14, distribute: bool = False,
                            text_color: RGBColor | None = None) -> None:
    """Render mixed paragraphs/bullets into one textbox.

    text_color: override the default INK_RGB for body text runs. Pass e.g.
        WHITE_RGB for dark-background slides. Bullet markers always use
        accent_rgb.
    """
    if not items:
        return
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)

    extra_pt = 0.0
    if distribute and len(items) > 1:
        line_spacing = 1.35
        total_h = 0.0
        for it in items:
            html = it.get("html", "") if isinstance(it, dict) else str(it).lstrip("• ").strip()
            text = _strip_html(html)
            total_h += _estimate_paragraph_height(text, width=width,
                                                  size=size,
                                                  line_spacing=line_spacing)
        total_h += (len(items) - 1) * (8 / 72.0)
        slack = max(0.0, height - total_h)
        if slack > 0:
            slack_pt = slack * 72.0
            extra_pt = min(24.0, slack_pt / (len(items) - 1))

    first = True
    for item in items:
        if isinstance(item, str):
            is_bullet = item.startswith("•  ") or item.startswith("• ")
            text = item.lstrip("• ").strip()
            item = {"kind": "bullet" if is_bullet else "paragraph",
                    "html": text}

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(8 + extra_pt)
        p.line_spacing = 1.35

        if item["kind"] == "bullet":
            r = p.add_run()
            r.text = "▸  "
            r.font.name = branding.MONO_FONT
            r.font.size = Pt(size)
            r.font.color.rgb = accent_rgb
            r.font.bold = True

        _add_runs_from_html(p, html_text=item["html"], size=size,
                            text_color=text_color)


def _get_image_aspect(path: Path) -> float:
    """Return image aspect ratio (width/height). Defaults to 1.0 if unreadable."""
    try:
        from PIL import Image
        with Image.open(str(path)) as im:
            return im.width / im.height if im.height else 1.0
    except Exception:
        return 1.0


def _add_table(slide, *, rows: list[list[str]], left: float, top: float,
               width: float, max_height: float, header_rgb: RGBColor) -> float:
    """Draw a table with a colored header row. Returns bottom y."""
    if not rows:
        return top
    n_rows = len(rows)
    n_cols = max(len(r) for r in rows)
    row_h = min(0.5, max_height / max(n_rows, 1))
    height = row_h * n_rows
    tbl_shape = slide.shapes.add_table(n_rows, n_cols,
                                       Inches(left), Inches(top),
                                       Inches(width), Inches(height))
    tbl = tbl_shape.table
    for i, row in enumerate(rows):
        for j in range(n_cols):
            cell = tbl.cell(i, j)
            text = row[j] if j < len(row) else ""
            cell.fill.solid()
            if i == 0:
                cell.fill.fore_color.rgb = header_rgb
                color = WHITE_RGB
                font_name = branding.MONO_FONT
                bold = True
                size = 12
            else:
                cell.fill.fore_color.rgb = PAPER_RGB if i % 2 == 0 else WHITE_RGB
                color = INK_RGB
                font_name = branding.SANS_FONT
                bold = False
                size = 12
            cell.margin_left = Emu(60000)
            cell.margin_right = Emu(60000)
            cell.margin_top = Emu(30000)
            cell.margin_bottom = Emu(30000)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = ""
            r = p.add_run()
            r.text = text
            r.font.name = font_name
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
    return top + height


def _render_media_block(slide, *, images: list[Path], tables: list[list[list[str]]],
                        left: float, top: float, width: float, height: float,
                        accent: RGBColor) -> None:
    """Render tables and images. Returns nothing (side-effects only)."""
    cursor = top
    for ti, tbl in enumerate(tables):
        header = accent if ti == 0 else INK_RGB
        max_h = min(2.5, (height / max(len(tables), 1)))
        cursor = _add_table(slide, rows=tbl, left=left, top=cursor,
                            width=width, max_height=max_h,
                            header_rgb=header)
        cursor += 0.2

    if not images:
        return
    remaining = (top + height) - cursor
    if remaining < 0.5:
        return

    n = len(images)
    if n == 1:
        try:
            pic = slide.shapes.add_picture(str(images[0]),
                                           Inches(left), Inches(cursor),
                                           width=Inches(width))
            pw = pic.width / 914400
            ph = pic.height / 914400
            if ph > remaining:
                new_h = max(0.5, remaining)
                new_w = max(0.5, width * (new_h / ph))
                pic.width = Inches(new_w)
                pic.height = Inches(new_h)
                pic.left = Inches(left + (width - new_w) / 2)
        except Exception as e:
            _add_text(slide, f"[image error: {Path(str(images[0])).name}: {e}]",
                      left=left, top=cursor, width=width, height=0.4,
                      size=10, color_rgb=DIM_RGB, font=branding.MONO_FONT)
    else:
        cols = 2 if n <= 4 else 3
        rows = (n + cols - 1) // cols
        gutter = 0.15
        sub_w = (width - gutter * (cols - 1)) / cols
        sub_h = (remaining - gutter * (rows - 1)) / rows
        sub_h = max(sub_h, 0.5)
        for i, img in enumerate(images):
            r, c = divmod(i, cols)
            x = left + c * (sub_w + gutter)
            y = cursor + r * (sub_h + gutter)
            try:
                pic = slide.shapes.add_picture(str(img),
                                               Inches(x), Inches(y),
                                               width=Inches(sub_w))
                pw = pic.width / 914400
                ph = pic.height / 914400
                if ph > sub_h:
                    new_h = max(0.5, sub_h)
                    new_w = max(0.5, sub_w * (new_h / ph))
                    pic.width = Inches(new_w)
                    pic.height = Inches(new_h)
                    pic.left = Inches(x + (sub_w - new_w) / 2)
            except Exception as e:
                _add_text(slide, f"[image error: {img.name}: {e}]",
                          left=x, top=y, width=sub_w, height=0.4,
                          size=9, color_rgb=DIM_RGB, font=branding.MONO_FONT)


# === Slide chrome helper ===
def _add_chrome(slide, *, title: str, lede: str, footer_kwargs: dict,
                accent: RGBColor,
                title_present: bool, title_wraps: bool,
                use_side_by_side: bool) -> tuple[float, float, float, float, float]:
    """Draw the standard content-slide chrome: left bar, title, hairline,
    lede (when not side-by-side), and footer.

    Returns (body_top, body_h, body_l, body_w, body_bottom) geometry tuple
    so each layout can position its body content correctly.
    """
    # Left accent bar
    _add_rect(slide, left=0, top=0, width=0.22, height=7.5, fill_rgb=accent)

    title_h = 1.05 if title_wraps else 0.55
    hairline_top = 0.30 + title_h + 0.10
    subtitle_top = hairline_top + 0.10

    if title_present:
        _add_text(slide, title, left=0.50, top=0.30, width=12.30, height=title_h,
                  size=28, color_rgb=INK_RGB, font=branding.MONO_FONT, bold=True)
        _add_rect(slide, left=0.50, top=hairline_top, width=12.30, height=0.005,
                  fill_rgb=RULE_RGB)

    # Estimate lede vertical room from char count so a long lede doesn't
    # overflow into the body region (slide 37 "Takeaways" had a 1.5-line
    # lede that ran into the top-row cards because the slot was a fixed
    # 0.40in). Cap at 1.0in so a sprawling lede doesn't eat the body.
    if lede and not use_side_by_side:
        est_h = _estimate_paragraph_height(lede, width=12.30, size=13,
                                           line_spacing=1.30)
        lede_h = min(1.0, max(0.40, est_h))
        _add_text(slide, lede, left=0.50, top=subtitle_top, width=12.30,
                  height=lede_h,
                  size=13, color_rgb=MUTED_RGB, font=branding.SANS_FONT)
    else:
        lede_h = 0.40

    if title_present:
        if use_side_by_side:
            body_top = hairline_top + 0.20
        else:
            # Push body_top below the dynamically-sized lede slot.
            body_top = subtitle_top + lede_h + 0.10
    else:
        body_top = 0.40
    body_bottom = 6.85
    body_h = body_bottom - body_top
    body_l = 0.50
    body_w = 12.30

    # Footer
    name = footer_kwargs.get("name", "")
    org = footer_kwargs.get("org", "")
    deck_title = footer_kwargs.get("deck_title", "")
    date = footer_kwargs.get("date", "")
    footer_parts = [p for p in (name, org, deck_title, date) if p]
    if footer_parts:
        _add_text(slide, "  ·  ".join(footer_parts),
                  left=0.50, top=7.12, width=12.30, height=0.30,
                  size=9, color_rgb=MUTED_RGB, font=branding.MONO_FONT)

    return body_top, body_h, body_l, body_w, body_bottom
