"""timeline layout: horizontal axis with milestone dots, labels above, dates below."""

from __future__ import annotations

import branding
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from ._common import (
    _add_chrome,
    _add_text,
    _set_bg,
    INK_RGB,
    WHITE_RGB,
)

# Visual constants
_AXIS_HEIGHT = 0.04        # thin horizontal axis line (in)
_DOT_SIZE    = 0.30        # diameter of milestone dot (in)
_LABEL_H     = 0.35        # height for label textbox above dot (in)
_DATE_H      = 0.28        # height for date textbox below dot (in)
_BODY_H      = 1.40        # height for body text below date (in)
_LABEL_SIZE  = 14          # label font size (pt)
_DATE_SIZE   = 11          # date font size (pt)
_BODY_SIZE   = 11          # body font size (pt)
# Vertical offset of axis from body_top
_AXIS_OFFSET_FRAC = 0.38   # axis sits at body_top + body_h * this fraction


def render(slide, *, params: dict, accent_rgb: RGBColor, footer_kwargs: dict) -> None:
    """Render a horizontal-timeline content slide.

    params keys:
        title       (str)
        lede        (str)
        milestones  list[{"date": str, "label": str, "body": str}]  — 2-6 items
    """
    title = params.get("title", "")
    lede = params.get("lede", "")
    milestones = list(params.get("milestones") or [])

    title_present = bool(title)
    title_wraps = len(title) > 30 if title_present else False

    _set_bg(slide, WHITE_RGB)

    body_top, body_h, body_l, body_w, body_bottom = _add_chrome(
        slide,
        title=title,
        lede=lede,
        footer_kwargs=footer_kwargs,
        accent=accent_rgb,
        title_present=title_present,
        title_wraps=title_wraps,
        use_side_by_side=False,
    )

    if not milestones:
        return

    n = len(milestones)

    # ── Axis line ───────────────────────────────────────────────────────────────
    axis_top = body_top + body_h * _AXIS_OFFSET_FRAC
    axis_center_y = axis_top + _AXIS_HEIGHT / 2

    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(body_l), Inches(axis_top),
        Inches(body_w), Inches(_AXIS_HEIGHT),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = accent_rgb
    shp.line.fill.background()
    shp.shadow.inherit = False

    # ── Milestone spacing ────────────────────────────────────────────────────────
    # Evenly distribute dot centres along body_w
    if n == 1:
        xs = [body_l + body_w / 2]
    else:
        step = body_w / (n - 1)
        xs = [body_l + i * step for i in range(n)]

    # Effective label+date width (centred on dot)
    item_w = min(body_w / n - 0.10, 2.20)
    item_w = max(item_w, 1.20)

    for i, milestone in enumerate(milestones):
        dot_cx = xs[i]
        dot_l = dot_cx - _DOT_SIZE / 2
        dot_t = axis_center_y - _DOT_SIZE / 2

        # Dot (filled circle)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(dot_l), Inches(dot_t),
            Inches(_DOT_SIZE), Inches(_DOT_SIZE),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent_rgb
        dot.line.fill.background()
        dot.shadow.inherit = False

        text_l = dot_cx - item_w / 2
        # Clamp to body bounds
        text_l = max(body_l, min(text_l, body_l + body_w - item_w))

        # Stagger: even indices → label above dot, odd → label shifted up more
        # (simple alternation: every milestone gets label above + date below)
        label_top = dot_t - _LABEL_H - 0.08
        date_top = dot_t + _DOT_SIZE + 0.08
        body_text_top = date_top + _DATE_H + 0.05

        # Label above dot (14pt mono bold INK)
        _add_text(
            slide, milestone.get("label", ""),
            left=text_l, top=label_top,
            width=item_w, height=_LABEL_H,
            size=_LABEL_SIZE, color_rgb=INK_RGB,
            font=branding.MONO_FONT, bold=True,
        )

        # Date below dot (11pt mono accent bold)
        _add_text(
            slide, milestone.get("date", ""),
            left=text_l, top=date_top,
            width=item_w, height=_DATE_H,
            size=_DATE_SIZE, color_rgb=accent_rgb,
            font=branding.MONO_FONT, bold=True,
        )

        # Body text below date (11pt sans INK)
        remaining_h = max(_BODY_H, body_bottom - (body_text_top + 0.10))
        _add_text(
            slide, milestone.get("body", ""),
            left=text_l, top=body_text_top,
            width=item_w, height=min(_BODY_H, remaining_h),
            size=_BODY_SIZE, color_rgb=INK_RGB,
            font=branding.SANS_FONT,
        )
