"""cards-heterogeneous layout: one large primary card left + 2-3 secondary cards stacked right.

Count-gated behavior:
  n_total ≤ 3 (1 primary + ≤2 secondaries): full-width stacked rows, vertically centered.
  n_total ≥ 4 (1 primary + ≥3 secondaries): primary on left 60%, secondaries tiled on right.
"""

from __future__ import annotations

from pathlib import Path

import branding
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from palette import LIGHT

from ._common import (
    _add_card,
    _add_chrome,
    _add_rect,
    _add_text,
    _estimate_paragraph_height,
    _set_bg,
)

# Stripe heights: primary card gets a thicker accent stripe to differentiate it.
_STRIPE_PRIMARY = 0.10   # primary card top accent stripe (stacked path)
_STRIPE_STD     = 0.06   # standard card top accent stripe


def _render_stacked_rows(
    slide, *,
    primary: dict,
    secondary_cards: list,
    body_top: float,
    body_h: float,
    body_l: float,
    body_w: float,
    accent_rgb: RGBColor,
    palette=LIGHT,
) -> None:
    """Render primary + secondaries as full-width stacked rows, vertically centered."""
    cards = [primary] + list(secondary_cards)
    n = len(cards)

    # --- Estimate each card's natural height ---
    _HEADER_PRIMARY = _STRIPE_PRIMARY + 0.14 + 0.50   # stripe + gap + label
    _HEADER_STD     = _STRIPE_STD     + 0.14 + 0.40
    _PAD_BOT        = 0.20
    _PAD_BODY_TOP_PRIMARY = _STRIPE_PRIMARY + 0.14 + 0.50   # same as header
    _PAD_BODY_TOP_STD     = _STRIPE_STD     + 0.14 + 0.40

    heights = []
    for idx, card in enumerate(cards):
        body_text = card.get("body", "")
        header = _HEADER_PRIMARY if idx == 0 else _HEADER_STD
        body_est = _estimate_paragraph_height(body_text, width=body_w - 0.36,
                                              size=13 if idx == 0 else 12,
                                              line_spacing=1.35)
        h = header + max(body_est, 0.30) + _PAD_BOT
        h = max(1.10, h)
        heights.append(h)

    gutter = 0.20
    total_h = sum(heights) + gutter * (n - 1)

    # Vertically center the stack if it fits; otherwise start at body_top.
    if total_h < body_h:
        cur_top = body_top + (body_h - total_h) / 2.0
    else:
        cur_top = body_top

    for idx, card in enumerate(cards):
        c_label = card.get("label", "")
        c_body  = card.get("body", "")
        c_icon  = card.get("icon")
        c_h     = heights[idx]
        is_primary = idx == 0

        stripe_h = _STRIPE_PRIMARY if is_primary else _STRIPE_STD

        # Card background
        _add_rect(slide, left=body_l, top=cur_top,
                  width=body_w, height=c_h, fill_rgb=palette.surface_rgb)

        if is_primary:
            # Primary card: full-width top stripe — the dominant accent element.
            _add_rect(slide, left=body_l, top=cur_top,
                      width=body_w, height=stripe_h, fill_rgb=accent_rgb)
        else:
            # Secondary cards: top-left partial stripe (1.0 in wide × 0.06 in
            # tall), flush with card top-left corner.  Signals association with
            # the primary accent without competing full-width stripes.
            _add_rect(slide, left=body_l, top=cur_top,
                      width=1.0, height=stripe_h, fill_rgb=accent_rgb)

        # Icon (optional)
        if c_icon is not None:
            c_icon = Path(c_icon) if not isinstance(c_icon, Path) else c_icon

        label_size = 16 if is_primary else 13
        body_size  = 13 if is_primary else 12
        label_top_offset = stripe_h + 0.14

        if c_icon is not None and c_icon.exists():
            try:
                slide.shapes.add_picture(
                    str(c_icon),
                    Inches(body_l + 0.18), Inches(cur_top + label_top_offset),
                    width=Inches(0.40 if is_primary else 0.32),
                    height=Inches(0.40 if is_primary else 0.32),
                )
                label_l = body_l + (0.70 if is_primary else 0.58)
                label_w = body_w - (0.88 if is_primary else 0.76)
            except Exception:
                label_l = body_l + 0.18
                label_w = body_w - 0.36
        else:
            label_l = body_l + 0.18
            label_w = body_w - 0.36

        label_h_alloc = 0.50 if is_primary else 0.40
        body_top_offset = stripe_h + 0.14 + label_h_alloc

        _add_text(slide, c_label,
                  left=label_l, top=cur_top + label_top_offset,
                  width=label_w, height=label_h_alloc,
                  size=label_size, color_rgb=accent_rgb,
                  font=branding.MONO_FONT, bold=True)
        _add_text(slide, c_body,
                  left=body_l + 0.18, top=cur_top + body_top_offset,
                  width=body_w - 0.36, height=c_h - body_top_offset - _PAD_BOT,
                  size=body_size, color_rgb=palette.text_rgb, font=branding.SANS_FONT)

        cur_top += c_h + gutter


def render(slide, *, params: dict, accent_rgb: RGBColor, footer_kwargs: dict, palette=LIGHT) -> None:
    """Render a heterogeneous-cards slide.

    params keys:
        title           (str)
        lede            (str)
        primary_card    {"label": str, "body": str, "icon": Path | None}
        secondary_cards list[{"label": str, "body": str, "icon": Path | None}]  — 2-3 items

    Layout is count-gated:
      n_total ≤ 3 → stacked full-width rows (primary on top, primary differentiated by
                    thicker stripe and larger label font).
      n_total ≥ 4 → original tile layout (primary on left 60%, secondaries tiled right).
    """
    title = params.get("title", "")
    lede = params.get("lede", "")
    primary = params.get("primary_card") or {}
    secondary_cards = list(params.get("secondary_cards") or [])

    n_total = 1 + len(secondary_cards)

    title_present = bool(title)
    title_wraps = len(title) > 30 if title_present else False

    _set_bg(slide, palette.canvas_rgb)

    body_top, body_h, body_l, body_w, body_bottom = _add_chrome(
        slide,
        title=title,
        lede=lede,
        footer_kwargs=footer_kwargs,
        accent=accent_rgb,
        title_present=title_present,
        title_wraps=title_wraps,
        use_side_by_side=False,
        on_dark=palette.on_dark,
        palette=palette,
    )

    # --- Count gate ---
    if n_total <= 3:
        _render_stacked_rows(
            slide,
            primary=primary,
            secondary_cards=secondary_cards,
            body_top=body_top,
            body_h=body_h,
            body_l=body_l,
            body_w=body_w,
            accent_rgb=accent_rgb,
            palette=palette,
        )
        return

    # --- n_total ≥ 4: original tile layout ---
    primary_w = body_w * 0.60
    right_gap = 0.20
    right_l = body_l + primary_w + right_gap
    right_w = body_w - primary_w - right_gap

    # --- Primary card (bigger label 15pt, body 13pt) ---
    p_label = primary.get("label", "")
    p_body = primary.get("body", "")
    p_icon = primary.get("icon")

    # Measure content height so the card sizes to its text rather than
    # bloating to fill the full column (Bug 2 fix).
    _LABEL_H  = 0.50   # label textbox allocation
    _HEADER   = 0.80   # top stripe + padding above label + label itself
    _PAD_BOT  = 0.20   # breathing room below body
    body_est = _estimate_paragraph_height(p_body, width=primary_w - 0.36,
                                          size=13, line_spacing=1.35)
    content_h = _HEADER + max(body_est, 0.30) + _PAD_BOT
    # Cap at body_h; never smaller than a sensible minimum
    card_h = min(body_h, max(1.20, content_h))
    # Vertically centre the card within the column
    card_top = body_top + (body_h - card_h) / 2.0

    _add_rect(slide, left=body_l, top=card_top, width=primary_w, height=card_h,
              fill_rgb=palette.surface_rgb)
    _add_rect(slide, left=body_l, top=card_top, width=primary_w, height=_STRIPE_STD,
              fill_rgb=accent_rgb)

    if p_icon is not None:
        p_icon = Path(p_icon) if not isinstance(p_icon, Path) else p_icon

    if p_icon is not None and p_icon.exists():
        try:
            slide.shapes.add_picture(
                str(p_icon),
                Inches(body_l + 0.18), Inches(card_top + 0.20),
                width=Inches(0.40), height=Inches(0.40),
            )
            label_l = body_l + 0.70
            label_w = primary_w - 0.88
        except Exception:
            label_l = body_l + 0.18
            label_w = primary_w - 0.36
    else:
        label_l = body_l + 0.18
        label_w = primary_w - 0.36

    _add_text(slide, p_label, left=label_l, top=card_top + 0.20,
              width=label_w, height=_LABEL_H,
              size=15, color_rgb=accent_rgb, font=branding.MONO_FONT, bold=True)
    _add_text(slide, p_body, left=body_l + 0.18, top=card_top + 0.80,
              width=primary_w - 0.36, height=card_h - 0.90,
              size=13, color_rgb=palette.text_rgb, font=branding.SANS_FONT)

    # --- Secondary cards stacked vertically ---
    n = max(1, len(secondary_cards))
    gutter = 0.15
    sec_h = (body_h - gutter * (n - 1)) / n

    for i, card in enumerate(secondary_cards):
        cy = body_top + i * (sec_h + gutter)
        _add_card(
            slide,
            label=card.get("label", ""),
            body=card.get("body", ""),
            left=right_l,
            top=cy,
            width=right_w,
            height=sec_h,
            accent_rgb=accent_rgb,
            icon_path=card.get("icon"),
            surface_rgb=palette.surface_rgb,
            text_rgb=palette.text_rgb,
        )
