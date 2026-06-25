"""Embed institutional logos on a deck's title slide from the superstack/logos convention.

Generic + machine-agnostic. Any PNGs in `superstack/logos/` (resolved relative
to this file, so it travels with the repo wherever superstack is checked out)
are embedded as a lower-right paired block (vertical-center) on the title slide.
If the folder is empty or absent, this is a silent no-op — so build-pptx stays
generic and decks built on machines without logos are unaffected. The logo image
files are gitignored, so each machine populates its own set (UCSF + Cal here, a
CurieDx logo on another machine, etc.).

Order and per-logo height scale come from an optional logos/manifest.json:

    {
      "title": [
        {"file": "ucsf-logo.png", "scale": 1.40},
        {"file": "cal-logo.png",  "scale": 1.0}
      ]
    }

Scale compensates for assets with canvas padding (the UCSF disc has ~30%
transparent padding, so 1.40x makes the visible mark match Cal's edge-to-edge
size). Without a manifest: all PNGs alphabetical, equal height.

Usage as a library (build-pptx render tail):

    from apply_logos import apply_logos_to_deck
    apply_logos_to_deck("deck.pptx")   # no-op if ~/arcadia/logos absent

Usage as a CLI (backfill an existing deck):

    python apply_logos.py deck.pptx
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

# superstack/logos — resolved relative to this file (skills/_shared/apply_logos.py),
# so it works regardless of where superstack is checked out on a given machine.
LOGOS_DIR = Path(__file__).resolve().parent.parent.parent / "logos"


def _logo_spec() -> list[tuple[Path, float]]:
    """Return [(path, scale), ...] in render order, or [] if no logos available."""
    if not LOGOS_DIR.is_dir():
        return []
    manifest_path = LOGOS_DIR / "manifest.json"
    if manifest_path.exists():
        try:
            manifest = json.loads(manifest_path.read_text())
            entries = manifest.get("title", [])
            spec = []
            for e in entries:
                p = LOGOS_DIR / e["file"]
                if p.exists():
                    spec.append((p, float(e.get("scale", 1.0))))
            if spec:
                return spec
        except Exception:
            pass  # fall through to convention
    # No manifest (or unusable): all PNGs alphabetical, equal height
    pngs = sorted(LOGOS_DIR.glob("*.png"))
    return [(p, 1.0) for p in pngs]


def apply_logos_to_deck(
    pptx_path: str | Path,
    *,
    base_height_in: float = 0.9,
    gap_in: float = 0.25,
    margin_in: float = 0.5,
    slide_index: int = 0,
) -> bool:
    """Embed logos on one slide of a .pptx. Returns True if any were applied.

    No-op (returns False) if ~/arcadia/logos/ is absent or empty. Saves in place.
    """
    spec = _logo_spec()
    if not spec:
        return False

    from pptx import Presentation
    from pptx.util import Emu, Inches

    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))
    if slide_index >= len(prs.slides):
        return False
    slide = prs.slides[slide_index]
    SLIDE_W, SLIDE_H = prs.slide_width, prs.slide_height
    base_h = Inches(base_height_in)
    gap = Inches(gap_in)

    pics = []
    for path, scale in spec:
        target_h = int(base_h * scale)
        pic = slide.shapes.add_picture(str(path), Inches(0), Inches(0))
        aspect = pic.width / pic.height
        pic.height = target_h
        pic.width = int(target_h * aspect)
        pics.append(pic)

    block_h = max(p.height for p in pics)
    total_w = sum(p.width for p in pics) + gap * (len(pics) - 1)
    block_left = SLIDE_W - total_w - Inches(margin_in)
    block_bottom = SLIDE_H - Inches(margin_in)
    block_center_y = (block_bottom - block_h) + block_h // 2

    cursor = block_left
    for pic in pics:
        pic.left = int(cursor)
        pic.top = int(block_center_y - pic.height // 2)
        cursor = int(cursor + pic.width + gap)

    prs.save(str(pptx_path))
    return True


def deck_opts_out(markdown_path: str | Path) -> bool:
    """True if the deck's frontmatter sets `logos: false` (or no/off)."""
    md = Path(markdown_path)
    if not md.exists():
        return False
    text = md.read_text()
    if not text.startswith("---"):
        return False
    end = text.find("---", 3)
    if end < 0:
        return False
    fm = text[3:end]
    for line in fm.splitlines():
        if line.strip().lower().replace(" ", "").startswith("logos:"):
            val = line.split(":", 1)[1].strip().lower()
            return val in ("false", "no", "off", "0", "none")
    return False


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("usage: python apply_logos.py <deck.pptx> [deck.md]", file=sys.stderr)
        sys.exit(2)
    pptx = sys.argv[1]
    md = sys.argv[2] if len(sys.argv) > 2 else None
    if md and deck_opts_out(md):
        print(f"deck opts out of logos (frontmatter logos: false) — skipping {pptx}")
        sys.exit(0)
    applied = apply_logos_to_deck(pptx)
    if applied:
        n = len(_logo_spec())
        print(f"applied {n} logo(s) to title slide of {Path(pptx).name}")
    else:
        print(f"no logos applied (~/arcadia/logos absent or empty)")
