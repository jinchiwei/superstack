"""Build a session results pptx for an autoresearch session.

Reads:
  results/{date}_{scope}/README.md          - session header + candidate table
  results/{date}_{scope}/iter-*/summary.md  - one slide per iteration
  results/{date}_{scope}/iter-*/figures/*   - optional: embedded on the slide

Writes:
  docs/runs/{date}_{scope}/SESSION_REPORT.pptx

Style: conference-presentation polish. 16:9 widescreen. Brand palette
(turquoise / deeppink / amber / blueviolet) with confident colored backgrounds
on title + section + summary slides; clean white per-iteration slides with
an accent rail and large headers. No internal references — slugs are
humanized into title-case for display.

Run:
    python docs/_build_pptx.py --date 2026-04-30 --scope fw-arch-sweep

Project-specific styling: extend or replace this template freely.
"""
from __future__ import annotations

import argparse
import re
from datetime import datetime
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt, Emu
except ImportError:
    raise SystemExit("python-pptx not installed. Run: pip install python-pptx")


# ----- Brand palette -----
TURQUOISE  = RGBColor(0x40, 0xE0, 0xD0)
DEEPPINK   = RGBColor(0xFF, 0x14, 0x93)
AMBER      = RGBColor(0xF0, 0xC8, 0x40)
BLUEVIOLET = RGBColor(0x8A, 0x2B, 0xE2)
INK        = RGBColor(0x14, 0x14, 0x18)   # near-black, slight cool tint
PAPER      = RGBColor(0xFA, 0xFA, 0xFC)   # near-white, slight cool tint
MUTED      = RGBColor(0x6B, 0x6B, 0x73)
RULE       = RGBColor(0xE5, 0xE5, 0xEA)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DEEP_BG    = RGBColor(0x18, 0x10, 0x28)   # title-slide deep purple-black

MONO = "Geist Mono"
BODY = "Geist"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ----- Helpers -----
def _humanize(slug: str) -> str:
    """fw-arch-sweep -> 'FW Arch Sweep'. iter-03_resnet50-lr1e-3 -> 'Resnet50 Lr1e 3'."""
    if not slug:
        return ""
    # Strip leading 'iter-NN_' if present
    slug = re.sub(r"^iter-\d+_", "", slug)
    parts = re.split(r"[-_]+", slug)
    return " ".join(p[:1].upper() + p[1:] if p else "" for p in parts).strip()


def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, *, left, top, width, height, fill=None, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    shp.shadow.inherit = False
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    return shp


def _add_text(slide, text, *, left, top, width, height, size=18, color=INK,
              font=BODY, bold=False, italic=False, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return tb


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ----- Slide builders -----
def _title_slide(prs, *, scope, scope_text, date_str, target):
    s = _blank(prs)
    _set_bg(s, DEEP_BG)
    # Accent rail at left
    _add_rect(s, left=0, top=0, width=0.45, height=7.5, fill=TURQUOISE)
    # A subtle accent stripe along the bottom
    _add_rect(s, left=0.45, top=7.20, width=12.88, height=0.06, fill=DEEPPINK)
    # Eyebrow
    _add_text(s, "RESEARCH SESSION", left=1.2, top=2.0, width=11, height=0.4,
              size=11, color=TURQUOISE, font=MONO, bold=True)
    # Big title
    _add_text(s, _humanize(scope), left=1.2, top=2.45, width=11, height=2.2,
              size=52, color=WHITE, font=BODY, bold=True)
    # Scope sentence
    _add_text(s, scope_text, left=1.2, top=4.55, width=11, height=1.1,
              size=18, color=RGBColor(0xC9, 0xC4, 0xD2), font=BODY)
    # Target chip
    if target and target.strip() and target.strip() != "(unspecified)":
        _add_rect(s, left=1.2, top=5.85, width=4.0, height=0.45, fill=DEEPPINK)
        _add_text(s, f"TARGET   {target}", left=1.35, top=5.91, width=3.85, height=0.35,
                  size=11, color=WHITE, font=MONO, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    # Date
    _add_text(s, date_str.upper(), left=1.2, top=6.7, width=10, height=0.4,
              size=10, color=AMBER, font=MONO, bold=True)


def _section_divider(prs, *, label, color):
    s = _blank(prs)
    _set_bg(s, color)
    _add_rect(s, left=0, top=0, width=0.45, height=7.5, fill=INK)
    _add_text(s, label.upper(), left=1.2, top=3.0, width=11, height=1.5,
              size=44, color=WHITE, font=BODY, bold=True)
    _add_text(s, "—", left=1.2, top=4.6, width=2, height=0.5,
              size=28, color=WHITE, font=BODY)


def _candidate_slide(prs, *, iter_num, candidate, summary_text, status, metric,
                     figures):
    s = _blank(prs)
    _set_bg(s, PAPER)

    # Left accent rail
    _add_rect(s, left=0, top=0, width=0.30, height=7.5, fill=TURQUOISE)

    # Status pill — color depends on outcome
    status_clean = (status or "").lower()
    if any(k in status_clean for k in ("complete", "success", "ok")):
        chip = TURQUOISE
    elif any(k in status_clean for k in ("running", "pending")):
        chip = AMBER
    elif any(k in status_clean for k in ("fail", "error", "halt")):
        chip = DEEPPINK
    else:
        chip = MUTED

    # Eyebrow: iter index
    _add_text(s, f"ITERATION {iter_num:02d}", left=0.85, top=0.55, width=4, height=0.35,
              size=10, color=MUTED, font=MONO, bold=True)

    # Big candidate title
    _add_text(s, _humanize(candidate), left=0.85, top=0.9, width=10.5, height=1.0,
              size=32, color=INK, font=BODY, bold=True)

    # Status chip + metric chip on the right
    if status:
        _add_rect(s, left=10.6, top=0.6, width=2.0, height=0.4, fill=chip)
        _add_text(s, (status or "").upper(), left=10.6, top=0.66, width=2.0, height=0.3,
                  size=10, color=WHITE, font=MONO, bold=True,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if metric and metric != "—":
        _add_text(s, "METRIC", left=10.6, top=1.1, width=2.0, height=0.25,
                  size=8, color=MUTED, font=MONO, bold=True, align=PP_ALIGN.CENTER)
        _add_text(s, str(metric), left=10.6, top=1.32, width=2.0, height=0.5,
                  size=22, color=BLUEVIOLET, font=MONO, bold=True,
                  align=PP_ALIGN.CENTER)

    # Hairline rule
    _add_rect(s, left=0.85, top=2.05, width=11.5, height=0.015, fill=RULE)

    # Body — split between text (left) and figure (right) if figure exists
    body_lines = summary_text.strip().splitlines() if summary_text else []
    body = "\n".join(body_lines[:24]) if body_lines else "—"

    if figures:
        _add_text(s, body, left=0.85, top=2.3, width=6.4, height=4.6,
                  size=11, color=INK, font=MONO)
        # Embed first figure
        fig = figures[0]
        try:
            slide_pic = s.shapes.add_picture(str(fig), Inches(7.5), Inches(2.3),
                                             width=Inches(5.0))
            # Constrain height
            if slide_pic.height > Inches(4.6):
                slide_pic.height = Inches(4.6)
                slide_pic.width = int(slide_pic.width * 4.6 / (slide_pic.height / Inches(1)))
        except Exception:
            pass
        # Caption
        _add_text(s, fig.name, left=7.5, top=7.0, width=5.0, height=0.3,
                  size=8, color=MUTED, font=MONO, align=PP_ALIGN.LEFT)
    else:
        _add_text(s, body, left=0.85, top=2.3, width=11.5, height=4.7,
                  size=12, color=INK, font=MONO)


def _summary_slide(prs, *, scope, candidates, target):
    s = _blank(prs)
    _set_bg(s, INK)
    _add_rect(s, left=0, top=0, width=0.45, height=7.5, fill=AMBER)

    _add_text(s, "RESULTS", left=1.2, top=0.55, width=10, height=0.4,
              size=11, color=AMBER, font=MONO, bold=True)
    _add_text(s, _humanize(scope), left=1.2, top=0.95, width=10, height=0.7,
              size=28, color=WHITE, font=BODY, bold=True)
    if target and target.strip() and target.strip() != "(unspecified)":
        _add_text(s, f"target — {target}", left=1.2, top=1.65, width=10, height=0.4,
                  size=12, color=TURQUOISE, font=MONO)

    # Hairline rule
    _add_rect(s, left=1.2, top=2.2, width=11.0, height=0.012, fill=MUTED)

    # Header row
    headers = [("ITER", 0.6), ("CANDIDATE", 5.5), ("STATUS", 1.6), ("METRIC", 2.0)]
    x = 1.2
    for label, w in headers:
        _add_text(s, label, left=x, top=2.35, width=w, height=0.3,
                  size=9, color=MUTED, font=MONO, bold=True)
        x += w

    # Body rows
    y = 2.75
    row_h = 0.42
    for c in candidates[:11]:
        x = 1.2
        # iter
        _add_text(s, f"{c['iter']:02d}", left=x, top=y, width=0.6, height=row_h,
                  size=12, color=AMBER, font=MONO, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        x += 0.6
        # candidate
        _add_text(s, _humanize(c["candidate"]), left=x, top=y, width=5.5, height=row_h,
                  size=12, color=WHITE, font=BODY, anchor=MSO_ANCHOR.MIDDLE)
        x += 5.5
        # status pill
        sc = (c.get("status") or "").lower()
        if any(k in sc for k in ("complete", "success", "ok")):
            chip = TURQUOISE
        elif any(k in sc for k in ("running", "pending")):
            chip = AMBER
        elif any(k in sc for k in ("fail", "error", "halt")):
            chip = DEEPPINK
        else:
            chip = MUTED
        _add_rect(s, left=x, top=y + 0.07, width=1.4, height=0.28, fill=chip)
        _add_text(s, (c.get("status") or "—").upper(), left=x, top=y + 0.08,
                  width=1.4, height=0.28, size=9, color=WHITE, font=MONO,
                  bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += 1.6
        # metric
        _add_text(s, str(c.get("metric") or "—"), left=x, top=y, width=2.0, height=row_h,
                  size=12, color=BLUEVIOLET, font=MONO, bold=True,
                  anchor=MSO_ANCHOR.MIDDLE)
        y += row_h

    if len(candidates) > 11:
        _add_text(s, f"+ {len(candidates) - 11} more iterations", left=1.2, top=y + 0.1,
                  width=11, height=0.3, size=10, color=MUTED, font=MONO, italic=True)


def _read_session(date_str, scope):
    session_dir = Path("results") / f"{date_str}_{scope}"
    readme = session_dir / "README.md"
    if not readme.exists():
        raise SystemExit(f"no session at {session_dir}")
    text = readme.read_text()
    scope_text = re.search(r"\*\*Scope:\*\*\s*(.+)", text)
    target = re.search(r"\*\*Target metric:\*\*\s*(.+)", text)

    candidates = []
    for d in sorted(session_dir.glob("iter-*")):
        if not d.is_dir():
            continue
        m = re.match(r"iter-(\d+)_(.+)", d.name)
        if not m:
            continue
        iter_num = int(m.group(1))
        cand_slug = m.group(2)
        summary = d / "summary.md"
        body = summary.read_text() if summary.exists() else ""
        metric_match = re.search(r"^\s*metric\s*[:|]\s*(.+)$", body, re.MULTILINE | re.IGNORECASE)
        status_match = re.search(r"^\s*status\s*[:|]\s*(.+)$", body, re.MULTILINE | re.IGNORECASE)
        # Pick up to 4 figures
        fig_dir = d / "figures"
        figures = sorted(fig_dir.glob("*"))[:4] if fig_dir.is_dir() else []
        figures = [f for f in figures if f.suffix.lower() in (".png", ".jpg", ".jpeg", ".gif")]
        candidates.append({
            "iter": iter_num,
            "candidate": cand_slug,
            "summary": body,
            "metric": metric_match.group(1).strip() if metric_match else "—",
            "status": status_match.group(1).strip() if status_match else "—",
            "figures": figures,
        })

    return {
        "scope_text": scope_text.group(1).strip() if scope_text else _humanize(scope),
        "target": target.group(1).strip() if target else "",
        "candidates": candidates,
    }


def main():
    p = argparse.ArgumentParser(description=__doc__)
    p.add_argument("--date", default=datetime.now().strftime("%Y-%m-%d"),
                   help="Session date (YYYY-MM-DD).")
    p.add_argument("--scope", required=True, help="Session scope slug.")
    args = p.parse_args()

    session = _read_session(args.date, args.scope)
    out_dir = Path("docs") / "runs" / f"{args.date}_{args.scope}"
    out_dir.mkdir(parents=True, exist_ok=True)
    out = out_dir / "SESSION_REPORT.pptx"

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Title
    _title_slide(prs, scope=args.scope, scope_text=session["scope_text"],
                 date_str=args.date, target=session["target"])
    # 2. Section divider — Iterations
    _section_divider(prs, label="Iterations", color=TURQUOISE)
    # 3..N. Per-iteration
    for c in session["candidates"]:
        _candidate_slide(prs, iter_num=c["iter"], candidate=c["candidate"],
                         summary_text=c["summary"], status=c["status"],
                         metric=c["metric"], figures=c["figures"])
    # N+1. Section divider — Results
    _section_divider(prs, label="Results", color=DEEPPINK)
    # N+2. Summary
    _summary_slide(prs, scope=args.scope, candidates=session["candidates"],
                   target=session["target"])

    prs.save(out)
    print(f"wrote {out}")


if __name__ == "__main__":
    main()
